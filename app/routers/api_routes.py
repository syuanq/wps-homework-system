# -*- coding: utf-8 -*-
"""
API路由
"""
import os
import re
import json
import uuid
import io
from datetime import datetime
from collections import defaultdict
from flask import Blueprint, request, jsonify, current_app
from app.services.scoring_engine import ScoringEngine
from app.services.scoring_rules import ALL_TASKS, MODULE_GROUPS, LEVEL_RULES, get_level
from app.services.ai_service import DeepSeekService
from app.services import student_service_db as student_service
from app.services import document_generator
from app.services import submission_service

api_bp = Blueprint('api', __name__)


def _clean_advice(text):
    """清理建议文本中的markdown代码块标记"""
    if not text:
        return text
    text = text.strip()
    # 去掉开头的 ``` 或 ```json
    if text.startswith('```'):
        first_nl = text.find('\n')
        if first_nl != -1:
            text = text[first_nl + 1:]
        else:
            text = text[3:]
    # 去掉结尾的 ```
    if text.rstrip().endswith('```'):
        text = text.rstrip()[:-3].rstrip()
    return text.strip()


def api_login_required(f):
    """API登录检查"""
    from functools import wraps
    @wraps(f)
    def wrapper(*args, **kwargs):
        from app.services.auth_db import is_logged_in
        if not is_logged_in():
            return jsonify({"success": False, "message": "请先登录"}), 401
        return f(*args, **kwargs)
    return wrapper


def api_teacher_required(f):
    """API教师权限检查"""
    from functools import wraps
    @wraps(f)
    def wrapper(*args, **kwargs):
        from app.services.auth_db import is_logged_in, is_teacher
        if not is_logged_in():
            return jsonify({"success": False, "message": "请先登录"}), 401
        if not is_teacher():
            return jsonify({"success": False, "message": "需要教师权限"}), 403
        return f(*args, **kwargs)
    return wrapper

# 内存存储(生产环境应使用数据库)
submissions = {}
DATA_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), 'data')
RESULT_FOLDER = os.path.join(DATA_DIR, 'results')
ANSWER_FOLDER = os.path.join(DATA_DIR, 'answers')

# 标准答案目录和索引
ANSWERS_DIR = os.path.join(DATA_DIR, 'answers')
ANSWERS_INDEX_FILE = os.path.join(ANSWERS_DIR, 'index.json')

# 学生信息索引(从students.json或users.json)
_STUDENTS_INDEX_CACHE = None


def _get_students_index():
    """获取学生信息索引,优先从students.json获取,如果没有则使用users.json"""
    """返回的索引以学号(student_no)为key,因为提交记录中的user_id是学号"""

    global _STUDENTS_INDEX_CACHE
    if _STUDENTS_INDEX_CACHE is not None:
        return _STUDENTS_INDEX_CACHE

    students_index = {}

    # 优先从students.json读取
    students_file = os.path.join(DATA_DIR, 'students.json')
    if os.path.exists(students_file):
        try:
            with open(students_file, 'r', encoding='utf-8') as f:
                students_data = json.load(f)
            if isinstance(students_data, list):
                for student in students_data:
                    # 使用学号作为key(与提交记录中的user_id对应?
                    student_no = student.get('student_no', '')
                    if student_no:
                        students_index[student_no] = {
                            'name': student.get('name', ''),
                            'class_name': student.get('class_name', ''),
                            'id': student.get('id', '')  # 保留内部ID
                        }
        except Exception:
            pass

    # 如果students.json为空或不存在,使用users.json作为备用
    if not students_index:
        users_file = os.path.join(DATA_DIR, 'users.json')
        if os.path.exists(users_file):
            try:
                with open(users_file, 'r', encoding='utf-8') as f:
                    users_data = json.load(f)
                for user_id, user_info in users_data.items():
                    if user_info.get('role') == 'student':
                        students_index[user_id] = {
                            'name': user_info.get('name', ''),
                            'class_name': user_info.get('class_name', '')
                        }
            except Exception:
                pass

    _STUDENTS_INDEX_CACHE = students_index
    return students_index


def _is_teacher_submission(submission):
    """判断提交记录是否为教师提交"""
    user_id = submission.get('user_id', '')
    students_index = _get_students_index()
    # 如果user_id不在学生索引中,可能不是学生提交
    # 检查数据库中该用户的角色
    users_file = os.path.join(DATA_DIR, 'users.json')
    if os.path.exists(users_file):
        try:
            with open(users_file, 'r', encoding='utf-8') as f:
                users_data = json.load(f)
            if user_id in users_data:
                return users_data[user_id].get('role') == 'teacher'
        except Exception:
            pass
    return False


def _enrich_submission_with_student_info(submission):
    """为提交记录补充学生姓名和班级信息(优先从学生基本信息获取)"""
    user_id = submission.get('user_id', '')
    students_index = _get_students_index()
    student_info = students_index.get(user_id, {})

    # 学生姓名和班级信息优先从学生基本信息获取(唯一来源)
    if student_info.get('name'):
        submission['student_name'] = student_info['name']
    elif not submission.get('student_name'):
        submission['student_name'] = user_id  # 备用:使用user_id

    if student_info.get('class_name'):
        submission['student_class'] = student_info['class_name']
    elif not submission.get('student_class'):
        submission['student_class'] = '未分班'

    return submission


def _get_all_students_for_stats():
    """获取所有学生列表(用于计算已交/待交统计)
    返回的学生列表包含学号(student_no),用于与提交记录中的user_id匹配
    """
    students = []

    # 从数据库读取
    students_file = os.path.join(DATA_DIR, 'students.json')
    if os.path.exists(students_file):
        try:
            with open(students_file, 'r', encoding='utf-8') as f:
                students_data = json.load(f)
            if isinstance(students_data, list):
                for student in students_data:
                    students.append({
                        'id': student.get('id', ''),
                        'student_no': student.get('student_no', ''),  # 学号,用于匹配user_id
                        'name': student.get('name', ''),
                        'class_name': student.get('class_name', '')
                    })
        except Exception:
            pass

    # 如果数据库为空
    if not students:
        users_file = os.path.join(DATA_DIR, 'users.json')
        if os.path.exists(users_file):
            try:
                with open(users_file, 'r', encoding='utf-8') as f:
                    users_data = json.load(f)
                for user_id, user_info in users_data.items():
                    if user_info.get('role') == 'student':
                        students.append({
                            'id': user_id,
                            'student_no': user_id,  # users.json的key就是学号
                            'name': user_info.get('name', ''),
                            'class_name': user_info.get('class_name', '')
                        })
            except Exception:
                pass

    return students


def _get_answers_index():
    """获取标准答案索引"""
    if os.path.exists(ANSWERS_INDEX_FILE):
        with open(ANSWERS_INDEX_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}


def _save_answers_index(index):
    """保存标准答案索引"""
    os.makedirs(ANSWERS_DIR, exist_ok=True)
    with open(ANSWERS_INDEX_FILE, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)


def _get_answer_file(task_id):
    """获取任务的标准答案文件路径
    优先从索引中查找,如果索引中没有则回退到旧方式
    """
    # 首先尝试从索引获取
    index = _get_answers_index()
    answers = index.get(task_id, [])
    if answers:
        filepath = answers[-1].get('filepath')
        if filepath:
            full_path = os.path.join(ANSWERS_DIR, filepath)
            if os.path.exists(full_path):
                return full_path
    
    # 回退到旧方式(兼容已有答案文件)
    if not os.path.exists(ANSWER_FOLDER):
        return None
    for f in os.listdir(ANSWER_FOLDER):
        if task_id in f and (f.endswith('.xlsx') or f.endswith('.docx') or f.endswith('.pptx')):
            return os.path.join(ANSWER_FOLDER, f)
    return None


def _parse_student_from_filename(filename):
    """

    支持格式:
    - 姓名+班级_xxxx.ext  (如:张三+计科1班_task_4_2.docx)
    - 姓名_班级_xxxx.ext  (如:张三_计科1班_task_4_2.docx)
    - 姓名+班级.ext       (如:张三+计科1班.docx)
    - 姓名_班级.ext       (如:张三_计科1班.docx)
    """
    # 去掉扩展名
    name = os.path.splitext(filename)[0]

    # 先尝试用 + 分隔
    if '+' in name:
        parts = name.split('+', 1)
        student_name = parts[0].strip()
        rest = parts[1].strip()
        # rest 可能"计科1班_task_4_2" "计科1?
        if '_' in rest:
            student_class = rest.split('_')[0].strip()
        else:
            student_class = rest
    # 再尝试用 _ 分隔(取前两段作为姓名和班级.
    elif '_' in name:
        parts = name.split('_')
        student_name = parts[0].strip()
        student_class = parts[1].strip()

    # 如果解析不到,用整个文件名作为姓?
    if not student_name:
        student_name = name

    return student_name, student_class


@api_bp.route('/tasks/list', methods=['GET'])
def get_tasks():
    """获取所有可用任务列表"""
    tasks = []
    for task_id, rules in ALL_TASKS.items():
        tasks.append({
            "task_id": task_id,
            "task_name": rules["task_name"],
            "module": rules["module"],
            "file_type": rules["file_type"],
            "description": rules["description"],
            "max_score": rules["max_score"],
            "check_items_count": len(rules["check_items"]),
        })
    return jsonify({"success": True, "tasks": tasks})


@api_bp.route('/tasks/grouped', methods=['GET'])
def get_tasks_grouped():
    """获取按模块分组的任务列表"""
    groups = {}
    for group_key, group_info in MODULE_GROUPS.items():
        tasks = []
        for task_id in group_info["tasks"]:
            if task_id in ALL_TASKS:
                rules = ALL_TASKS[task_id]
                tasks.append({
                    "task_id": task_id,
                    "task_name": rules["task_name"],
                    "module": rules["module"],
                    "file_type": rules["file_type"],
                    "description": rules["description"],
                    "max_score": rules["max_score"],
                    "check_items_count": len(rules["check_items"]),
                })
        groups[group_key] = {
            "name": group_info["name"],
            "icon": group_info["icon"],
            "color": group_info["color"],
            "tasks": tasks,
        }
    return jsonify({"success": True, "groups": groups})


@api_bp.route('/submit', methods=['POST'])
@api_login_required
def submit_homework():
    """提交作业并评分"""
    try:
        task_id = request.form.get('task_id')
        file = request.files.get('file')

        if not task_id or not file:
            return jsonify({"success": False, "message": "请选择任务并上传文件"}), 400

        if task_id not in ALL_TASKS:
            return jsonify({"success": False, "message": f"无效的任务ID: {task_id}"}), 400

        # 检查文件类型(支持旧格.doc/.xls/.ppt?
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        expected_ext = f".{ALL_TASKS[task_id]['file_type']}"
        # 旧格式映?
        old_ext_map = {'.doc': '.docx', '.xls': '.xlsx', '.ppt': '.pptx'}
        if ext != expected_ext and ext not in old_ext_map:
            return jsonify({"success": False, "message": f"文件格式错误,请上传 {expected_ext} 格式文件"}), 400

        # 获取当前登录用户
        from app.services.auth_db import get_current_user
        user = get_current_user()

        # 从文件名解析姓名和班?
        student_name, student_class = _parse_student_from_filename(filename)

        # 如果文件名中无法解析,从登录用户信息补充
        if not student_name and user and user.get('name'):
            student_name = user['name']
        if not student_class and user and user.get('class_name'):
            student_class = user['class_name']

        # 检查学号是否与登录用户一?
        if user and user.get('role') == 'student':
            login_username = user.get('username', '')
            # 从文件名提取学号(文件名第一段,?2025122755_任务.docx"中的2025122755?
            file_student_id = os.path.splitext(filename)[0].split('_')[0].split('+')[0].strip()
            if not file_student_id or not file_student_id.isdigit():
                return jsonify({
                    "success": False,
                    "message": f"文件名格式错误！文件名必须以学号开头,格式如:<strong>{login_username}_任务.docx</strong>.当前文件名:{filename}"
                }), 400
            if file_student_id != login_username:
                return jsonify({
                    "success": False,
                    "message": f"学号不一致！登录学号是<strong>{login_username}</strong>,但文件名中学号是<strong>{file_student_id}</strong>.请确认文件名以您的学号开头,或使用正确的账号登录."
                }), 400

        # 保存文件
        submission_id = str(uuid.uuid4())[:8]
        upload_dir = current_app.config['UPLOAD_FOLDER']
        os.makedirs(upload_dir, exist_ok=True)
        filepath = os.path.join(upload_dir, f"{submission_id}_{filename}")
        file.save(filepath)

        # 如果是旧格式(.doc/.xls/.ppt),自动转换为新格式
        if ext in old_ext_map:
            try:
                import subprocess, tempfile, glob
                new_ext = old_ext_map[ext]
                output_dir = tempfile.mkdtemp()
                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', new_ext[1:], '--outdir', output_dir, filepath],
                    capture_output=True, text=True, timeout=30
                )
                # LibreOffice保持原文件名,查找转换后的文?
                converted_files = glob.glob(os.path.join(output_dir, f"*{new_ext}"))
                if converted_files:
                    filepath = converted_files[0]
                else:
                    return jsonify({"success": False, "message": f"文件转换失败,请手动另存为 {new_ext} 格式后重新上传"}), 400
            except subprocess.TimeoutExpired:
                return jsonify({"success": False, "message": "文件转换超时,请手动另存为格式后重新上传"}), 400
            except Exception as e:
                return jsonify({"success": False, "message": f"文件转换出错:{str(e)}"}), 400

        # 执行评分
        engine = ScoringEngine(task_id, filepath)
        result = engine.score()

        # 如果有答案文件,优先使用答案对比评分(更准确?
        answer_file = _get_answer_file(task_id)
        if answer_file and os.path.exists(answer_file):
            try:
                from app.services.answer_comparator import ExcelAnswerComparator
                comp = ExcelAnswerComparator(answer_file, filepath)
                comp_result = comp.compare_task(task_id)
                
                # 计算归一化分数(百分制)
                comp_total = comp_result['total_score']
                comp_max = comp_result['max_score']
                normalized_score = round(comp_total / comp_max * 100) if comp_max > 0 else 0
                
                # 重新计算等级
                from app.services.scoring_rules import get_level
                level_key, level_info = get_level(normalized_score)
                
                # 生成 category_summary(按category分组统计算
                comp_details = comp_result.get('details', [])
                cat_summary = {}
                for d in comp_details:
                    cat = d.get('category', '答案对比')
                    if cat not in cat_summary:
                        cat_summary[cat] = {"passed": 0, "total": 0, "score": 0, "max_score": 0}
                    cat_summary[cat]["total"] += 1
                    if d.get('passed'):
                        cat_summary[cat]["passed"] += 1
                        cat_summary[cat]["score"] += d.get('score', 0)
                    cat_summary[cat]["max_score"] += d.get('score', 0)
                
                # 用答案对比结果完全替换评分结?
                result['total_score'] = normalized_score
                result['max_score'] = 100
                result['percentage'] = round(normalized_score / 100 * 100, 1)
                result['level'] = level_key
                result['level_name'] = level_info['name']
                result['level_color'] = level_info['color']
                result['level_desc'] = level_info['description']
                result['details'] = comp_details
                result['category_summary'] = cat_summary
                result['scoring_method'] = comp_result.get('scoring_method', 'answer_compare')
            except Exception as e:
                # 对比评分失败,保留规则评分结?
                import traceback
                traceback.print_exc()
                result['scoring_method'] = 'rule_check'

        # 保存提交记录
        user_id = user.get('username', '') if user else ''

        submission = {
            "submission_id": submission_id,
            "task_id": task_id,
            "student_name": student_name,
            "student_class": student_class,
            "user_id": user_id,
            "filename": filename,
            "filepath": filepath,
            "score_result": result,
            "submitted_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        }
        submissions[submission_id] = submission

        # 同一用户同一任务只保留最新提交,删除旧记录
        if user_id:
            for existing in submission_service.get_submissions_by_user(user_id):
                if existing.get('task_id') == task_id and existing.get('submission_id') != submission_id:
                    submission_service.delete_submission(existing['submission_id'])

        # 保存结果到数据库
        submission_service.save_submission(submission)

        return jsonify({"success": True, "submission_id": submission_id, "result": result})

    except Exception as e:
        return jsonify({"success": False, "message": f"评分出错: {str(e)}"}), 500


@api_bp.route('/result/<submission_id>', methods=['GET'])
@api_login_required
def get_result(submission_id):
    """获取评分结果"""
    # 先从内存查找
    if submission_id in submissions:
        return jsonify({"success": True, "data": submissions[submission_id]})

    # 从数据库查询
    data = submission_service.get_submission_full_data(submission_id)
    if data:
        submissions[submission_id] = data
        return jsonify({"success": True, "data": data})

    return jsonify({"success": False, "message": "未找到该提交记录"}), 404


# 异步任务状态存储
_task_status = {}  # {submission_id: {"status": "pending"/"done"/"error", "result": ..., "error": ...}}

def _run_ai_generation(app_ref, submission_id, submission, api_key, api_url, model):
    """后台线程执行AI生成"""
    try:
        from app.services.ai_service import DeepSeekService
        ai_service = DeepSeekService(api_key=api_key, api_url=api_url, model=model)
        result = submission["score_result"]

        tasks_result = ai_service.generate_learning_tasks(
            result["task_id"], result["total_score"], result["details"]
        )
        advice = tasks_result.get('advice', '')
        if not advice:
            advice = ai_service.generate_study_advice(
                result["task_id"], result["total_score"], result["details"]
            )

        submission["learning_tasks"] = tasks_result
        submission["study_advice"] = advice

        # 保存结果到数据库
        with app_ref.app_context():
            submission_service.save_submission(submission)

        _task_status[submission_id] = {"status": "done", "result": {
            "learning_tasks": tasks_result, "study_advice": _clean_advice(advice)
        }}
    except Exception as e:
        import traceback
        traceback.print_exc()
        _task_status[submission_id] = {"status": "error", "error": str(e)}


@api_bp.route('/generate-tasks', methods=['POST'])
@api_login_required
def generate_learning_tasks():
    """AI生成分层学习任务(异步)"""
    try:
        data = request.get_json()
        submission_id = data.get('submission_id')

        # 从内存或文件加载提交记录
        if submission_id and submission_id in submissions:
            submission = submissions[submission_id]
        elif submission_id:
            # 从数据库加载提交记录
            submission = submission_service.get_submission_full_data(submission_id)
            if submission:
                submissions[submission_id] = submission
            else:
                return jsonify({"success": False, "message": "未找到提交记录"}), 404
        else:
            return jsonify({"success": False, "message": "缺少提交ID"}), 400

        # 如果已经有结果,直接返回
        if submission.get("learning_tasks"):
            return jsonify({
                "success": True,
                "learning_tasks": submission["learning_tasks"],
                "study_advice": _clean_advice(submission.get("study_advice", "")),
            })

        # 如果正在生成中,返回pending
        if _task_status.get(submission_id, {}).get("status") == "pending":
            return jsonify({"success": True, "status": "pending", "message": "正在生成?.."})

        # 启动后台线程生成
        api_key = current_app.config.get('DEEPSEEK_API_KEY', '')
        api_url = current_app.config.get('DEEPSEEK_API_URL', '')
        model = current_app.config.get('DEEPSEEK_MODEL', '')

        _task_status[submission_id] = {"status": "pending"}
        import threading
        app_ref = current_app._get_current_object()
        t = threading.Thread(target=_run_ai_generation, args=(app_ref, submission_id, submission, api_key, api_url, model))
        t.daemon = True
        t.start()

        return jsonify({"success": True, "status": "pending", "message": "AI正在生成,请稍?.."})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "message": f"生成任务出错: {str(e)}"}), 500


@api_bp.route('/generate-tasks/status', methods=['POST'])
@api_login_required
def generate_tasks_status():
    """查询AI生成任务状态"""
    data = request.get_json() or {}
    submission_id = data.get('submission_id')
    if not submission_id:
        return jsonify({"success": False, "message": "缺少提交ID"}), 400

    status_info = _task_status.get(submission_id, {})
    status = status_info.get("status", "")

    if status == "done":
        result = status_info.get("result", {})
        return jsonify({"success": True, "status": "done", **result})
    elif status == "error":
        return jsonify({"success": False, "status": "error", "message": status_info.get("error", "生成失败")})
    else:
        return jsonify({"success": True, "status": "pending", "message": "正在生成?.."})


@api_bp.route('/submission/<submission_id>/download', methods=['GET'])
@api_login_required
def download_single_submission(submission_id):
    """下载单个作业文件"""
    upload_dir = current_app.config['UPLOAD_FOLDER']

    # 从数据库获取提交记录
    data = submission_service.get_submission_full_data(submission_id)
    if data:
        filepath = data.get('filepath', '') or data.get('file_path', '')
        if filepath and os.path.exists(filepath):
            # 从filepath中提取原始文件名
            import os as _os
            filename = _os.path.basename(filepath)
            # 去掉submission_id前缀
            if '_' in filename:
                parts = filename.split('_', 1)
                if len(parts[0]) == 8:  # submission_id?位hex
                    filename = parts[1]
            from flask import send_file
            return send_file(filepath, as_attachment=True, download_name=filename)

    # 从上传目录查询
    for ext in ['.xlsx', '.docx', '.pptx']:
        filepath = os.path.join(upload_dir, f"{submission_id}{ext}")
        if os.path.exists(filepath):
            from flask import send_file
            return send_file(filepath, as_attachment=True, download_name=f"{submission_id}{ext}")

    return jsonify({"success": False, "message": "文件不存在"}), 404


@api_bp.route('/generate-practice-from-history', methods=['POST'])
@api_login_required
def generate_practice_from_history():
    """根据历史作业评分结果生成针对性练习"""
    data = request.get_json()
    task_id = data.get('task_id', '')
    submission_id = data.get('submission_id', '')

    if not task_id:
        return jsonify({"success": False, "message": "缺少任务ID"}), 400

    # 加载历史评分结果作为上下
    score_context = ""
    if submission_id:
        # 从数据库获取
        result_data = submission_service.get_submission_full_data(submission_id)
        if result_data:
            score_result = result_data.get('score_result', {})
            total_score = score_result.get('total_score', 0)
            details = score_result.get('details', [])
            failed_items = [d for d in details if not d.get('passed', True)]
            if failed_items:
                score_context = f"学生本次得分{total_score}分,以下检查项未通过:"
                for item in failed_items:
                    score_context += f"\n- {item.get('name', '')}: {item.get('message', '')}"

    # 复用现有的AI生成逻辑
    from app.services.ai_service import DeepSeekService
    ai_service = DeepSeekService()

    # 启动后台生成线程
    import threading
    gen_task_id = str(uuid.uuid4())[:8]

    def run_generation():
        try:
            from app.services import document_generator
            # 获取学生得分作为级别参?
            score_level = 'pass'
            if submission_id:
                # 从数据库获取提交记录
                result_data = submission_service.get_submission_full_data(submission_id)
                if result_data:
                    score_result = result_data.get('score_result', {})
                    total_score = score_result.get('total_score', 0)
                    if total_score >= 90:
                        score_level = 'excellent'
                    elif total_score >= 75:
                        score_level = 'good'
                    elif total_score >= 60:
                        score_level = 'pass'
                    else:
                        score_level = 'fail'

            tasks_result = ai_service.generate_learning_tasks(task_id, 60, [])
            advice_result = ai_service.generate_study_advice(task_id)

            task_info = ALL_TASKS.get(task_id, {})
            student_info = {'score_context': score_context, 'level': score_level}

            package_data = document_generator.generate_practice_package(
                task_id=task_id,
                level_key=score_level,
                tasks=tasks_result,
                student_info=student_info
            )

            # 保存生成结果
            _task_status[gen_task_id] = {
                'status': 'done',
                'download_url': f'/api/practice/download/{package_data["filename"]}',
                'filename': package_data['filename']
            }
        except Exception as e:
            _task_status[gen_task_id] = {
                'status': 'error',
                'error': str(e)
            }
            import traceback
            traceback.print_exc()

    _task_status[gen_task_id] = {'status': 'pending'}
    thread = threading.Thread(target=run_generation, daemon=True)
    thread.start()

    return jsonify({"success": True, "task_id": gen_task_id})


@api_bp.route('/submissions', methods=['GET'])
@api_teacher_required
def get_submissions():
    """获取所有提交记录(教师后台用)- 排除教师提交,增加学生姓名班级信息"""
    # 从数据库获取所有提交记录
    all_submissions_raw = submission_service.get_all_submissions_full()
    all_submissions = []

    for data in all_submissions_raw:
        # 排除教师提交
        if _is_teacher_submission(data):
            continue
        # 补充学生姓名和班级信息
        data = _enrich_submission_with_student_info(data)
        all_submissions.append(data)

    # 按时间倒序
    all_submissions.sort(key=lambda x: x.get('submitted_at', ''), reverse=True)

    return jsonify({"success": True, "submissions": all_submissions, "total": len(all_submissions)})


@api_bp.route('/submissions/clear', methods=['POST'])
@api_teacher_required
def clear_submissions():
    """清空所有提交记录(教师后台用)"""
    import shutil
    upload_dir = current_app.config['UPLOAD_FOLDER']

    count = 0
    # 清空数据库中的提交记录
    count = submission_service.get_submission_count()
    submission_service.delete_all_submissions()

    # 清空上传文件
    if os.path.exists(upload_dir):
        for f in os.listdir(upload_dir):
            if f.endswith(('.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt')):
                os.remove(os.path.join(upload_dir, f))
                count += 1

    return jsonify({"success": True, "message": f"已清空{count}个文件"})


@api_bp.route('/my-submissions', methods=['GET'])
@api_login_required
def get_my_submissions():
    """获取当前学生的提交记录"""
    from app.services.auth_db import get_current_user
    user = get_current_user()
    if not user:
        return jsonify({"success": False, "message": "请先登录"}), 401

    student_name = user.get('name', '')
    user_id = user.get('username', '')
    # 从数据库获取当前用户的提交记录
    all_submissions = submission_service.get_all_submissions_full()
    my_submissions = []

    for data in all_submissions:
        # 优先按user_id匹配,其次按student_name匹配
        if data.get('user_id') == user_id or data.get('student_name') == student_name:
            my_submissions.append(data)

    my_submissions.sort(key=lambda x: x.get('submitted_at', ''), reverse=True)
    return jsonify({"success": True, "submissions": my_submissions, "total": len(my_submissions)})


@api_bp.route('/stats', methods=['GET'])
@api_teacher_required
def get_stats():
    """获取统计信息 - 排除教师提交"""
    # 从数据库获取所有提交记录
    stats = {
        "total_submissions": 0,
        "average_score": 0,
        "level_distribution": {"excellent": 0, "good": 0, "pass": 0, "fail": 0},
        "task_distribution": {},
    }

    all_submissions = submission_service.get_all_submissions_full()
    total_score = 0
    count = 0
    for data in all_submissions:
        # 排除教师提交
        if _is_teacher_submission(data):
            continue

        stats["total_submissions"] += 1
        score = data.get("score_result", {}).get("total_score", 0)
        total_score += score
        count += 1

        level = data.get("score_result", {}).get("level", "fail")
        if level in stats["level_distribution"]:
            stats["level_distribution"][level] += 1

        task_id = data.get("task_id", "unknown")
        task_name = data.get("score_result", {}).get("task_name", task_id)
        if task_name not in stats["task_distribution"]:
            stats["task_distribution"][task_name] = 0
        stats["task_distribution"][task_name] += 1

        if count > 0:
            stats["average_score"] = round(total_score / count, 1)

    return jsonify({"success": True, "stats": stats})


@api_bp.route('/analysis', methods=['GET'])
@api_teacher_required
def get_analysis():
    """获取学情分析数据 - 排除教师提交"""
    analysis = {
        "task_avg_scores": [],       # 各任务平均分
        "task_pass_rates": [],       # 各任务通过率
        "weak_skills": [],           # 薄弱技能点排行(通过率最低的5个检查项)
        "score_distribution": [],    # 学生成绩分布直方图
        "teaching_suggestions": [],  # 分层教学建议
    }

    # 从数据库获取所有提交记录
    all_submissions_raw = submission_service.get_all_submissions_full()
    all_submissions = []
    for data in all_submissions_raw:
        # 排除教师提交
        if _is_teacher_submission(data):
            continue
        all_submissions.append(data)

    if not all_submissions:
        return jsonify({"success": True, "analysis": analysis})

    # ---- a) 各任务平均分对比 ----
    task_scores = defaultdict(list)
    for sub in all_submissions:
        task_id = sub.get("task_id", "unknown")
        task_name = sub.get("score_result", {}).get("task_name", task_id)
        max_score = sub.get("score_result", {}).get("max_score", 100)
        score = sub.get("score_result", {}).get("total_score", 0)
        task_scores[task_name].append({"score": score, "max_score": max_score, "task_id": task_id})

    for task_name, records in task_scores.items():
        avg = round(sum(r["score"] for r in records) / len(records), 1)
        max_score = records[0]["max_score"]
        analysis["task_avg_scores"].append({
            "task_name": task_name,
            "avg_score": avg,
            "max_score": max_score,
            "count": len(records),
        })

    # ---- b) 各任务通过?----
    for task_name, records in task_scores.items():
        passed = sum(1 for r in records if r["score"] >= 60)
        rate = round(passed / len(records) * 100, 1)
        analysis["task_pass_rates"].append({
            "task_name": task_name,
            "pass_rate": rate,
            "passed": passed,
            "total": len(records),
        })

    # ---- c) 薄弱技能点排行(通过率最低的5个检查项) ----
    item_stats = defaultdict(lambda: {"name": "", "total": 0, "pass_count": 0})
    for sub in all_submissions:
        details = sub.get("score_result", {}).get("details", [])
        for item in details:
            item_id = item.get("id", "")
            item_name = item.get("name", item_id)
            if item_id:
                item_stats[item_id]["name"] = item_name
                item_stats[item_id]["total"] += 1
                if item.get("passed"):
                    item_stats[item_id]["pass_count"] += 1

    # 计算通过率并排序(通过率越低越薄弱)
    item_rates = []
    for item_id, stats in item_stats.items():
        if stats["total"] >= 3:  # 至少3人次才参与排行
            pass_rate = round(stats["pass_count"] / stats["total"] * 100, 1)
            item_rates.append({
                "item_id": item_id,
                "item_name": stats["name"],
                "rate": pass_rate,
                "pass_count": stats["pass_count"],
                "fail_count": stats["total"] - stats["pass_count"],
                "count": stats["total"],
            })

    item_rates.sort(key=lambda x: x["rate"])
    analysis["weak_skills"] = item_rates[:5]

    # ---- d) 学生成绩分布直方?----
    # 按学生汇总平均分(同一学生取最新一次提交或所有提交的平均分)
    student_scores = defaultdict(list)
    for sub in all_submissions:
        student_name = sub.get("student_name", "匿名")
        score = sub.get("score_result", {}).get("total_score", 0)
        student_scores[student_name].append(score)

    # 每个学生的平均分
    student_avg = []
    for name, scores in student_scores.items():
        avg = round(sum(scores) / len(scores), 1)
        student_avg.append(avg)

    # 分段统计
    bins = [
        {"label": "0-19", "min": 0, "max": 19, "count": 0},
        {"label": "20-39", "min": 20, "max": 39, "count": 0},
        {"label": "40-59", "min": 40, "max": 59, "count": 0},
        {"label": "60-69", "min": 60, "max": 69, "count": 0},
        {"label": "70-79", "min": 70, "max": 79, "count": 0},
        {"label": "80-89", "min": 80, "max": 89, "count": 0},
        {"label": "90-100", "min": 90, "max": 100, "count": 0},
    ]
    for avg in student_avg:
        for b in bins:
            if b["min"] <= avg <= b["max"]:
                b["count"] += 1
                break

    analysis["score_distribution"] = bins

    # ---- e) 分层教学建议 ----
    suggestions = []

    # 整体平均?
    all_scores = [sub.get("score_result", {}).get("total_score", 0) for sub in all_submissions]
    overall_avg = round(sum(all_scores) / len(all_scores), 1)
    overall_pass_rate = round(sum(1 for s in all_scores if s >= 60) / len(all_scores) * 100, 1)
    overall_excellent_rate = round(sum(1 for s in all_scores if s >= 85) / len(all_scores) * 100, 1)

    if overall_avg < 60:
        suggestions.append({
            "level": "urgent",
            "title": "整体基础薄弱,需加强基础教学",
            "content": f"全班平均分仅{overall_avg}分,通过率{overall_pass_rate}%.建议回归教材基础操作,放慢教学进度,增加课堂实操练习时间,确保每位同学掌握基本操作."
        })
    elif overall_avg < 75:
        suggestions.append({
            "level": "warning",
            "title": "整体水平中等,需巩固提升",
            "content": f"全班平均分{overall_avg}分,通过率{overall_pass_rate}%.大部分同学掌握了基本操作,但综合运用能力不足.建议增加综合性练习任务,强化操作熟练度."
        })
    else:
        suggestions.append({
            "level": "good",
            "title": "整体水平良好,可适当拓展",
            "content": f"全班平均分{overall_avg}分,通过率{overall_pass_rate}%,优秀率{overall_excellent_rate}%.大部分同学掌握较好,可适当增加拓展任务和创意设计类作业."
        })

    # 薄弱技能点建议
    if analysis["weak_skills"]:
        weak_names = [s["item_name"] for s in analysis["weak_skills"][:3]]
        suggestions.append({
            "level": "warning",
            "title": "重点关注薄弱技能点",
            "content": f"以下技能点得分率最低,建议在课堂上演示并重点练习:{', '.join(weak_names)}.可以设计针对性的专项练习,帮助同学逐一突破."
        })

    # 模块差异建议
    module_avgs = defaultdict(list)
    for sub in all_submissions:
        task_id = sub.get("task_id", "")
        score = sub.get("score_result", {}).get("total_score", 0)
        if task_id in ALL_TASKS:
            module = ALL_TASKS[task_id].get("module_group", "unknown")
            module_avgs[module].append(score)

    module_results = {}
    for module, scores in module_avgs.items():
        module_results[module] = round(sum(scores) / len(scores), 1)

    if len(module_results) >= 2:
        weakest_module = min(module_results, key=module_results.get)
        weakest_name = MODULE_GROUPS.get(weakest_module, {}).get("name", weakest_module)
        module_scores_str = ', '.join(
            '{} {}分'.format(MODULE_GROUPS.get(m, {}).get('name', m), v)
            for m, v in module_results.items()
        )
        suggestions.append({
            "level": "info",
            "title": f"{weakest_name}模块相对薄弱",
            "content": f"各模块平均分:{module_scores_str}.{weakest_name}模块得分最低,建议增加该模块的教学课时和练习机会."
        })

    # 两极分化建议
    if len(student_avg) >= 5:
        high_count = sum(1 for a in student_avg if a >= 85)
        low_count = sum(1 for a in student_avg if a < 60)
        if high_count > 0 and low_count > 0:
            gap = round(max(student_avg) - min(student_avg), 1)
            suggestions.append({
                "level": "info",
                "title": "成绩存在两极分化现象",
                "content": f"最高分{max(student_avg)}分,最低分{min(student_avg)}分,分差{gap}分.优秀{high_count}人,需努力{low_count}人.建议实施分层教学:为优秀学生设计拓展任务,为基础薄弱学生安排补弱练习和一对一辅导."
            })

    analysis["teaching_suggestions"] = suggestions

    return jsonify({"success": True, "analysis": analysis})


# ============================================================
# 作业管理 API
# ============================================================

@api_bp.route('/homework-management/overview', methods=['GET'])
@api_teacher_required
def homework_overview():
    """作业管理概览:按任务分组统计,增加按班级已交/待交统计"""
    class_filter = request.args.get('class_name', '').strip()

    # 获取所有学生信息
    all_students = _get_all_students_for_stats()

    # 从数据库获取所有提交记录(排除教师提交,补充学生信息)
    all_submissions_raw = submission_service.get_all_submissions_full()
    all_submissions = []
    for data in all_submissions_raw:
        # 排除教师提交
        if _is_teacher_submission(data):
            continue
        # 补充学生姓名和班级信息
        data = _enrich_submission_with_student_info(data)
        all_submissions.append(data)

    # 按班级筛选
    if class_filter:
        all_submissions = [s for s in all_submissions if s.get('student_class') == class_filter]

    # 按任务分组统计
    task_stats = {}
    for sub in all_submissions:
        tid = sub.get('task_id', '')
        if tid not in task_stats:
            task_info = ALL_TASKS.get(tid, {})
            task_stats[tid] = {
                'task_id': tid,
                'task_name': task_info.get('task_name', tid),
                'module': task_info.get('module', ''),
                'file_type': task_info.get('file_type', ''),
                'submissions': [],
                'total_score': 0,
                'count': 0,
            }
        task_stats[tid]['submissions'].append(sub)
        score = (sub.get('score_result') or {}).get('total_score', 0)
        task_stats[tid]['total_score'] += score
        task_stats[tid]['count'] += 1

    # 计算平均分和班级统计
    result_list = []
    for tid, stats in task_stats.items():
        avg_score = round(stats['total_score'] / stats['count']) if stats['count'] > 0 else 0

        # 按班级统计已待交(班级名称从学生基本信息获取消
        class_stats = {}
        for sub in stats['submissions']:
            # 班级名称优先从学生基本信息获?
            user_id = sub.get('user_id', '')
            students_index = _get_students_index()
            student_info = students_index.get(user_id, {})
            cls = student_info.get('class_name') or sub.get('student_class', '未分班')
            
            if cls not in class_stats:
                class_stats[cls] = {
                    'submitted_students': [],  # 已提交学生列表(含姓名)
                    'not_submitted_students': []  # 待提交学生列表(含姓名)
                }
            # 学生姓名优先从学生基本信息获?
            student_name = student_info.get('name') or sub.get('student_name', sub.get('user_id', '未知'))
            class_stats[cls]['submitted_students'].append({
                'name': student_name,
                'user_id': user_id
            })

        # 班级列表完全从学生基本信息获?
        all_classes_from_students = set(s.get('class_name', '') for s in all_students if s.get('class_name'))

        # 计算待交学生(从学生基本信息获取消
        for cls in all_classes_from_students:
            if cls not in class_stats:
                class_stats[cls] = {'submitted_students': [], 'not_submitted_students': []}
            submitted_user_ids = set(s['user_id'] for s in class_stats[cls]['submitted_students'])
            # 从all_students获取该班级的学生
            cls_students = [s for s in all_students if s.get('class_name') == cls]
            for student in cls_students:
                # 使用学号匹配(student_no 对应提交记录中的 user_id?
                student_user_id = student.get('student_no', '') or student.get('id', '')
                student_name = student.get('name', student_user_id)
                if student_user_id and student_user_id not in submitted_user_ids:
                    class_stats[cls]['not_submitted_students'].append({
                        'name': student_name,
                        'user_id': student_user_id
                    })

        # 转换为要求的格式
        class_stats_list = []
        for cls, cls_data in class_stats.items():
            total = len(cls_data['submitted_students']) + len(cls_data['not_submitted_students'])
            class_stats_list.append({
                'class_name': cls,
                'total': total,
                'submitted': len(cls_data['submitted_students']),
                'not_submitted': len(cls_data['not_submitted_students']),
                'submitted_students': cls_data['submitted_students'],  # 包含姓名和user_id
                'not_submitted_students': cls_data['not_submitted_students']  # 包含姓名和user_id
            })

        result_list.append({
            'task_id': tid,
            'task_name': stats['task_name'],
            'module': stats['module'],
            'file_type': stats['file_type'],
            'avg_score': avg_score,
            'submitted_count': stats['count'],
            'submission_ids': [s['submission_id'] for s in stats['submissions']],
            'class_stats': class_stats_list
        })

    # 按任务序号排序(从task_id中提取数字进行排序)
    def extract_task_number(tid):
        """从task_id中提取数字用于排序,?task_5_4_1 -> (5, 4, 1)"""
        import re
        numbers = re.findall(r'\d+', tid)
        return tuple(int(n) for n in numbers) if numbers else (0,)
    
    result_list.sort(key=lambda x: extract_task_number(x['task_id']))

    # 获取班级列表(唯一来源:学生基本信息)
    all_classes = sorted(set(s.get('class_name', '') for s in all_students if s.get('class_name')))

    return jsonify({
        'success': True,
        'tasks': result_list,
        'classes': all_classes,
        'total_students': len(all_students)
    })


@api_bp.route('/homework-management/detail', methods=['GET'])
@api_teacher_required
def homework_detail():
    """某个任务的提交详情 - 排除教师提交,增加学生姓名班级信息"""
    task_id = request.args.get('task_id', '')
    class_filter = request.args.get('class_name', '').strip()

    if not task_id:
        return jsonify({"success": False, "message": "缺少task_id"}), 400

    # 从数据库获取提交记录
    all_submissions = submission_service.get_all_submissions_full()
    submissions = []
    for data in all_submissions:
        if data.get('task_id') == task_id:
            # 排除教师提交
            if _is_teacher_submission(data):
                continue
            # 补充学生姓名和班级信息
            data = _enrich_submission_with_student_info(data)
            if not class_filter or data.get('student_class') == class_filter:
                submissions.append(data)

    submissions.sort(key=lambda x: x.get('submitted_at', ''), reverse=True)

    task_info = ALL_TASKS.get(task_id, {})
    avg_score = 0
    if submissions:
        avg_score = round(sum((s.get('score_result') or {}).get('total_score', 0) for s in submissions) / len(submissions))

    return jsonify({
        'success': True,
        'task_id': task_id,
        'task_name': task_info.get('task_name', task_id),
        'submissions': submissions,
        'avg_score': avg_score,
        'total': len(submissions)
    })


@api_bp.route('/homework-management/download/<submission_id>', methods=['GET'])
@api_teacher_required
def download_student_homework(submission_id):
    """下载学生提交的作业文件"""
    # 从数据库获取提交记录
    data = submission_service.get_submission_full_data(submission_id)
    if not data:
        return jsonify({"success": False, "message": "记录不存在"}), 404

    filepath = data.get('filepath', '')
    filename = data.get('filename', f'{submission_id}.docx')

    if not filepath or not os.path.exists(filepath):
        return jsonify({"success": False, "message": "文件不存在"}), 404

    from flask import send_file
    return send_file(filepath, as_attachment=True, download_name=filename)


# ============================================================
# 学生管理 API
# ============================================================

@api_bp.route('/students', methods=['GET'])
@api_teacher_required
def get_students():
    """获取学生列表"""
    keyword = request.args.get('keyword', '').strip()
    class_name = request.args.get('class_name', '').strip()

    if keyword:
        students = student_service.search_students(keyword)
    elif class_name:
        students = student_service.get_students_by_class(class_name)
    else:
        students = student_service.get_all_students()

    stats = student_service.get_student_stats()
    return jsonify({
        "success": True,
        "students": students,
        "stats": stats,
    })


@api_bp.route('/students/classes', methods=['GET'])
@api_teacher_required
def get_classes():
    """获取班级列表"""
    classes = student_service.get_class_list()
    return jsonify({"success": True, "classes": classes})


@api_bp.route('/students', methods=['POST'])
@api_teacher_required
def add_student():
    """添加单个学生"""
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "请求数据为空"}), 400

    name = data.get('name', '').strip()
    student_no = data.get('student_no', '').strip()
    class_name = data.get('class_name', '').strip()

    if not name:
        return jsonify({"success": False, "message": "姓名不能为空"}), 400

    student = student_service.add_student(name, student_no, class_name)
    return jsonify({"success": True, "student": student, "message": f"学生 {name} 添加成功"})


@api_bp.route('/students/<student_id>', methods=['PUT'])
@api_teacher_required
def update_student(student_id):
    """更新学生信息"""
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "请求数据为空"}), 400

    student = student_service.update_student(student_id, **data)
    if student:
        return jsonify({"success": True, "student": student, "message": "更新成功"})
    return jsonify({"success": False, "message": "学生不存在"}), 404


@api_bp.route('/students/<student_id>', methods=['DELETE'])
@api_teacher_required
def delete_student(student_id):
    """删除单个学生"""
    if student_service.delete_student(student_id):
        return jsonify({"success": True, "message": "删除成功"})
    return jsonify({"success": False, "message": "学生不存在"}), 404


@api_bp.route('/students/batch-delete', methods=['POST'])
@api_teacher_required
def batch_delete_students():
    """批量删除学生"""
    data = request.get_json()
    ids = data.get('ids', []) if data else []
    if not ids:
        return jsonify({"success": False, "message": "未选择学生"}), 400

    deleted = student_service.batch_delete_students(ids)
    return jsonify({"success": True, "message": f"已删除 {deleted} 名学生", "deleted": deleted})


@api_bp.route('/students/import', methods=['POST'])
@api_teacher_required
def import_students():
    """导入学生名单(支持CSV和Excel)"""
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "请上传文件"}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({"success": False, "message": "文件名为空"}), 400

    # 保存临时文件
    ext = os.path.splitext(file.filename)[1].lower()
    temp_path = os.path.join(DATA_DIR, f'temp_import{ext}')
    os.makedirs(DATA_DIR, exist_ok=True)
    file.save(temp_path)

    try:
        if ext in ('.csv',):
            result = student_service.import_from_csv(temp_path)
        elif ext in ('.xlsx', '.xls'):
            result = student_service.import_from_excel(temp_path)
        else:
            return jsonify({"success": False, "message": "仅支持CSV和Excel文件"}), 400

        if result['errors'] and result['imported'] == 0:
            return jsonify({"success": False, "message": "导入失败", "detail": result}), 400

        msg = f"成功导入 {result['imported']} 名学生"
        if result['skipped'] > 0:
            msg += f",跳过 {result['skipped']} 条(学号重复或姓名为空)"
        if result['errors']:
            msg += f",{len(result['errors'])} 条错误。"

        return jsonify({"success": True, "message": msg, "detail": result})
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


@api_bp.route('/students/sync-accounts', methods=['POST'])
@api_teacher_required
def sync_accounts():
    """为已导入学生同步创建登录账号"""
    created = student_service.sync_user_accounts()
    return jsonify({"success": True, "message": f"已为 {created} 名学生同步登录账号(密码为学号后6位)"})


# ============================================================
# 分层练习文档生成 API
# ============================================================

@api_bp.route('/practice/generate', methods=['POST'])
@api_login_required
def generate_practice():
    """学生端:按需生成分层练习文档并下载"""
    data = request.get_json()
    submission_id = data.get('submission_id', '').strip()
    level_key = data.get('level', '').strip()

    if not submission_id:
        return jsonify({"success": False, "message": "缺少提交ID"}), 400

    # 获取提交记录
    sub = submissions.get(submission_id)
    if not sub:
        # 从数据库获取提交记录
        sub = submission_service.get_submission_full_data(submission_id)
        if not sub:
            return jsonify({"success": False, "message": "未找到提交记录"}), 404
        submissions[submission_id] = sub

    task_id = sub.get('task_id', '')
    score_result = sub.get('score_result', {})
    details = score_result.get('details', [])

    # 确定等级
    if not level_key:
        score = score_result.get('total_score', 0)
        level_key, _ = get_level(score)

    # 获取AI分层任务
    ai_service = DeepSeekService(
        api_key=current_app.config.get('DEEPSEEK_API_KEY', ''),
        api_url=current_app.config.get('DEEPSEEK_API_URL', ''),
        model=current_app.config.get('DEEPSEEK_MODEL', '')
    )

    ai_result = ai_service.generate_learning_tasks(task_id, score_result.get('total_score', 0), details)
    all_tasks = ai_result.get('tasks') or {}

    # 获取对应等级的任?
    level_tasks = all_tasks.get(level_key, [])
    if not level_tasks:
        # 尝试其他等级
        for key in ['excellent', 'good', 'pass', 'fail']:
            if all_tasks.get(key):
                level_tasks = all_tasks[key]
                break

    if not level_tasks:
        error_msg = ai_result.get('error', '') or ai_result.get('raw_response', '') or ''
        return jsonify({"success": False, "message": f"AI未生成有效的练习任务: {error_msg[:200]}"}), 500

    # 生成练习
    student_info = {
        'name': sub.get('student_name', ''),
        'class_name': sub.get('student_class', ''),
        'score': f"{score_result.get('total_score', 0)}/{score_result.get('max_score', 100)}",
        'filepath': sub.get('filepath', '')
    }

    try:
        zip_bytes, filename = document_generator.generate_practice_package(
            task_id, level_key, level_tasks, student_info
        )

        # 保存到磁?
        document_generator.save_generated_package(task_id, level_key, level_tasks, student_info)

        return jsonify({
            "success": True,
            "message": "练习文档生成成功",
            "filename": filename,
            "download_url": f"/api/practice/download?submission_id={submission_id}&level={level_key}"
        })
    except Exception as e:
        return jsonify({"success": False, "message": f"文档生成失败: {str(e)}"}), 500


@api_bp.route('/practice/download', methods=['GET', 'POST'])
@api_login_required
def download_practice():
    """下载练习文档"""
    if request.method == 'POST':
        data = request.get_json(silent=True) or {}
        submission_id = data.get('submission_id', '').strip()
        level_key = data.get('level', '').strip()
    else:
        submission_id = request.args.get('submission_id', '').strip()
        level_key = request.args.get('level', '').strip()

    if not submission_id:
        return jsonify({"success": False, "message": "缺少提交ID"}), 400

    # 获取提交记录
    sub = submissions.get(submission_id)
    if not sub:
        # 从数据库获取提交记录
        sub = submission_service.get_submission_full_data(submission_id)
        if not sub:
            return jsonify({"success": False, "message": "未找到提交记录"}), 404

    task_id = sub.get('task_id', '')
    score_result = sub.get('score_result', {})
    details = score_result.get('details', [])

    if not level_key:
        score = score_result.get('total_score', 0)
        level_key, _ = get_level(score)

    # 优先使用已有的learning_tasks(generateTasks已生成的?
    level_tasks = []
    cached_tasks = sub.get('learning_tasks', {})
    if cached_tasks:
        all_tasks = cached_tasks.get('tasks') or cached_tasks
        level_tasks = all_tasks.get(level_key, [])
        if not level_tasks:
            for key in ['excellent', 'good', 'pass', 'fail']:
                if all_tasks.get(key):
                    level_tasks = all_tasks[key]
                    break

    # 如果没有缓存,才调用AI
    if not level_tasks:
        ai_service = DeepSeekService(
            api_key=current_app.config.get('DEEPSEEK_API_KEY', ''),
            api_url=current_app.config.get('DEEPSEEK_API_URL', ''),
            model=current_app.config.get('DEEPSEEK_MODEL', '')
        )
        ai_result = ai_service.generate_learning_tasks(task_id, score_result.get('total_score', 0), details)
        all_tasks = ai_result.get('tasks') or {}
        level_tasks = all_tasks.get(level_key, [])
        if not level_tasks:
            for key in ['excellent', 'good', 'pass', 'fail']:
                if all_tasks.get(key):
                    level_tasks = all_tasks[key]
                    break

    if not level_tasks:
        return jsonify({"success": False, "message": "无有效练习任务,请先点击\"AI生成个性化作业\""}), 500

    student_info = {
        'name': sub.get('student_name', ''),
        'class_name': sub.get('student_class', ''),
        'score': f"{score_result.get('total_score', 0)}/{score_result.get('max_score', 100)}",
        'filepath': sub.get('filepath', '')
    }

    try:
        zip_bytes, filename = document_generator.generate_practice_package(
            task_id, level_key, level_tasks, student_info
        )

        from flask import send_file
        from urllib.parse import quote as url_quote
        encoded_name = url_quote(filename, safe='')
        response = send_file(
            io.BytesIO(zip_bytes),
            mimetype='application/zip',
            as_attachment=True,
            download_name=filename
        )
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_name}"
        return response
    except Exception as e:
        return jsonify({"success": False, "message": f"下载失败: {str(e)}"}), 500


# ============================================================
# 教师端:预生成练习文档管?API
# ============================================================

@api_bp.route('/admin/practice/generate', methods=['POST'])
@api_teacher_required
def admin_generate_practice():
    """教师端:为指定任务和等级预生成练习文档"""
    data = request.get_json()
    task_id = data.get('task_id', '').strip()
    level_key = data.get('level', '').strip()

    if not task_id or not level_key:
        return jsonify({"success": False, "message": "缺少任务ID或等级"}), 400

    task_info = ALL_TASKS.get(task_id)
    if not task_info:
        return jsonify({"success": False, "message": f"未找到任?{task_id}"}), 404

    # 使用AI生成分层任务
    ai_service = DeepSeekService(
        api_key=current_app.config.get('DEEPSEEK_API_KEY', ''),
        api_url=current_app.config.get('DEEPSEEK_API_URL', ''),
        model=current_app.config.get('DEEPSEEK_MODEL', '')
    )

    # 构造模拟的评分详情(全部通过,用于生成该等级的任务)
    mock_details = [{"name": item["name"], "score": item["score"], "passed": True, "message": "通过"}
                    for item in task_info.get('check_items', [])]

    ai_result = ai_service.generate_learning_tasks(task_id, 100, mock_details)
    all_tasks = ai_result.get('tasks') or {}
    level_tasks = all_tasks.get(level_key, [])

    if not level_tasks:
        error_msg = ai_result.get('error', '') or ai_result.get('raw_response', '') or '未知错误'
        return jsonify({"success": False, "message": f"AI未生成有效任? {error_msg[:200]}"}), 500

    try:
        filepath = document_generator.save_generated_package(task_id, level_key, level_tasks)
        filename = os.path.basename(filepath)
        return jsonify({
            "success": True,
            "message": f"练习文档已生成:{filename}",
            "filename": filename,
            "download_url": f"/api/admin/practice/download?filename={filename}"
        })
    except Exception as e:
        return jsonify({"success": False, "message": f"生成失败: {str(e)}"}), 500


@api_bp.route('/admin/practice/list', methods=['GET'])
@api_teacher_required
def admin_practice_list():
    """获取已生成的练习文档列表"""
    packages = document_generator.get_generated_packages()
    return jsonify({"success": True, "packages": packages})


@api_bp.route('/admin/practice/download', methods=['GET', 'POST'])
@api_teacher_required
def admin_practice_download():
    """下载已生成的练习文档"""
    # 支持GET和POST两种方式,POST方式避免中文URL编码问题
    if request.method == 'POST':
        data = request.get_json(silent=True) or {}
        filename = data.get('filename', '').strip()
    else:
        filename = request.args.get('filename', '').strip()
    if not filename:
        return jsonify({"success": False, "message": "缺少文件"}), 400

    # 安全检?
    if '..' in filename or '/' in filename:
        return jsonify({"success": False, "message": "非法文件"}), 400

    filepath = os.path.join(DATA_DIR, 'generated', filename)
    if not os.path.exists(filepath):
        return jsonify({"success": False, "message": "文件不存在"}), 404

    from flask import send_file
    from urllib.parse import quote as url_quote
    # 中文文件名需要RFC 5987编码
    encoded_name = url_quote(filename, safe='')
    response = send_file(filepath, mimetype='application/zip', as_attachment=True, download_name=filename)
    response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_name}"
    return response


@api_bp.route('/admin/practice/delete', methods=['POST'])
@api_teacher_required
def admin_practice_delete():
    """删除已生成的练习文档"""
    data = request.get_json()
    filename = data.get('filename', '').strip() if data else ''

    if not filename:
        return jsonify({"success": False, "message": "缺少文件"}), 400

    if document_generator.delete_generated_package(filename):
        return jsonify({"success": True, "message": "删除成功"})
    return jsonify({"success": False, "message": "文件不存在"}), 404


# ============================================================
# 系统配置 API
# ============================================================

@api_bp.route('/admin/config/ai', methods=['GET'])
@api_teacher_required
def get_ai_config_api():
    """获取AI接口配置"""
    from config.settings import get_ai_config
    cfg = get_ai_config()
    # 隐藏API Key中间部分
    key = cfg.get('api_key', '')
    if key and len(key) > 8:
        masked_key = key[:4] + '*' * (len(key) - 8) + key[-4:]
    else:
        masked_key = key
    return jsonify({
        "success": True,
        "config": {
            "api_key": masked_key,
            "api_key_set": bool(key),
            "api_url": cfg.get('api_url', ''),
            "model": cfg.get('model', ''),
        }
    })


@api_bp.route('/admin/config/ai', methods=['POST'])
@api_teacher_required
def save_ai_config_api():
    """保存AI接口配置"""
    from config.settings import get_ai_config, save_ai_config
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "请求数据为空"}), 400

    current = get_ai_config()
    api_key = data.get('api_key', '').strip()
    api_url = data.get('api_url', '').strip()
    model = data.get('model', '').strip()

    # 如果API Key没有变化(还是掩码形式),保留原来的
    if api_key and '*' in api_key:
        api_key = current['api_key']

    if not api_url:
        api_url = current['api_url']
    if not model:
        model = current['model']

    save_ai_config(api_key, api_url, model)

    # 更新Flask运行时配?
    current_app.config['DEEPSEEK_API_KEY'] = api_key
    current_app.config['DEEPSEEK_API_URL'] = api_url
    current_app.config['DEEPSEEK_MODEL'] = model

    # 测试连接
    test_result = None
    if api_key:
        try:
            import requests
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            test_payload = {
                'model': model,
                'messages': [{'role': 'user', 'content': '你好'}],
                'max_tokens': 5
            }
            resp = requests.post(api_url, headers=headers, json=test_payload, timeout=10)
            if resp.status_code == 200:
                test_result = {'success': True, 'message': '连接成功'}
            else:
                test_result = {'success': False, 'message': f'连接失败: HTTP {resp.status_code}'}
        except Exception as e:
            test_result = {'success': False, 'message': f'连接失败: {str(e)}'}

    return jsonify({
        "success": True,
        "message": "配置已保存",
        "test_result": test_result
    })


@api_bp.route('/admin/config/ai/test', methods=['POST'])
@api_teacher_required
def test_ai_connection():
    """测试AI接口连接"""
    from config.settings import get_ai_config
    cfg = get_ai_config()
    api_key = cfg.get('api_key', '')

    if not api_key:
        return jsonify({"success": False, "message": "请先配置API Key"}), 400

    try:
        import requests
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        test_payload = {
            'model': cfg.get('model', 'deepseek-chat'),
            'messages': [{'role': 'user', 'content': '请回复连接成功'}],
            'max_tokens': 20
        }
        resp = requests.post(cfg.get('api_url', ''), headers=headers, json=test_payload, timeout=15)

        if resp.status_code == 200:
            result = resp.json()
            reply = result.get('choices', [{}])[0].get('message', {}).get('content', '')
            return jsonify({"success": True, "message": f"连接成功！模型回复:{reply}"})
        else:
            return jsonify({"success": False, "message": f"连接失败: HTTP {resp.status_code} - {resp.text[:200]}"})
    except Exception as e:
        return jsonify({"success": False, "message": f"连接失败: {str(e)}"})


# ==================== 认证相关API ====================

@api_bp.route('/auth/login', methods=['POST'])
def auth_login():
    """用户登录"""
    data = request.get_json()
    username = (data.get('username') or '').strip()
    password = (data.get('password') or '').strip()

    if not username or not password:
        return jsonify({"success": False, "message": "请输入工号和密码"})

    from app.services.auth_db import authenticate, login_user
    user = authenticate(username, password)
    if not user:
        return jsonify({"success": False, "message": "工号或密码错误"}), 401

    login_user(user)
    return jsonify({"success": True, "user": user})


@api_bp.route('/auth/logout', methods=['POST'])
def auth_logout():
    """用户登出"""
    from app.services.auth_db import logout_user
    logout_user()
    return jsonify({"success": True})


@api_bp.route('/auth/current', methods=['GET'])
def auth_current():
    """获取当前登录用户信息"""
    from app.services.auth_db import get_current_user, is_logged_in
    if is_logged_in():
        return jsonify({"success": True, "user": get_current_user()})
    return jsonify({"success": False, "message": "未登录"}), 401


@api_bp.route('/auth/change-password', methods=['POST'])
def auth_change_password():
    """修改密码"""
    from app.services.auth_db import get_current_user, change_password
    user = get_current_user()
    if not user:
        return jsonify({"success": False, "message": "请先登录"}), 401

    data = request.get_json()
    old_pwd = data.get('old_password', '')
    new_pwd = data.get('new_password', '')

    if not old_pwd or not new_pwd:
        return jsonify({"success": False, "message": "请输入原密码和新密码"})

    ok, msg = change_password(user['username'], old_pwd, new_pwd)
    return jsonify({"success": ok, "message": msg})


# ==================== 知识库管理API ====================

@api_bp.route('/knowledge-base', methods=['GET'])
@api_teacher_required
def get_knowledge_base_api():
    """获取知识库数据"""
    from app.services.knowledge_base import get_knowledge_base
    kb = get_knowledge_base()
    return jsonify({"success": True, "data": kb})


@api_bp.route('/knowledge-base/module/<module_key>', methods=['GET'])
@api_teacher_required
def get_kb_module(module_key):
    """获取单个模块的知识库"""
    from app.services.knowledge_base import get_knowledge_base
    kb = get_knowledge_base()
    module = kb.get(module_key)
    if not module:
        return jsonify({"success": False, "message": "模块不存在"}), 404
    return jsonify({"success": True, "data": module})


@api_bp.route('/knowledge-base/module/<module_key>', methods=['PUT'])
@api_teacher_required
def update_kb_module(module_key):
    """更新单个模块的知识库"""
    from app.services.knowledge_base import get_knowledge_base, update_knowledge_base
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "无数据"}), 400
    
    kb = get_knowledge_base()
    if module_key not in kb:
        return jsonify({"success": False, "message": "模块不存在"}), 404
    
    kb[module_key] = data
    update_knowledge_base(kb)
    return jsonify({"success": True, "message": "保存成功"})


@api_bp.route('/knowledge-base/task/<task_id>', methods=['GET'])
@api_teacher_required
def get_kb_task(task_id):
    """获取单个任务的知识点"""
    from app.services.knowledge_base import get_knowledge_base
    kb = get_knowledge_base()
    
    # 确定模块
    from app.services.scoring_rules import ALL_TASKS
    task_info = ALL_TASKS.get(task_id, {})
    module_group = task_info.get('module_group', '')
    
    module_kb = kb.get(module_group, {})
    task_kb = module_kb.get('task_knowledge', {}).get(task_id)
    if not task_kb:
        return jsonify({"success": False, "message": "未找到该任务的知识点"}), 404
    
    return jsonify({"success": True, "data": task_kb, "module": module_group})


@api_bp.route('/knowledge-base/task/<task_id>', methods=['PUT'])
@api_teacher_required
def update_kb_task(task_id):
    """更新单个任务的知识点"""
    from app.services.knowledge_base import get_knowledge_base, update_knowledge_base
    from app.services.scoring_rules import ALL_TASKS
    
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "无数据"}), 400
    
    task_info = ALL_TASKS.get(task_id, {})
    module_group = task_info.get('module_group', '')
    
    kb = get_knowledge_base()
    if module_group not in kb:
        return jsonify({"success": False, "message": "模块不存在"}), 404
    
    if 'task_knowledge' not in kb[module_group]:
        kb[module_group]['task_knowledge'] = {}
    
    kb[module_group]['task_knowledge'][task_id] = data
    update_knowledge_base(kb)
    return jsonify({"success": True, "message": "保存成功"})


@api_bp.route('/knowledge-base/reset', methods=['POST'])
@api_teacher_required
def reset_knowledge_base():
    """重置知识库为默认数据"""
    from app.services.knowledge_base import KNOWLEDGE_BASE, update_knowledge_base
    update_knowledge_base(KNOWLEDGE_BASE)
    return jsonify({"success": True, "message": "已重置为默认知识库"})


# ==================== 素材文件管理API ====================

MATERIALS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'data', 'materials')


def _get_materials_index():
    """读取素材文件索引"""
    index_path = os.path.join(MATERIALS_DIR, 'index.json')
    if os.path.exists(index_path):
        with open(index_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}


def _save_materials_index(index):
    """保存素材文件索引"""
    os.makedirs(MATERIALS_DIR, exist_ok=True)
    index_path = os.path.join(MATERIALS_DIR, 'index.json')
    with open(index_path, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)


def _extract_file_keywords(filepath):
    """从素材文件中提取关键词(用于匹配学生提交的文件)"""
    ext = os.path.splitext(filepath)[1].lower()
    keywords = []

    try:
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            ws = wb.active
            # 提取消行的所有文?
            for row in ws.iter_rows(min_row=1, max_row=min(3, ws.max_row), values_only=True):
                for val in row:
                    if val is not None:
                        text = str(val).strip()
                        if text and len(text) <= 20:
                            keywords.append(text)
            wb.close()
        elif ext == '.docx':
            from docx import Document
            doc = Document(filepath)
            # 提取消0个段落的文本
            for i, para in enumerate(doc.paragraphs[:10]):
                text = para.text.strip()
                if text and len(text) <= 50:
                    keywords.append(text)
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and len(text) <= 50:
                                keywords.append(text)
    except Exception:
        pass

    return keywords


@api_bp.route('/materials', methods=['GET'])
def get_materials():
    """获取所有任务的素材文件列表"""
    index = _get_materials_index()
    return jsonify({"success": True, "data": index})


@api_bp.route('/materials/<task_id>', methods=['GET'])
def get_task_material(task_id):
    """获取指定任务的素材文件信息"""
    index = _get_materials_index()
    task_materials = index.get(task_id, [])
    return jsonify({"success": True, "data": task_materials})


@api_bp.route('/materials/<task_id>', methods=['POST'])
@api_teacher_required
def upload_material(task_id):
    """教师上传素材文件"""
    if task_id not in ALL_TASKS:
        return jsonify({"success": False, "message": "无效的任务ID"}), 400

    file = request.files.get('file')
    if not file:
        return jsonify({"success": False, "message": "请选择文件"}), 400

    # 检查文件类
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    expected_ext = f".{ALL_TASKS[task_id]['file_type']}"
    if ext != expected_ext:
        return jsonify({"success": False, "message": f"文件格式错误,请上传 {expected_ext} 格式文件"}), 400

    # 保存文件
    os.makedirs(MATERIALS_DIR, exist_ok=True)
    material_id = str(uuid.uuid4())[:8]
    filepath = os.path.join(MATERIALS_DIR, f"{material_id}_{filename}")
    file.save(filepath)

    # 提取关键
    keywords = _extract_file_keywords(filepath)

    # 更新索引
    index = _get_materials_index()
    if task_id not in index:
        index[task_id] = []

    # 如果已有素材,先删除旧文?
    for old in index[task_id]:
        old_path = os.path.join(MATERIALS_DIR, old['filepath'])
        if os.path.exists(old_path):
            os.remove(old_path)

    index[task_id] = [{
        "material_id": material_id,
        "filename": filename,
        "filepath": f"{material_id}_{filename}",
        "keywords": keywords,
        "uploaded_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }]
    _save_materials_index(index)

    return jsonify({"success": True, "message": "素材文件上传成功", "keywords": keywords})


@api_bp.route('/materials/<task_id>', methods=['DELETE'])
@api_teacher_required
def delete_material(task_id):
    """删除指定任务的素材文件"""
    index = _get_materials_index()
    if task_id in index:
        for mat in index[task_id]:
            filepath = os.path.join(MATERIALS_DIR, mat['filepath'])
            if os.path.exists(filepath):
                os.remove(filepath)
        del index[task_id]
        _save_materials_index(index)
    return jsonify({"success": True, "message": "素材文件已删除"})


def get_material_keywords(task_id):
    """获取指定任务的素材关键词(供评分引擎使用)"""
    index = _get_materials_index()
    materials = index.get(task_id, [])
    if not materials:
        return None
    # 返回最新素材的关键
    return materials[-1].get('keywords', [])


# ========== 标准答案管理 API ==========

@api_bp.route('/answers', methods=['GET'])
def get_all_answers():
    """获取所有标准答案列表"""
    index = _get_answers_index()
    return jsonify({"success": True, "data": index})


@api_bp.route('/answers/<task_id>', methods=['GET'])
def get_task_answer(task_id):
    """获取指定任务的标准答案信息"""
    index = _get_answers_index()
    task_answers = index.get(task_id, [])
    return jsonify({"success": True, "data": task_answers})


@api_bp.route('/answers/<task_id>', methods=['POST'])
@api_teacher_required
def upload_answer(task_id):
    """教师上传标准答案文件"""
    if task_id not in ALL_TASKS:
        return jsonify({"success": False, "message": "无效的任务ID"}), 400

    file = request.files.get('file')
    if not file:
        return jsonify({"success": False, "message": "请选择文件"}), 400

    # 检查文件类
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    expected_ext = f".{ALL_TASKS[task_id]['file_type']}"
    if ext != expected_ext:
        return jsonify({"success": False, "message": f"文件格式错误,请上传 {expected_ext} 格式文件"}), 400

    # 保存文件
    os.makedirs(ANSWERS_DIR, exist_ok=True)
    answer_id = str(uuid.uuid4())[:8]
    filepath = os.path.join(ANSWERS_DIR, f"{answer_id}_{filename}")
    file.save(filepath)

    # 更新索引
    index = _get_answers_index()
    if task_id not in index:
        index[task_id] = []

    # 如果已有答案,先删除旧文?
    for old in index[task_id]:
        old_path = os.path.join(ANSWERS_DIR, old['filepath'])
        if os.path.exists(old_path):
            os.remove(old_path)

    index[task_id] = [{
        "answer_id": answer_id,
        "filename": filename,
        "filepath": f"{answer_id}_{filename}",
        "uploaded_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }]
    _save_answers_index(index)

    return jsonify({"success": True, "message": "标准答案上传成功"})


@api_bp.route('/answers/<task_id>', methods=['DELETE'])
@api_teacher_required
def delete_answer(task_id):
    """删除指定任务的标准答案"""
    index = _get_answers_index()
    if task_id in index:
        for ans in index[task_id]:
            filepath = os.path.join(ANSWERS_DIR, ans['filepath'])
            if os.path.exists(filepath):
                os.remove(filepath)
        del index[task_id]
        _save_answers_index(index)
    return jsonify({"success": True, "message": "标准答案已删除"})


def get_answer_file(task_id):
    """获取指定任务的标准答案文件路径(供评分引擎使用)"""
    index = _get_answers_index()
    answers = index.get(task_id, [])
    if not answers:
        return None
    # 返回最新答案的文件路径
    filepath = answers[-1].get('filepath')
    if filepath:
        return os.path.join(ANSWERS_DIR, filepath)
    return None


# ========== 教师管理 API ==========

@api_bp.route('/teachers', methods=['GET'])
@api_teacher_required
def get_teachers():
    """获取教师列表"""
    from app.services.auth_db import _load_users
    users = _load_users()
    teachers = []
    for uid, u in users.items():
        if u.get('role') in ('teacher', 'admin'):
            teachers.append({'id': uid, 'name': u.get('name', ''), 'role': u.get('role', '')})
    return jsonify({"success": True, "teachers": teachers})


@api_bp.route('/teachers', methods=['POST'])
@api_teacher_required
def add_teacher():
    """添加教师"""
    from app.services.auth_db import _load_users, _hash_password
    data = request.get_json()
    uid = data.get('id', '').strip()
    name = data.get('name', '').strip()
    password = data.get('password', '123456').strip()
    
    if not uid or not name:
        return jsonify({"success": False, "message": "工号和姓名不能为空"}), 400
    
    users = _load_users()
    if uid in users:
        return jsonify({"success": False, "message": "该工号已存在"}), 400
    
    users[uid] = {
        'password_hash': _hash_password(password),
        'name': name,
        'role': 'teacher',
        'title': '教师'
    }
    
    # 保存
    users_file = current_app.config.get('USERS_FILE') or os.path.join(
        os.path.dirname(os.path.dirname(__file__)), 'data', 'users.json')
    os.makedirs(os.path.dirname(users_file), exist_ok=True)
    with open(users_file, 'w', encoding='utf-8') as f:
        json.dump(users, f, ensure_ascii=False, indent=2)
    
    return jsonify({"success": True, "message": f"教师 {name} 添加成功"})


@api_bp.route('/teachers/<teacher_id>', methods=['DELETE'])
@api_teacher_required
def delete_teacher(teacher_id):
    """删除教师"""
    if teacher_id == 'admin':
        return jsonify({"success": False, "message": "不能删除管理员账号"}), 400
    
    from app.services.auth_db import _load_users
    users = _load_users()
    
    if teacher_id not in users:
        return jsonify({"success": False, "message": "教师不存在"}), 404
    
    # 不允许删除自己
    from app.services.auth_db import get_current_user
    current = get_current_user()
    if current and current.get('username') == teacher_id:
        return jsonify({"success": False, "message": "不能删除自己的账号"}), 400
    
    del users[teacher_id]
    
    users_file = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'data', 'users.json')
    with open(users_file, 'w', encoding='utf-8') as f:
        json.dump(users, f, ensure_ascii=False, indent=2)
    
    return jsonify({"success": True, "message": "删除成功"})


# ========== 角色检?API ==========

@api_bp.route('/auth/check-role', methods=['GET'])
@api_login_required
def check_role():
    """检查当前用户角色"""
    from app.services.auth_db import get_current_user
    user = get_current_user()
    is_admin = user and (user.get('username') == 'admin' or user.get('role') == 'admin')
    return jsonify({"success": True, "is_admin": is_admin})


# ========== 教案生成 API ==========

@api_bp.route('/generate-lesson-plan', methods=['GET'])
@api_teacher_required
def generate_lesson_plan():
    """根据薄弱技能点生成教案"""
    class_filter = request.args.get('class_name', '').strip()
    task_filter = request.args.get('task_id', '').strip()
    
    # 1. 分析薄弱技能点
    # 从数据库获取所有提交记录
    all_submissions = submission_service.get_all_submissions_full()
    for data in all_submissions:
                # 排除教师提交
                if _is_teacher_submission(data):
                    continue
                # 补充学生姓名和班级信息(从学生基本信息获取)
                data = _enrich_submission_with_student_info(data)
                all_submissions.append(data)
    
    # 筛选(班级信息已从学生基本信息获取消
    if class_filter:
        all_submissions = [s for s in all_submissions if s.get('student_class') == class_filter]
    if task_filter:
        all_submissions = [s for s in all_submissions if s.get('task_id') == task_filter]
    
    # 统计每个检查项的失败率
    point_stats = {}
    for sub in all_submissions:
        details = (sub.get('score_result') or {}).get('details', [])
        task_id = sub.get('task_id', '')
        task_name = (sub.get('score_result') or {}).get('task_name', task_id)
        for d in details:
            key = f"{task_id}:{d.get('name', '')}"
            if key not in point_stats:
                point_stats[key] = {
                    'name': d.get('name', ''),
                    'task_name': task_name,
                    'task_id': task_id,
                    'total': 0,
                    'fail_count': 0,
                }
            point_stats[key]['total'] += 1
            if not d.get('passed', True):
                point_stats[key]['fail_count'] += 1
    
    # 计算失败率并排序
    weak_points = []
    for key, ps in point_stats.items():
        if ps['total'] > 0:
            fail_rate = round(ps['fail_count'] / ps['total'] * 100)
            if fail_rate > 0:
                weak_points.append({
                    'name': ps['name'],
                    'task_name': ps['task_name'],
                    'fail_count': ps['fail_count'],
                    'total': ps['total'],
                    'fail_rate': fail_rate,
                })
    
    weak_points.sort(key=lambda x: x['fail_rate'], reverse=True)
    weak_points = weak_points[:10]  # 取前10?    
    if not weak_points:
        return jsonify({"success": True, "weak_points": [], "lesson_plan": "所有技能点掌握良好,无需生成针对性教案."})
    
    # 2. 调用AI生成教案
    try:
        from app.services.ai_service import DeepSeekService
        api_key = current_app.config.get('DEEPSEEK_API_KEY', '')
        api_url = current_app.config.get('DEEPSEEK_API_URL', '')
        model = current_app.config.get('DEEPSEEK_MODEL', '')
        ai_service = DeepSeekService(api_key=api_key, api_url=api_url, model=model)
        
        weak_points_text = '\n'.join([
            f"- {wp['name']}({wp['task_name']}):{wp['fail_rate']}%的学生未通过({wp['fail_count']}/{wp['total']}人)"
            for wp in weak_points
        ])
        
        prompt = f"""你是一位经验丰富的信息技术课程教师. 根据以下学生作业的薄弱技能点分析结果, 生成一份针对性教案.
薄弱技能点(按错误率从高到低排序): {weak_points_text}

请生成教案, 包含以下部分:
1. 教学目标(针对薄弱技能点)
2. 教学重点与难点
3. 教学过程(分步骤, 每步包含教师活动和学生活动)
4. 课堂练习设计
5. 教学反思建议
要求:
- 教案要具体, 可操作
- 每个薄弱技能点都要有对应的教学环节
- 适合课堂教学使用
- 语言简洁专业"""
        
        result = ai_service._call_api([{"role": "user", "content": prompt}])
        
        if result:
            # 生成docx文件
            from app.services.document_generator import generate_lesson_plan_docx
            docx_path = generate_lesson_plan_docx(result, weak_points)
            download_url = f'/api/lesson-plan/download/{os.path.basename(docx_path)}'
            
            return jsonify({
                "success": True,
                "weak_points": weak_points,
                "lesson_plan": result,
                "download_url": download_url
            })
        else:
            return jsonify({"success": False, "message": "AI生成失败,请检查API配置"})
    except Exception as e:
        return jsonify({"success": False, "message": f"生成失败: {str(e)}"})


@api_bp.route('/lesson-plan/download/<filename>', methods=['GET'])
@api_teacher_required
def download_lesson_plan(filename):
    """下载生成的教案文件"""
    from flask import send_file
    filepath = os.path.join(current_app.config['GENERATED_FOLDER'], filename)
    if not os.path.exists(filepath):
        return jsonify({"success": False, "message": "文件不存在"}), 404
    return send_file(filepath, as_attachment=True, download_name=filename)


# ========== 作业移动 API ==========

@api_bp.route('/submission/<submission_id>/move', methods=['POST'])
@api_teacher_required
def move_submission(submission_id):
    """移动作业到另一个任务"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "message": "请求数据为空"}), 400

        target_task_id = data.get('target_task_id', '').strip()
        if not target_task_id:
            return jsonify({"success": False, "message": "缺少目标任务ID"}), 400

        # 验证目标task_id存在
        if target_task_id not in ALL_TASKS:
            return jsonify({"success": False, "message": f"无效的目标任务ID: {target_task_id}"}), 400

        # 获取提交记录
        submission = None
        if submission_id in submissions:
            submission = submissions[submission_id]
        else:
            submission = submission_service.get_submission_full_data(submission_id)
            if not submission:
                return jsonify({"success": False, "message": "未找到该提交记录"}), 404
            submissions[submission_id] = submission

        # 获取原任务ID
        original_task_id = submission.get('task_id', '')
        if original_task_id == target_task_id:
            return jsonify({"success": False, "message": "目标任务与原任务相同"}), 400

        # 获取文件路径
        filepath = submission.get('filepath', '')
        if not filepath or not os.path.exists(filepath):
            return jsonify({"success": False, "message": "作业文件不存在"}), 404

        # 检查文件类型是否匹配目标任?
        expected_ext = f".{ALL_TASKS[target_task_id]['file_type']}"
        file_ext = os.path.splitext(filepath)[1].lower()
        if file_ext != expected_ext:
            return jsonify({
                "success": False,
                "message": f"文件类型不匹配.目标任务需?{expected_ext} 格式,当前文件是 {file_ext} 格式"
            }), 400

        # 更新提交记录的task_id
        submission['task_id'] = target_task_id
        submission['original_task_id'] = original_task_id  # 记录原始任务ID
        submission['moved_at'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        # 重新评分(调用评分引擎)
        engine = ScoringEngine(target_task_id, filepath)
        result = engine.score()

        # 如果有目标任务的答案文件,优先使用答案对比评?
        answer_file = _get_answer_file(target_task_id)
        if answer_file and os.path.exists(answer_file):
            try:
                from app.services.answer_comparator import ExcelAnswerComparator
                comp = ExcelAnswerComparator(answer_file, filepath)
                comp_result = comp.compare_task(target_task_id)

                # 计算归一化分数(百分制)
                comp_total = comp_result['total_score']
                comp_max = comp_result['max_score']
                normalized_score = round(comp_total / comp_max * 100) if comp_max > 0 else 0

                # 重新计算等级
                from app.services.scoring_rules import get_level
                level_key, level_info = get_level(normalized_score)

                # 生成 category_summary(按category分组统计算
                comp_details = comp_result.get('details', [])
                cat_summary = {}
                for d in comp_details:
                    cat = d.get('category', '答案对比')
                    if cat not in cat_summary:
                        cat_summary[cat] = {"passed": 0, "total": 0, "score": 0, "max_score": 0}
                    cat_summary[cat]["total"] += 1
                    if d.get('passed'):
                        cat_summary[cat]["passed"] += 1
                        cat_summary[cat]["score"] += d.get('score', 0)
                    cat_summary[cat]["max_score"] += d.get('score', 0)

                # 用答案对比结果完全替换评分结?
                result['total_score'] = normalized_score
                result['max_score'] = 100
                result['percentage'] = round(normalized_score / 100 * 100, 1)
                result['level'] = level_key
                result['level_name'] = level_info['name']
                result['level_color'] = level_info['color']
                result['level_desc'] = level_info['description']
                result['details'] = comp_details
                result['category_summary'] = cat_summary
                result['scoring_method'] = comp_result.get('scoring_method', 'answer_compare')
            except Exception as e:
                # 对比评分失败,保留规则评分结?
                import traceback
                traceback.print_exc()
                result['scoring_method'] = 'rule_check'

        # 更新评分结果
        submission['score_result'] = result
        submission['score_result']['original_task_id'] = original_task_id  # 在评分结果中也记录原始任务
        # 保存更新后的提交记录到数据库
        submission_service.save_submission(submission)

        # 更新内存中的记录
        submissions[submission_id] = submission

        return jsonify({
            "success": True,
            "message": f"作业已移动到 {target_task_id}",
            "data": {
                "submission_id": submission_id,
                "target_task_id": target_task_id,
                "original_task_id": original_task_id,
                "new_score": result.get('total_score', 0),
                "new_level": result.get('level_name', '--')
            }
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "message": f"移动作业失败: {str(e)}"}), 500


# ========== 打包下载全部作业 API ==========

@api_bp.route('/homework-management/download-all', methods=['GET'])
@api_teacher_required
def download_all_submissions():
    """打包下载某个任务的所有学生作业"""
    import zipfile
    import io

    task_id = request.args.get('task_id', '').strip()
    if not task_id:
        return jsonify({"success": False, "message": "未指定任务ID"}), 400

    upload_dir = current_app.config['UPLOAD_FOLDER']

    # 从数据库获取该任务的所有提交记录
    all_submissions = submission_service.get_all_submissions_full()
    sub_list = []
    for data in all_submissions:
        if data.get('task_id') != task_id:
            continue
        if _is_teacher_submission(data):
            continue
        sub_list.append(data)

    if not sub_list:
        return jsonify({"success": False, "message": "该任务暂无学生提交"}), 404

    # 创建ZIP文件
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        for sub in sub_list:
            student_name = sub.get('student_name', sub.get('user_id', '未知'))
            student_class = sub.get('student_class', '')
            user_id = sub.get('user_id', '')
            original_filename = sub.get('original_filename', '')
            ext = os.path.splitext(original_filename)[1] if original_filename else '.xlsx'

            if student_class:
                zip_name = f"{student_class}_{student_name}_{user_id}{ext}"
            else:
                zip_name = f"{student_name}_{user_id}{ext}"

            file_path = sub.get('file_path', '') or sub.get('filepath', '')
            if file_path and os.path.exists(file_path):
                zf.write(file_path, zip_name)
            else:
                sid = sub.get('submission_id', '')
                for ext_try in ['.xlsx', '.docx', '.pptx']:
                    candidate = os.path.join(upload_dir, f"{sid}{ext_try}")
                    if os.path.exists(candidate):
                        zf.write(candidate, zip_name)
                        break

    zip_buffer.seek(0)

    task_name = ALL_TASKS.get(task_id, {}).get('task_name', task_id)
    from flask import send_file
    return send_file(
        zip_buffer,
        mimetype='application/zip',
        as_attachment=True,
        download_name=f"{task_name}_全部作业.zip"
    )


@api_bp.route('/homework-management/recalculate', methods=['POST'])
@api_teacher_required
def recalculate_all_scores():
    """
    重新计算所有学生的成绩

    请求参数:
    {
        "task_id": "task_5_4_1"  # 可选,如果指定则只计算该任务,否则计算所有
    }

    返回:
    {
        "success": true,
        "message": "已重新计算15 份作业的成绩",
        "details": {
            "total": 15,
            "success": 15,
            "failed": 0
        }
    }
    """
    data = request.get_json(silent=True) or {}
    task_id_filter = data.get('task_id', '').strip()

    total = 0
    success_count = 0
    failed_count = 0
    failed_details = []

    # 从数据库获取所有提交记录
    all_submissions = submission_service.get_all_submissions_full()
    if not all_submissions:
        return jsonify({
            "success": True,
            "message": "暂无提交记录",
            "details": {"total": 0, "success": 0, "failed": 0}
        })

    for submission in all_submissions:
        try:
            # 排除教师提交
            if _is_teacher_submission(submission):
                continue

            # 如果指定了task_id,只处理该任务的提交
            submission_task_id = submission.get('task_id', '')
            if task_id_filter and submission_task_id != task_id_filter:
                continue

            total += 1

            # 获取文件路径
            filepath = submission.get('filepath', '')
            if not filepath or not os.path.exists(filepath):
                failed_count += 1
                failed_details.append(f"{filename}: 文件不存在")
                continue

            # 确定要使用的task_id(优先使用original_task_id,如果存在的话)
            target_task_id = submission.get('original_task_id', submission_task_id)
            if not target_task_id:
                failed_count += 1
                failed_details.append(f"{filename}: 无法确定任务ID")
                continue

            # 调用评分引擎重新评分
            engine = ScoringEngine(target_task_id, filepath)
            result = engine.score()

            # 如果有答案文件,优先使用答案对比评分
            answer_file = _get_answer_file(target_task_id)
            if answer_file and os.path.exists(answer_file):
                try:
                    from app.services.answer_comparator import ExcelAnswerComparator
                    comp = ExcelAnswerComparator(answer_file, filepath)
                    comp_result = comp.compare_task(target_task_id)

                    # 计算归一化分数(百分制)
                    comp_total = comp_result['total_score']
                    comp_max = comp_result['max_score']
                    normalized_score = round(comp_total / comp_max * 100) if comp_max > 0 else 0

                    # 重新计算等级
                    from app.services.scoring_rules import get_level
                    level_key, level_info = get_level(normalized_score)

                    # 生成 category_summary
                    comp_details = comp_result.get('details', [])
                    cat_summary = {}
                    for d in comp_details:
                        cat = d.get('category', '答案对比')
                        if cat not in cat_summary:
                            cat_summary[cat] = {"passed": 0, "total": 0, "score": 0, "max_score": 0}
                        cat_summary[cat]["total"] += 1
                        if d.get('passed'):
                            cat_summary[cat]["passed"] += 1
                            cat_summary[cat]["score"] += d.get('score', 0)
                        cat_summary[cat]["max_score"] += d.get('score', 0)

                    # 用答案对比结果完全替换评分结?
                    result['total_score'] = normalized_score
                    result['max_score'] = 100
                    result['percentage'] = round(normalized_score / 100 * 100, 1)
                    result['level'] = level_key
                    result['level_name'] = level_info['name']
                    result['level_color'] = level_info['color']
                    result['level_desc'] = level_info['description']
                    result['details'] = comp_details
                    result['category_summary'] = cat_summary
                    result['scoring_method'] = comp_result.get('scoring_method', 'answer_compare')
                except Exception as e:
                    # 对比评分失败,保留规则评分结?
                    import traceback
                    traceback.print_exc()
                    result['scoring_method'] = 'rule_check'

            # 更新评分结果
            submission['score_result'] = result
            submission['recalculated_at'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

            # 保存更新后的提交记录到数据库
            submission_id = submission.get('submission_id', '')
            submission_service.save_submission(submission)

            # 更新内存中的记录
            submissions[submission_id] = submission

            success_count += 1

        except Exception as e:
            failed_count += 1
            failed_details.append(f"{submission.get('submission_id', 'unknown')}: {str(e)}")
            import traceback
            traceback.print_exc()

    # 生成返回消息
    if total == 0:
        message = "没有找到需要重新计算的作业"
    else:
        message = f"已重新计算{success_count} 份作业的成绩"
        if failed_count > 0:
            message += f",{failed_count} 份失败."

    return jsonify({
        "success": True,
        "message": message,
        "details": {
            "total": total,
            "success": success_count,
            "failed": failed_count,
            "failed_details": failed_details[:10]  # 最多返回10条错误详情
        }
    })
