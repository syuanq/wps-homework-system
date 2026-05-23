# -*- coding: utf-8 -*-
"""
分层练习文档生成服务
根据AI生成的分层任务，创建Word练习文档和素材包，打包为ZIP供下载
"""
import os
import json
import zipfile
import io
from datetime import datetime

# 数据目录
DATA_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), 'data')
GENERATED_DIR = os.path.join(DATA_DIR, 'generated')


def _ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def generate_practice_docx(task_info, level_key, tasks, student_info=None):
    """生成Word练习文档

    Args:
        task_info: 任务信息（task_id, task_name, module, description）
        level_key: 等级键（excellent/good/pass/fail）
        tasks: 该等级的任务列表 [{"name": "", "description": "", "difficulty": ""}]
        student_info: 学生信息 {"name": "", "class_name": "", "score": ""}

    Returns:
        bytes: Word文档的二进制内容
    """
    from docx import Document
    from docx.shared import Pt, Cm, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn

    LEVEL_NAMES = {
        "excellent": "优秀（拓展挑战）",
        "good": "良好（巩固提升）",
        "pass": "及格（基础强化）",
        "fail": "需努力（入门学习）"
    }
    LEVEL_COLORS = {
        "excellent": "28a745",
        "good": "17a2b8",
        "pass": "ffc107",
        "fail": "dc3545"
    }

    doc = Document()

    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')

    # ===== 标题 =====
    title = doc.add_heading('', level=0)
    run = title.add_run(f'{task_info["module"]} 分层练习任务')
    run.font.size = Pt(26)
    run.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ===== 副标题 =====
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(f'任务：{task_info["task_name"]}')
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)

    # ===== 等级标签 =====
    level_para = doc.add_paragraph()
    level_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    level_run = level_para.add_run(f'【{LEVEL_NAMES.get(level_key, level_key)}】')
    level_run.font.size = Pt(16)
    level_run.font.bold = True
    color_hex = LEVEL_COLORS.get(level_key, "333333")
    level_run.font.color.rgb = RGBColor(int(color_hex[:2], 16), int(color_hex[2:4], 16), int(color_hex[4:], 16))

    # ===== 学生信息（如果有）=====
    if student_info:
        doc.add_paragraph()
        info_table = doc.add_table(rows=1, cols=3)
        info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cells = info_table.rows[0].cells
        info_data = [
            f'姓名：{student_info.get("name", "--")}',
            f'班级：{student_info.get("class_name", "--")}',
            f'得分：{student_info.get("score", "--")}'
        ]
        for i, text in enumerate(info_data):
            cells[i].text = text
            for paragraph in cells[i].paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(11)

    # ===== 分隔线 =====
    doc.add_paragraph('─' * 50).alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ===== 练习说明 =====
    doc.add_heading('练习说明', level=1)
    intro = doc.add_paragraph()
    intro.add_run('本练习文档根据你的作业评分结果，由AI智能生成。请按照以下任务要求完成练习，完成后可重新提交作业进行评分。').font.size = Pt(11)

    # 素材文件说明
    material_tip = doc.add_paragraph()
    material_tip_run = material_tip.add_run('📎 素材文件说明：')
    material_tip_run.font.bold = True
    material_tip_run.font.size = Pt(11)
    material_tip_run.font.color.rgb = RGBColor(0xDC, 0x35, 0x45)
    
    module = task_info.get('module', '')
    file_type = task_info.get('file_type', 'docx')
    if file_type == 'xlsx' or '表格' in module:
        material_name = '练习素材.xlsx（WPS表格文件）'
    elif file_type == 'pptx' or '演示' in module:
        material_name = '练习素材.pptx（WPS演示文件）'
    else:
        material_name = '练习素材.docx（WPS文档文件）'
    
    material_tip.add_run(f'练习包中包含 {material_name}，请在该素材文件上按照下方操作步骤进行练习。所有操作步骤均基于该素材文件设计，请确保打开正确的文件。').font.size = Pt(11)

    tips = doc.add_paragraph()
    tips_run = tips.add_run('温馨提示：')
    tips_run.font.bold = True
    tips_run.font.size = Pt(11)
    tips.add_run('请仔细阅读每个任务的操作要求，参照教材对应步骤完成练习。遇到困难时可以回顾教材或向老师请教。').font.size = Pt(11)

    # ===== 练习任务 =====
    doc.add_heading('练习任务', level=1)

    for i, task in enumerate(tasks, 1):
        # 任务标题
        task_heading = doc.add_heading(f'任务{i}：{task.get("name", "未命名任务")}', level=2)

        # 难度标签
        diff = task.get("difficulty", "")
        if diff:
            diff_para = doc.add_paragraph()
            diff_run = diff_para.add_run(f'难度：{diff}')
            diff_run.font.size = Pt(10)
            diff_run.font.color.rgb = RGBColor(0x71, 0x71, 0x71)
            diff_run.font.italic = True

        # 任务描述
        desc = task.get("description", "")
        if desc:
            desc_para = doc.add_paragraph()
            desc_para.add_run('任务要求：').font.bold = True
            desc_para.add_run(desc)

        # AI生成的具体操作步骤
        steps = task.get("steps", [])
        if steps:
            doc.add_paragraph()
            steps_title = doc.add_paragraph()
            steps_title_run = steps_title.add_run('📋 操作步骤（请按顺序完成）：')
            steps_title_run.font.bold = True
            steps_title_run.font.size = Pt(11)
            steps_title_run.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)

            for j, step in enumerate(steps, 1):
                step_para = doc.add_paragraph()
                step_para.paragraph_format.left_indent = Cm(0.5)
                step_para.paragraph_format.space_after = Pt(4)
                # 步骤编号圆圈
                num_run = step_para.add_run(f'  {j}. ')
                num_run.font.bold = True
                num_run.font.size = Pt(11)
                num_run.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)
                # 步骤内容
                step_text = step.strip()
                if step_text.startswith(('步骤', '第')):
                    # 去掉重复的"步骤X"前缀
                    step_text = step_text.lstrip('步骤第0123456789：:、. ')
                    step_text = step_text.lstrip('步')
                    if step_text.startswith('骤'):
                        step_text = step_text[1:]
                    step_text = step_text.lstrip('：:、. ')
                content_run = step_para.add_run(step_text)
                content_run.font.size = Pt(11)
                # 完成勾选框
                check_run = step_para.add_run('  ☐')
                check_run.font.size = Pt(12)
                check_run.font.color.rgb = RGBColor(0x28, 0xa7, 0x45)

        # 完成情况勾选
        doc.add_paragraph()
        check_para = doc.add_paragraph()
        check_para.add_run('完成情况：').font.bold = True
        check_para.add_run('  □ 已完成    □ 部分完成    □ 未完成')

        # 备注/问题记录区域
        note_para = doc.add_paragraph()
        note_para.add_run('遇到的问题或心得：').font.bold = True
        note_para.add_run(' ________________________________________________________________________')

        doc.add_paragraph()  # 空行

    # ===== 学习反思 =====
    doc.add_heading('学习反思', level=1)
    reflection_questions = [
        '通过本次练习，我掌握了哪些新技能？',
        '哪些操作我还不太熟练，需要继续练习？',
        '我计划如何改进自己的学习方法？'
    ]
    for q in reflection_questions:
        q_para = doc.add_paragraph()
        q_para.add_run(f'{q}').font.bold = True
        for _ in range(2):
            doc.add_paragraph('    ______________________________________________________________________')
        doc.add_paragraph()

    # ===== 页脚信息 =====
    doc.add_paragraph('─' * 50).alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_para.add_run(f'文档生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M")}  |  信息技术基础作业智能评价系统')
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 保存到内存
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_material_file(task_info, level_key, tasks):
    """根据任务类型生成素材文件
    优先使用AI返回的material字段动态生成，否则使用预设模板
    WPS文字 → docx, WPS表格 → xlsx, WPS演示 → pptx
    """
    # 优先使用AI返回的material字段
    if tasks and isinstance(tasks, list) and len(tasks) > 0:
        task = tasks[0]
        material = task.get("material")
        if material and isinstance(material, dict):
            mat_type = material.get("type", "")
            try:
                if mat_type == "xlsx" or '表格' in task_info.get('module', ''):
                    return _generate_xlsx_from_material(material)
                elif mat_type == "pptx" or '演示' in task_info.get('module', ''):
                    return _generate_pptx_from_material(material)
                elif mat_type == "docx":
                    return _generate_docx_from_material(material)
            except Exception as e:
                print(f"AI素材生成失败，回退到预设模板: {e}")

    # 回退到预设模板
    module = task_info.get('module', '')
    file_type = task_info.get('file_type', 'docx')

    if file_type == 'xlsx' or '表格' in module:
        return _generate_xlsx_material(task_info, level_key, tasks)
    elif file_type == 'pptx' or '演示' in module:
        return _generate_pptx_material(task_info, level_key, tasks)
    else:
        return _generate_docx_material(task_info, level_key, tasks)


def _generate_xlsx_from_material(material):
    """根据AI返回的material字段动态生成xlsx素材文件"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side

    wb = Workbook()
    sheets = material.get("sheets", [])

    for i, sheet_def in enumerate(sheets):
        if i == 0:
            ws = wb.active
            ws.title = sheet_def.get("name", "Sheet1")
        else:
            ws = wb.create_sheet(title=sheet_def.get("name", f"Sheet{i+1}"))

        headers = sheet_def.get("headers", [])
        data = sheet_def.get("data", [])

        # 写入表头
        if headers:
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True, size=11)
                cell.alignment = Alignment(horizontal='center')

        # 写入数据
        for row_idx, row_data in enumerate(data, 2):
            if isinstance(row_data, list):
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)

        # 自动调整列宽
        for col in ws.columns:
            max_length = 0
            col_letter = col[0].column_letter
            for cell in col:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            ws.column_dimensions[col_letter].width = min(max(max_length + 2, 8), 30)

    # 保存到内存
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue(), "练习素材.xlsx"


def _generate_docx_from_material(material):
    """根据AI返回的material字段动态生成docx素材文件"""
    from docx import Document
    from docx.shared import Pt
    from docx.oxml.ns import qn

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style.font.size = Pt(12)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')

    paragraphs = material.get("paragraphs", [])
    for p_def in paragraphs:
        if isinstance(p_def, str):
            doc.add_paragraph(p_def)
        elif isinstance(p_def, dict):
            p = doc.add_paragraph()
            text = p_def.get("text", "")
            bold = p_def.get("bold", False)
            size = p_def.get("size", 12)
            align = p_def.get("align", "left")
            run = p.add_run(text)
            run.font.size = Pt(size)
            run.font.bold = bold
            if align == "center":
                p.alignment = 1  # WD_ALIGN_PARAGRAPH.CENTER
            elif align == "right":
                p.alignment = 2

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue(), "练习素材.docx"


def _generate_pptx_from_material(material):
    """根据AI返回的material字段动态生成pptx素材文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slides_def = material.get("slides", [])

    for slide_def in slides_def:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
        title = slide_def.get("title", "")
        content = slide_def.get("content", "")

        slide.shapes.title.text = title
        if content and slide.placeholders[1]:
            slide.placeholders[1].text = content

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), "练习素材.pptx"


def _generate_docx_material(task_info, level_key, tasks):
    """生成WPS文字素材文件（docx）"""
    from docx import Document
    from docx.shared import Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    doc = Document()

    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = '宋体'
    font.size = Pt(12)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')

    # 标题
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(task_info['task_name'])
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.name = '黑体'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), '黑体')

    doc.add_paragraph()  # 空行

    # 根据任务生成对应的素材内容
    task_id = task_info.get('task_id', '')
    if '4_1' in task_id or '信' in task_info.get('task_name', ''):
        # 致大一新生的一封信
        doc.add_paragraph('亲爱的学弟学妹们：')
        doc.add_paragraph('')
        doc.add_paragraph('    你们好！我是你们的学长/学姐。时光飞逝，转眼间我已经在大学度过了难忘的时光。在这里，我想和你们分享一些我的大学生活感悟，希望能对你们有所帮助。')
        doc.add_paragraph('')
        doc.add_paragraph('    大学是一个全新的开始。在这里，你们将遇到来自五湖四海的同学，学习各种有趣的知识，参加丰富多彩的社团活动。希望你们能珍惜这段宝贵的时光，努力学习，积极进取。')
        doc.add_paragraph('')
        doc.add_paragraph('    首先，学习是大学生的首要任务。要学会自主学习，合理安排时间，养成良好的学习习惯。其次，要多参加实践活动，锻炼自己的综合能力。最后，要学会与人相处，建立良好的人际关系。')
        doc.add_paragraph('')
        doc.add_paragraph('    祝你们在大学里度过充实而美好的时光！')
        doc.add_paragraph('')
        doc.add_paragraph('    此致')
        doc.add_paragraph('    敬礼！')
        doc.add_paragraph('')
        doc.add_paragraph('    一个关心你们的学长/学姐')
        doc.add_paragraph(f'    {datetime.now().strftime("%Y年%m月%d日")}')

    elif '4_2' in task_id or '感动' in task_info.get('task_name', ''):
        # 感动中国人物
        doc.add_paragraph('感动中国人物事迹介绍')
        doc.add_paragraph('')
        for name in ['张桂梅', '袁隆平', '钟南山']:
            doc.add_paragraph(f'【{name}】')
            doc.add_paragraph(f'    {name}，一位令人敬佩的{name}。他/她用自己的行动诠释了什么是真正的奉献精神。在平凡的岗位上，他/她做出了不平凡的贡献，感动了无数中国人。')
            doc.add_paragraph(f'    主要事迹：他/她长期致力于自己的事业，克服了重重困难，取得了卓越的成就。他/她的精神激励着一代又一代人不断前行。')
            doc.add_paragraph('')

    elif '4_3' in task_id or '职业' in task_info.get('task_name', ''):
        # 职业生涯规划
        doc.add_paragraph('一、前言')
        doc.add_paragraph('')
        doc.add_paragraph('    在当今竞争激烈的社会中，做好职业生涯规划对于每一位大学生来说都至关重要。职业生涯规划不仅能帮助我们明确发展方向，还能让我们在大学期间有针对性地提升自己的能力。')
        doc.add_paragraph('')
        doc.add_paragraph('二、自我分析')
        doc.add_paragraph('')
        doc.add_paragraph('    1、个人优势：性格开朗，善于沟通，具有较强的学习能力和团队合作精神。')
        doc.add_paragraph('    2、个人劣势：缺乏实践经验，专业知识不够扎实，需要进一步加强。')
        doc.add_paragraph('    3、兴趣爱好：喜欢阅读、运动、参加社团活动。')
        doc.add_paragraph('')
        doc.add_paragraph('三、职业目标')
        doc.add_paragraph('')
        doc.add_paragraph('    短期目标：掌握专业知识，考取相关证书。')
        doc.add_paragraph('    中期目标：进入理想的企业实习，积累工作经验。')
        doc.add_paragraph('    长期目标：成为行业内的专业人才。')
        doc.add_paragraph('')
        doc.add_paragraph('四、实施计划')
        doc.add_paragraph('')
        doc.add_paragraph('    1、大学期间：认真学习专业课程，积极参加实践活动。')
        doc.add_paragraph('    2、毕业前：完成毕业论文，准备好求职材料。')
        doc.add_paragraph('    3、工作后：不断学习，提升自己的专业能力。')
        doc.add_paragraph('')

    elif '4_4' in task_id or '简历' in task_info.get('task_name', ''):
        # 个人简历
        doc.add_paragraph('个人简历')
        doc.add_paragraph('')
        doc.add_paragraph('基本信息')
        doc.add_paragraph('    姓名：张三        性别：男')
        doc.add_paragraph('    出生日期：2005年6月    籍贯：广东广州')
        doc.add_paragraph('    联系电话：13800138000    电子邮箱：zhangsan@example.com')
        doc.add_paragraph('    求职意向：软件开发工程师')
        doc.add_paragraph('')
        doc.add_paragraph('教育背景')
        doc.add_paragraph('    2023.09 - 至今    XX大学    计算机科学与技术专业    本科')
        doc.add_paragraph('    主修课程：数据结构、操作系统、计算机网络、数据库原理')
        doc.add_paragraph('')
        doc.add_paragraph('专业技能')
        doc.add_paragraph('    - 熟悉Python、Java编程语言')
        doc.add_paragraph('    - 了解MySQL数据库操作')
        doc.add_paragraph('    - 熟练使用Office办公软件')
        doc.add_paragraph('')
        doc.add_paragraph('实践经历')
        doc.add_paragraph('    2024.06 - 2024.09    XX科技公司    实习生')
        doc.add_paragraph('    - 参与公司内部系统的开发与维护')
        doc.add_paragraph('    - 协助完成项目文档编写')
        doc.add_paragraph('')
        doc.add_paragraph('自我评价')
        doc.add_paragraph('    本人性格开朗，学习能力强，有良好的团队合作精神。在校期间积极参加各类活动，具有较强的组织协调能力。')
        doc.add_paragraph('')
    else:
        # 通用素材
        doc.add_paragraph('请根据练习文档中的任务要求，在此文件中完成练习。')
        doc.add_paragraph('')
        doc.add_paragraph('注意事项：')
        doc.add_paragraph('1. 仔细阅读练习文档中的操作步骤')
        doc.add_paragraph('2. 按照步骤要求完成每项操作')
        doc.add_paragraph('3. 完成后保存文件并提交')

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue(), '练习素材.docx'


def _generate_xlsx_material(task_info, level_key, tasks):
    """生成WPS表格素材文件（xlsx）"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = '素材数据'

    # 设置列宽
    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 15
    ws.column_dimensions['C'].width = 15
    ws.column_dimensions['D'].width = 15
    ws.column_dimensions['E'].width = 15

    header_font = Font(bold=True, size=12)
    header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
    header_font_white = Font(bold=True, size=12, color='FFFFFF')
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    task_id = task_info.get('task_id', '')

    if '5' in task_id or '二手房' in task_info.get('task_name', ''):
        # 二手房房源信息
        headers = ['序号', '小区名称', '户型', '面积(㎡)', '单价(元/㎡)', '总价(万元)', '楼层', '朝向', '装修情况', '建成年份']
        ws.column_dimensions['F'].width = 15
        ws.column_dimensions['G'].width = 10
        ws.column_dimensions['H'].width = 10
        ws.column_dimensions['I'].width = 15
        ws.column_dimensions['J'].width = 12

        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = header_font_white
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border

        data = [
            [1, '阳光花园', '3室2厅', 120, 25000, 300, '中层', '南北通透', '精装修', 2018],
            [2, '翠湖名苑', '2室1厅', 85, 28000, 238, '高层', '朝南', '简装修', 2020],
            [3, '金地华府', '3室2厅', 135, 22000, 297, '低层', '东西通透', '毛坯', 2015],
            [4, '碧桂园', '4室2厅', 160, 20000, 320, '中层', '南北通透', '精装修', 2019],
            [5, '万科城', '2室2厅', 95, 26000, 247, '高层', '朝南', '简装修', 2021],
            [6, '保利花园', '3室1厅', 110, 23000, 253, '低层', '朝东', '精装修', 2017],
            [7, '恒大名都', '2室1厅', 78, 30000, 234, '中层', '朝南', '毛坯', 2022],
            [8, '龙湖天街', '3室2厅', 125, 24000, 300, '高层', '南北通透', '精装修', 2016],
            [9, '中海锦城', '4室2厅', 150, 21000, 315, '低层', '东西通透', '简装修', 2014],
            [10, '融创御湖', '2室2厅', 90, 27000, 243, '中层', '朝南', '精装修', 2020],
        ]
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
    else:
        # 通用表格素材
        headers = ['序号', '项目', '数据1', '数据2', '备注']
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = header_font_white
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border

        for i in range(1, 6):
            for col in range(1, 6):
                cell = ws.cell(row=i+1, column=col, value=f'示例数据{i}-{col}')
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.getvalue(), '练习素材.xlsx'


def _generate_pptx_material(task_info, level_key, tasks):
    """生成WPS演示素材文件（pptx）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    task_id = task_info.get('task_id', '')

    if '6' in task_id or '红色' in task_info.get('task_name', '') or '旅游' in task_info.get('task_name', ''):
        # 红色旅游地宣传演示
        slides_content = [
            ('红色旅游地宣传', '探索革命圣地  传承红色精神', 'center'),
            ('井冈山', '中国革命的摇篮\n\n• 位于江西省吉安市\n• 1927年毛泽东在此创建第一个农村革命根据地\n• 主要景点：黄洋界、茨坪、茅坪', 'left'),
            ('延安', '革命圣地\n\n• 位于陕西省北部\n• 1935-1948年中共中央所在地\n• 主要景点：宝塔山、杨家岭、枣园', 'left'),
            ('遵义', '转折之城\n\n• 位于贵州省\n• 1935年遵义会议在此召开\n• 主要景点：遵义会议会址、红军山', 'left'),
            ('西柏坡', '新中国从这里走来\n\n• 位于河北省平山县\n• 1948-1949年中共中央所在地\n• 主要景点：西柏坡纪念馆、中共中央旧址', 'left'),
        ]
    else:
        slides_content = [
            ('演示文稿标题', '请根据练习要求完成制作', 'center'),
            ('第一页内容', '这是素材幻灯片\n请在此基础上进行编辑', 'left'),
            ('第二页内容', '可以添加文字、图片、图表等元素', 'left'),
        ]

    for title_text, content_text, align in slides_content:
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)  # 红色
        p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT

        # 内容
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for line in content_text.split('\n'):
            p2 = tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p2.space_after = Pt(12)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue(), '练习素材.pptx'


def generate_practice_package(task_id, level_key, tasks, student_info=None):
    """生成完整的练习包（ZIP）

    Args:
        task_id: 任务ID
        level_key: 等级键
        tasks: 该等级的任务列表
        student_info: 学生信息（可选）

    Returns:
        tuple: (zip_bytes, filename)
    """
    from .scoring_rules import ALL_TASKS

    task_info = ALL_TASKS.get(task_id, {})
    if not task_info:
        task_info = {"task_id": task_id, "task_name": "未知任务", "module": "未知模块", "description": ""}

    _ensure_dir(GENERATED_DIR)

    # 生成Word文档（操作步骤）
    docx_bytes = generate_practice_docx(task_info, level_key, tasks, student_info)

    # 生成素材文件（根据任务类型：docx/xlsx/pptx）
    material_bytes, material_filename = generate_material_file(task_info, level_key, tasks)

    # 打包为ZIP
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        zf.writestr('操作步骤.docx', docx_bytes)
        zf.writestr(material_filename, material_bytes)

    zip_buffer.seek(0)
    zip_bytes = zip_buffer.getvalue()

    # 生成文件名
    level_names = {"excellent": "优秀", "good": "良好", "pass": "及格", "fail": "入门"}
    level_name = level_names.get(level_key, level_key)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    student_suffix = f"_{student_info.get('name', '')}" if student_info and student_info.get('name') else ""
    filename = f"{task_info['task_name']}_分层练习_{level_name}{student_suffix}_{timestamp}.zip"

    return zip_bytes, filename


def save_generated_package(task_id, level_key, tasks, student_info=None):
    """保存生成的练习包到磁盘（教师预生成用）

    Returns:
        str: 保存的文件路径
    """
    _ensure_dir(GENERATED_DIR)
    zip_bytes, filename = generate_practice_package(task_id, level_key, tasks, student_info)
    filepath = os.path.join(GENERATED_DIR, filename)
    with open(filepath, 'wb') as f:
        f.write(zip_bytes)
    return filepath


def get_generated_packages():
    """获取所有已生成的练习包列表"""
    _ensure_dir(GENERATED_DIR)
    packages = []
    for fname in sorted(os.listdir(GENERATED_DIR), reverse=True):
        if fname.endswith('.zip'):
            fpath = os.path.join(GENERATED_DIR, fname)
            stat = os.stat(fpath)
            packages.append({
                'filename': fname,
                'filepath': fpath,
                'size': stat.st_size,
                'created_at': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                'size_str': f"{stat.st_size / 1024:.1f}KB"
            })
    return packages


def delete_generated_package(filename):
    """删除已生成的练习包"""
    filepath = os.path.join(GENERATED_DIR, filename)
    if os.path.exists(filepath):
        os.remove(filepath)
        return True
    return False


def generate_lesson_plan_docx(content, weak_points):
    """生成教案Word文档"""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.oxml.ns import qn
    import io, os
    from datetime import datetime
    
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style.font.size = Pt(12)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    
    # 标题
    title = doc.add_heading('针对性教学教案', level=0)
    title.alignment = 1
    
    # 副标题
    sub = doc.add_paragraph()
    sub.alignment = 1
    run = sub.add_run(f'基于薄弱技能点分析 | 生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M")}')
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    
    # 薄弱技能点摘要
    doc.add_heading('一、薄弱技能点分析', level=1)
    for i, wp in enumerate(weak_points[:5], 1):
        p = doc.add_paragraph()
        run = p.add_run(f'{i}. {wp["name"]}（{wp["task_name"]}）')
        run.font.bold = True
        p.add_run(f' — 错误率 {wp["fail_rate"]}%（{wp["fail_count"]}/{wp["total"]}人未通过）')
    
    # 教案正文
    doc.add_heading('二、教案内容', level=1)
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            doc.add_paragraph()
        elif line.startswith('# ') or line.startswith('## '):
            level = 2 if line.startswith('## ') else 1
            doc.add_heading(line.lstrip('# ').strip(), level=level)
        elif line.startswith('**') and line.endswith('**'):
            p = doc.add_paragraph()
            run = p.add_run(line.strip('*'))
            run.font.bold = True
        elif line.startswith('- ') or line.startswith('• '):
            doc.add_paragraph(line[2:], style='List Bullet')
        else:
            doc.add_paragraph(line)
    
    # 保存
    gen_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), 'data', 'generated')
    os.makedirs(gen_dir, exist_ok=True)
    filename = f'教案_{datetime.now().strftime("%Y%m%d_%H%M%S")}.docx'
    filepath = os.path.join(gen_dir, filename)
    doc.save(filepath)
    return filepath
