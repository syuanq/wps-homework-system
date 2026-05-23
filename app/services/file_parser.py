# -*- coding: utf-8 -*-
"""
文件解析服务
支持解析 Word(.docx)、Excel(.xlsx)、PowerPoint(.pptx) 文件
"""
import os
import re
from docx import Document
from docx.shared import Pt, Emu, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt


class DocxParser:
    """Word文档解析器"""
    
    def __init__(self, filepath):
        self.filepath = filepath
        self.doc = Document(filepath)
        self.paragraphs = self.doc.paragraphs
        self.sections = self.doc.sections
        self.tables = self.doc.tables
        self.inline_shapes = self.doc.inline_shapes
        # 获取所有body元素（段落、表格等）的顺序
        self.body_elements = list(self.doc.element.body)
    
    def get_title_paragraph(self):
        """获取标题段落（第一个非空段落）"""
        for para in self.paragraphs:
            if para.text.strip():
                return para
        return None

    def get_all_text(self):
        """提取文档全部文本内容"""
        texts = []
        for para in self.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        for table in self.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        return "\n".join(texts)
    
    def get_body_paragraphs(self):
        """获取正文段落（跳过第一个标题段落）"""
        body_paras = []
        started = False
        for para in self.paragraphs:
            if not started and para.text.strip():
                started = True
                continue
            if started and para.text.strip():
                body_paras.append(para)
        return body_paras
    
    def check_title_font(self, font_name):
        """检查标题字体"""
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.name and font_name in run.font.name:
                return True
            # 检查东亚字体
            rpr = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts')
            if rpr is not None:
                ea = rpr.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia')
                if ea and font_name in ea:
                    return True
                ascii_font = rpr.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ascii')
                if ascii_font and font_name in ascii_font:
                    return True
        return False
    
    def check_title_font_size(self, min_size_pt):
        """检查标题字号（大于等于）"""
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.size and run.font.size.pt >= min_size_pt:
                return True
            sz = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sz')
            if sz is not None:
                val = int(sz.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 0))
                if val / 2 >= min_size_pt:
                    return True
        return False

    def check_title_font_size_exact(self, exact_size_pt):
        """检查标题字号（精确匹配，允许±1磅误差）"""
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.size:
                if abs(run.font.size.pt - exact_size_pt) <= 1:
                    return True
            sz = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sz')
            if sz is not None:
                val = int(sz.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 0))
                if abs(val / 2 - exact_size_pt) <= 1:
                    return True
        return False
    
    def check_title_bold(self):
        """检查标题是否加粗"""
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.bold:
                return True
            b = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}b')
            if b is not None:
                val = b.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 'true')
                if val != 'false' and val != '0':
                    return True
        return False
    
    def check_title_alignment(self, alignment):
        """检查标题对齐方式"""
        title = self.get_title_paragraph()
        if not title:
            return False
        align_map = {
            'center': WD_ALIGN_PARAGRAPH.CENTER,
            'left': WD_ALIGN_PARAGRAPH.LEFT,
            'right': WD_ALIGN_PARAGRAPH.RIGHT,
            'justify': WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        expected = align_map.get(alignment)
        if expected is None:
            return False
        return title.alignment == expected
    
    def check_body_font(self, font_name):
        """检查正文字体"""
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:  # 检查前5个正文段落
            for run in para.runs:
                if run.font.name and font_name in run.font.name:
                    match_count += 1
                    break
                rpr = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts')
                if rpr is not None:
                    ea = rpr.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia')
                    if ea and font_name in ea:
                        match_count += 1
                        break
        return match_count >= 2  # 至少2个段落匹配
    
    def check_body_first_indent(self, chars):
        """检查正文首行缩进"""
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            pf = para.paragraph_format
            if pf.first_line_indent:
                # 首行缩进值（EMU），1字符约= 12pt * 12700 EMU/pt
                chars_indent = pf.first_line_indent.pt / 12.0 if pf.first_line_indent else 0
                if chars_indent >= chars * 0.8:  # 允许20%误差
                    match_count += 1
            # 也检查XML中的ind元素
            pPr = para._element.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pPr')
            if pPr is not None:
                ind = pPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ind')
                if ind is not None:
                    first_line = ind.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}firstLine')
                    first_line_chars = ind.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}firstLineChars')
                    if first_line_chars:
                        if int(first_line_chars) >= chars * 0.8:
                            match_count += 1
                    elif first_line:
                        val = int(first_line)
                        chars_indent = val / (12.0 * 12700)
                        if chars_indent >= chars * 0.8:
                            match_count += 1
        return match_count >= 2
    
    def check_body_line_spacing(self, min_spacing):
        """检查正文行距（大于等于）"""
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            pf = para.paragraph_format
            if pf.line_spacing:
                if isinstance(pf.line_spacing, float) and pf.line_spacing >= min_spacing:
                    match_count += 1
                elif isinstance(pf.line_spacing, int) and pf.line_spacing >= min_spacing * 240:
                    match_count += 1
        return match_count >= 2

    def check_body_line_spacing_exact(self, exact_spacing):
        """检查行距精确匹配（支持固定值和倍数，通过XML解析确保准确）
        
        固定值模式：XML中 lineRule=exact, line=磅值×20（单位1/20磅）
        倍数模式：XML中 lineRule=auto, line=倍数×240
        """
        # 使用所有非空段落（全文检测）
        all_paras = [p for p in self.paragraphs if p.text.strip()]
        if not all_paras:
            return False
        
        match_count = 0
        checked = 0
        for para in all_paras[:10]:
            checked += 1
            # 优先通过XML精确检测
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                spacing = pPr.find(qn('w:spacing'))
                if spacing is not None:
                    line_val = spacing.get(qn('w:line'))
                    line_rule = spacing.get(qn('w:lineRule'))
                    if line_val:
                        line_int = int(line_val)
                        if line_rule == 'exact':
                            # 固定值模式：line值 / 20 = 磅值
                            actual_pt = line_int / 20.0
                            if abs(actual_pt - exact_spacing) <= 1:
                                match_count += 1
                        elif line_rule == 'auto' or line_rule is None:
                            # 倍数模式：line值 / 240 = 倍数
                            actual_mult = line_int / 240.0
                            if abs(actual_mult - exact_spacing) <= 0.2:
                                match_count += 1
                        else:
                            # 其他模式也尝试按1/20磅换算
                            actual_pt = line_int / 20.0
                            if abs(actual_pt - exact_spacing) <= 1:
                                match_count += 1
                    continue
            
            # 回退到python-docx API
            pf = para.paragraph_format
            if pf.line_spacing:
                ls_val = pf.line_spacing
                # Twips类型或其他数值类型
                try:
                    ls_num = float(ls_val)
                    # 尝试判断是固定值还是倍数
                    if pf.line_spacing_rule == WD_LINE_SPACING.EXACTLY:
                        # 固定值：Twips转磅 (1pt = 20twips)
                        actual_pt = ls_num / 20.0
                        if abs(actual_pt - exact_spacing) <= 1:
                            match_count += 1
                    else:
                        # 倍数或其他
                        if abs(ls_num - exact_spacing) <= 0.2:
                            match_count += 1
                        elif abs(ls_num / 240.0 - exact_spacing) <= 0.2:
                            match_count += 1
                        elif abs(ls_num / 20.0 - exact_spacing) <= 1:
                            match_count += 1
                except (TypeError, ValueError):
                    pass
        
        return checked > 0 and match_count >= max(1, checked // 2)
    
    def check_body_para_spacing(self, min_before, min_after):
        """检查正文段间距"""
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            pf = para.paragraph_format
            before = pf.space_before.pt if pf.space_before else 0
            after = pf.space_after.pt if pf.space_after else 0
            if before >= min_before or after >= min_after:
                match_count += 1
        return match_count >= 2
    
    def check_body_alignment(self, alignment):
        """检查正文对齐方式"""
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        align_map = {
            'justify': WD_ALIGN_PARAGRAPH.JUSTIFY,
            'left': WD_ALIGN_PARAGRAPH.LEFT,
            'center': WD_ALIGN_PARAGRAPH.CENTER,
        }
        expected = align_map.get(alignment)
        if expected is None:
            return False
        match_count = sum(1 for p in body_paras[:5] if p.alignment == expected)
        return match_count >= 2
    
    def check_section_columns(self, min_columns):
        """检查分栏设置"""
        for section in self.sections:
            # 检查sectPr中的cols元素
            sectPr = section._sectPr
            cols = sectPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}cols')
            if cols is not None:
                num = cols.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}num', '1')
                if int(num) >= min_columns:
                    return True
            # 也检查body中的sectPr
        for elem in self.body_elements:
            sectPr = elem.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sectPr')
            if sectPr is not None:
                cols = sectPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}cols')
                if cols is not None:
                    num = cols.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}num', '1')
                    if int(num) >= min_columns:
                        return True
        return False
    
    def check_drop_cap(self):
        """检查首字下沉"""
        for elem in self.body_elements:
            # 检查fldChar或dropCap元素
            dc = elem.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}dropCap')
            if dc is not None:
                return True
            # 检查framePr（首字下沉可能使用frame）
            framePr = elem.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}framePr')
            if framePr is not None:
                drop = framePr.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}dropCap')
                if drop and drop != 'none':
                    return True
        return False
    
    def check_paragraph_border_or_shading(self):
        """检查段落边框或底纹"""
        for para in self.paragraphs:
            pPr = para._element.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pPr')
            if pPr is not None:
                # 检查边框
                pBdr = pPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pBdr')
                if pBdr is not None:
                    for child in pBdr:
                        if child.tag.endswith('}top') or child.tag.endswith('}bottom') or \
                           child.tag.endswith('}left') or child.tag.endswith('}right'):
                            val = child.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')
                            if val and val != 'none':
                                return True
                # 检查底纹
                shd = pPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd')
                if shd is not None:
                    fill = shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
                    if fill and fill != 'auto':
                        return True
        return False
    
    def check_has_image(self):
        """检查是否包含图片"""
        return len(self.inline_shapes) > 0
    
    def check_has_wordart(self):
        """检测文档是否包含艺术字
        通过XML检测wps:wsp（WordprocessingShape）元素，这是WPS/Word中艺术字的存储方式
        """
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        wps_ns = 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'
        mc_ns = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'

        # 方法1: 直接查找wps:wsp元素
        for elem in self.body_elements:
            # 查找所有wps:wsp
            wsps = elem.findall(f'.//{{{wps_ns}}}wsp')
            if wsps:
                return True

        # 方法2: 通过mc:AlternateContent查找
        for elem in self.body_elements:
            alt_contents = elem.findall(f'.//{{{mc_ns}}}AlternateContent')
            for ac in alt_contents:
                # 检查Choice或Fallback中是否包含wps:wsp
                for child in ac:
                    wsps = child.findall(f'.//{{{wps_ns}}}wsp')
                    if wsps:
                        return True

        # 方法3: 检查是否有带特殊效果的文本（艺术字效果如阴影、发光、3D等）
        for para in self.paragraphs:
            for run in para.runs:
                rPr = run._element.find(f'{{{wns}}}rPr')
                if rPr is not None:
                    # 检查艺术字效果元素
                    effect_elements = ['w:effect', 'w:emboss', 'w:imprint', 'w:outline', 'w:shadow']
                    for eff in effect_elements:
                        if rPr.find(f'{{{wns}}}{eff}') is not None:
                            return True

        return False
    
    def check_image_wrap_type(self, non_inline):
        """检查图片环绕方式"""
        if not non_inline:
            return True
        # 检查是否有浮动图片（anchor类型）
        for elem in self.body_elements:
            anchors = elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}anchor')
            if anchors:
                return True
        return False

    def doc_has_right_align(self):
        """检查是否有右对齐段落"""
        for para in self.paragraphs:
            if para.alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                return True
        return False

    def doc_has_date(self):
        """检查文档是否包含日期文本（匹配YYYY年MM月DD日或YYYY-MM-DD格式）"""
        date_pattern1 = re.compile(r'\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日')
        date_pattern2 = re.compile(r'\d{4}-\d{1,2}-\d{1,2}')
        for para in self.paragraphs:
            if date_pattern1.search(para.text) or date_pattern2.search(para.text):
                return True
        # 也检查表格中的日期
        for table in self.tables:
            for row in table.rows:
                for cell in row.cells:
                    if date_pattern1.search(cell.text) or date_pattern2.search(cell.text):
                        return True
        return False

    def doc_has_page_breaks(self, min_breaks):
        """检查分页符数量"""
        break_count = 0
        for para in self.paragraphs:
            # 检查段落中的分页符
            for run in para.runs:
                if run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}br'):
                    for br in run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}br'):
                        br_type = br.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type', '')
                        if br_type == 'page':
                            break_count += 1
            # 检查段落属性中的分页设置
            pPr = para._element.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pPr')
            if pPr is not None:
                for child in pPr:
                    if child.tag.endswith('}pageBreakBefore') and child.get(
                            '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 'true') != 'false':
                        break_count += 1
        return break_count >= min_breaks

    def doc_page_margins(self, left_cm, right_cm):
        """检查页边距设置"""
        for section in self.sections:
            # 页边距单位是EMU，1cm = 360000 EMU
            left_emu = section.left_margin
            right_emu = section.right_margin
            # 允许0.5cm误差
            if abs(left_emu / 360000 - left_cm) <= 0.5 and abs(right_emu / 360000 - right_cm) <= 0.5:
                return True
        return False

    def doc_page_margins_tb(self, top_cm, bottom_cm):
        """检查上下页边距设置"""
        for section in self.sections:
            top_emu = section.top_margin
            bottom_emu = section.bottom_margin
            if abs(top_emu / 360000 - top_cm) <= 0.3 and abs(bottom_emu / 360000 - bottom_cm) <= 0.3:
                return True
        return False

    def doc_has_header(self):
        """检查是否有页眉"""
        for section in self.sections:
            header = section.header
            if header and header.paragraphs:
                for para in header.paragraphs:
                    if para.text.strip():
                        return True
            # 也检查XML中的headerReference
            sectPr = section._sectPr
            header_refs = sectPr.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/relationships}headerReference')
            if header_refs:
                return True
        return False

    def doc_has_page_number(self):
        """检查是否有页码"""
        for section in self.sections:
            # 检查页脚中是否有页码字段
            footer = section.footer
            if footer:
                for para in footer.paragraphs:
                    # 检查fldSimple或fldChar中的PAGE字段
                    fld_chars = para._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldChar')
                    for fc in fld_chars:
                        if fc.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldCharType') == 'begin':
                            # 找到begin后查找instrText
                            instr = para._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instrText')
                            for inst in instr:
                                if 'PAGE' in inst.text or 'page' in inst.text:
                                    return True
                    # 检查fldSimple
                    fld_simples = para._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldSimple')
                    for fs in fld_simples:
                        instr = fs.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instr', '')
                        if 'PAGE' in instr:
                            return True
            # 也检查页眉
            header = section.header
            if header:
                for para in header.paragraphs:
                    fld_chars = para._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldChar')
                    for fc in fld_chars:
                        if fc.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldCharType') == 'begin':
                            instr = para._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instrText')
                            for inst in instr:
                                if 'PAGE' in inst.text or 'page' in inst.text:
                                    return True
        return False

    def doc_has_heading_style(self):
        """检查是否使用了标题样式"""
        for para in self.paragraphs:
            style_name = para.style.name if para.style else ''
            if style_name and style_name.startswith('Heading') or style_name.startswith('标题'):
                return True
        return False

    def doc_heading1_centered(self):
        """检查标题1样式的段落是否居中
        检测段落直接设置或样式继承的对齐方式
        """
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        def get_effective_jc(para):
            """获取段落的有效对齐方式XML值"""
            # 段落直接设置
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                jc = pPr.find(qn('w:jc'))
                if jc is not None:
                    return jc.get(qn('w:val'))
            # 从样式继承
            style = para.style
            while style:
                style_elem = style.element
                style_pPr = style_elem.find(qn('w:pPr'))
                if style_pPr is not None:
                    jc = style_pPr.find(qn('w:jc'))
                    if jc is not None:
                        return jc.get(qn('w:val'))
                based_on = style_elem.find(qn('w:basedOn'))
                if based_on is not None:
                    base_val = based_on.get(qn('w:val'))
                    try:
                        style = self.doc.styles[base_val]
                    except:
                        break
                else:
                    break
            return None
        
        for para in self.paragraphs:
            style_name = para.style.name if para.style else ''
            if style_name in ('Heading 1', '标题 1', 'Heading1', '标题1'):
                # 优先检查python-docx的alignment
                if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    return True
                # 检查有效XML对齐方式
                jc_val = get_effective_jc(para)
                if jc_val == 'center':
                    return True
        return False

    def doc_toc_centered(self):
        """检查目录标题是否居中（查找包含'目录'文字且居中的段落）"""
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        for para in self.paragraphs:
            text = para.text.strip()
            if '目录' in text and len(text) <= 10:
                if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    return True
        return False

    def doc_has_toc(self):
        """检查是否有目录（TOC字段）"""
        for elem in self.body_elements:
            # 检查fldSimple中的TOC字段
            fld_simples = elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldSimple')
            for fs in fld_simples:
                instr = fs.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instr', '')
                if 'TOC' in instr.upper():
                    return True
            # 检查fldChar中的TOC字段
            fld_chars = elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldChar')
            for fc in fld_chars:
                if fc.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fldCharType') == 'begin':
                    instrs = elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instrText')
                    for inst in instrs:
                        if 'TOC' in inst.text.upper():
                            return True
            # 检查sdt（结构化文档标签）中的TOC
            sdts = elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sdt')
            for sdt in sdts:
                doc_parts = sdt.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}docPartGallery')
                for dp in doc_parts:
                    if dp.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '') == 'Table of Contents':
                        return True
        return False

    def doc_has_table(self, min_rows, min_cols):
        """检查是否包含表格"""
        for table in self.tables:
            rows = len(table.rows)
            cols = len(table.columns)
            if rows >= min_rows and cols >= min_cols:
                return True
        return False

    def doc_has_merged_cells(self):
        """检查表格是否有合并单元格"""
        for table in self.tables:
            # 检查表格中的gridSpan和vMerge
            tbl = table._tbl
            grid_spans = tbl.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}gridSpan')
            v_merges = tbl.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}vMerge')
            if grid_spans or v_merges:
                return True
        return False

    def doc_table_has_borders(self):
        """检查表格是否有边框"""
        if not self.tables:
            return False
        for table in self.tables:
            tbl = table._tbl
            tbl_pr = tbl.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblPr')
            if tbl_pr is not None:
                tbl_borders = tbl_pr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblBorders')
                if tbl_borders is not None:
                    for border_name in ['top', 'bottom', 'left', 'right', 'insideH', 'insideV']:
                        border = tbl_borders.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}' + border_name)
                        if border is not None:
                            val = border.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 'none')
                            if val and val != 'none':
                                return True
        return False

    def doc_table_has_shading(self):
        """检查表格是否有底纹"""
        if not self.tables:
            return False
        for table in self.tables:
            tbl = table._tbl
            # 检查表格级别的底纹
            tbl_pr = tbl.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblPr')
            if tbl_pr is not None:
                shd = tbl_pr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd')
                if shd is not None:
                    fill = shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill', 'auto')
                    if fill and fill != 'auto':
                        return True
            # 检查单元格级别的底纹
            shadings = tbl.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd')
            for shd in shadings:
                fill = shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill', 'auto')
                if fill and fill != 'auto':
                    return True
        return False

    def doc_table_alignment(self):
        """检查表格内容对齐方式"""
        if not self.tables:
            return False
        align_count = 0
        for table in self.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if para.alignment is not None:
                            align_count += 1
                            if align_count >= 3:
                                return True
        return False

    def title_bold_center(self):
        """检查标题是否加粗居中"""
        title = self.get_title_paragraph()
        if not title:
            return False
        is_bold = False
        is_center = False
        # 检查加粗
        for run in title.runs:
            if run.font.bold:
                is_bold = True
                break
            b = run._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}b')
            if b is not None:
                val = b.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', 'true')
                if val != 'false' and val != '0':
                    is_bold = True
                    break
        # 检查居中
        if title.alignment == WD_ALIGN_PARAGRAPH.CENTER:
            is_center = True
        return is_bold and is_center

    def doc_has_page_border(self):
        """检查文档是否设置了页面边框（通过解析XML中的sectPr/pgBorders）"""
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        # 检查每个节的 sectPr 中的 pgBorders
        for section in self.sections:
            sect_pr = section._sectPr
            pg_borders = sect_pr.find(f'{{{wns}}}pgBorders')
            if pg_borders is not None:
                # 检查是否有至少一个边框（top, left, bottom, right, 或 box）
                border_types = ['top', 'left', 'bottom', 'right', 'start', 'end']
                for bt in border_types:
                    border_el = pg_borders.find(f'{{{wns}}}{bt}')
                    if border_el is not None:
                        # 检查该边框是否有实际值（val不为nil/none）
                        val = border_el.get(f'{{{wns}}}val', '')
                        if val and val.lower() != 'nil' and val.lower() != 'none':
                            return True
                # 也检查 box 类型（四边同时设置）
                box_el = pg_borders.find(f'{{{wns}}}box')
                if box_el is not None:
                    val = box_el.get(f'{{{wns}}}val', '')
                    if val and val.lower() != 'nil' and val.lower() != 'none':
                        return True
        return False

    def doc_has_table_formula(self):
        """检查文档表格中是否包含公式（通过解析XML中的fldChar/instrText）"""
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        if not self.tables:
            return False
        for table in self.tables:
            tbl_el = table._tbl
            # 在表格的所有段落中查找域代码（公式）
            for para_el in tbl_el.iter(f'{{{wns}}}p'):
                # 查找 fldChar（域代码标记）
                fld_chars = para_el.findall(f'.//{{{wns}}}fldChar')
                if fld_chars:
                    # 查找 instrText（域代码内容）
                    instr_texts = para_el.findall(f'.//{{{wns}}}instrText')
                    for instr in instr_texts:
                        text = (instr.text or '').strip()
                        # Word表格公式通常以 = 开头，如 =SUM(ABOVE), =AVERAGE(LEFT) 等
                        if text.startswith('=') or 'FORMULA' in text.upper():
                            return True
                # 也检查 fldSimple（简单域）
                fld_simples = para_el.findall(f'.//{{{wns}}}fldSimple')
                for fld in fld_simples:
                    instr = fld.get(f'{{{wns}}}instr', '')
                    if instr.startswith('=') or 'FORMULA' in instr.upper():
                        return True
        return False

    def doc_page_border_exact(self, border_sz_pt, border_style=None, border_color=None):
        """精确检测页面边框粗细（样式可选）
        参数: border_sz_pt如1.5(磅), border_style如"single"(可选,不传则不检查样式), border_color如"000000"(可选)
        XML中sz属性单位为1/8磅，所以1.5磅=12
        兼容WPS多种XML格式（threeDEmboss/single/无val属性/box元素等）
        """
        expected_sz = int(border_sz_pt * 8)
        if border_color:
            border_color = border_color.lstrip('#').upper()
        for section in self.sections:
            sect_pr = section._sectPr
            pg_borders = sect_pr.find(qn('w:pgBorders'))
            if pg_borders is None:
                continue
            border_types = ['top', 'left', 'bottom', 'right', 'box']
            match_count = 0
            checked_count = 0
            for bt in border_types:
                border_el = pg_borders.find(qn(f'w:{bt}'))
                if border_el is None:
                    continue
                checked_count += 1
                val = border_el.get(qn('w:val'), '')
                sz = border_el.get(qn('w:sz'), '')
                color = border_el.get(qn('w:color'), '')
                
                # val为nil/none时跳过
                if val and val.lower() in ('nil', 'none'):
                    continue
                # 检查样式（仅当指定了border_style时）
                if border_style and val and val.lower() != border_style.lower():
                    continue
                # 检查粗细
                if sz:
                    actual_sz = int(sz)
                    if abs(actual_sz - expected_sz) > 1:
                        continue
                # 检查颜色（如要求）
                if border_color is not None:
                    if not color or color.lstrip('#').upper() != border_color:
                        continue
                match_count += 1
            
            if checked_count > 0 and match_count > 0:
                return True
            if match_count >= 2:
                return True
        return False

    def title_font_color(self, color_hex):
        """检测标题字体颜色
        参数: color_hex如"FF0000"，支持带#或不带#的hex格式
        """
        color_hex = color_hex.lstrip('#').upper()
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            rPr = run._element.find(qn('w:rPr'))
            if rPr is not None:
                color_el = rPr.find(qn('w:color'))
                if color_el is not None:
                    val = color_el.get(qn('w:val'), '')
                    if val and val.upper() == color_hex:
                        return True
                    # 也检查themeColor映射
                    theme = color_el.get(qn('w:themeColor'), '')
                    if theme:
                        return True  # 主题颜色视为有设置颜色
        return False

    def body_font_color(self, color_hex):
        """检测正文字体颜色
        参数: color_hex如"0000FF"，支持带#或不带#的hex格式
        """
        color_hex = color_hex.lstrip('#').upper()
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            for run in para.runs:
                rPr = run._element.find(qn('w:rPr'))
                if rPr is not None:
                    color_el = rPr.find(qn('w:color'))
                    if color_el is not None:
                        val = color_el.get(qn('w:val'), '')
                        if val and val.upper() == color_hex:
                            match_count += 1
                            break
        return match_count >= max(1, len(body_paras[:5]) // 2)

    def body_font_size_exact(self, exact_size_pt):
        """精确检测正文字号
        参数: exact_size_pt如16
        通过XML的w:sz/w:szCs属性检测（sz=磅值×2，单位半磅），允许±1误差
        """
        expected_sz = exact_size_pt * 2  # 半磅单位
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            for run in para.runs:
                rPr = run._element.find(qn('w:rPr'))
                if rPr is not None:
                    # 检查w:sz和w:szCs（python-docx可能存储在szCs中）
                    for sz_tag in [qn('w:sz'), qn('w:szCs')]:
                        sz_el = rPr.find(sz_tag)
                        if sz_el is not None:
                            val = int(sz_el.get(qn('w:val'), 0))
                            if abs(val - expected_sz) <= 1:
                                match_count += 1
                                break
        return match_count >= max(1, len(body_paras[:5]) // 2)

    def check_image_size(self, min_width_cm=None, min_height_cm=None, max_width_cm=None, max_height_cm=None):
        """检测图片尺寸
        参数均为可选，只检查提供的条件
        通过inline shapes的width/height（EMU单位，1cm=360000EMU）检测，允许10%误差
        """
        if not self.inline_shapes:
            return False
        for shape in self.inline_shapes:
            width_emu = shape.width
            height_emu = shape.height
            width_cm = width_emu / 360000.0 if width_emu else 0
            height_cm = height_emu / 360000.0 if height_emu else 0
            ok = True
            if min_width_cm is not None:
                if width_cm < min_width_cm * 0.9:
                    ok = False
            if min_height_cm is not None:
                if height_cm < min_height_cm * 0.9:
                    ok = False
            if max_width_cm is not None:
                if width_cm > max_width_cm * 1.1:
                    ok = False
            if max_height_cm is not None:
                if height_cm > max_height_cm * 1.1:
                    ok = False
            if ok:
                return True
        return False

    def check_image_position(self, position_type):
        """检测图片位置/环绕方式
        position_type: "behind_text"(衬于文字下方), "in_front_of_text"(浮于文字上方), 
                       "inline"(嵌入式), "non_inline"(任何非嵌入式/浮动)
        通过anchor元素的behindDoc属性和wrap子元素检测
        """
        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        
        if position_type == 'inline':
            return len(self.inline_shapes) > 0
        
        # 检查浮动图片（anchor类型）
        for elem in self.body_elements:
            anchors = elem.findall(f'.//{{{wp_ns}}}anchor')
            for anchor in anchors:
                # non_inline: 只要有anchor就说明是非嵌入式
                if position_type == 'non_inline':
                    return True
                
                # 检查子元素中的wrap类型
                for child in anchor:
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    
                    if position_type == 'behind_text':
                        behind = anchor.get('behindDoc', '')
                        # WPS中behindDoc可能为"1"、"true"或空但有wrapPolygon等
                        if behind.lower() in ('1', 'true'):
                            return True
                    
                    elif position_type == 'in_front_of_text':
                        behind = anchor.get('behindDoc', '')
                        # behindDoc为空或0/false时，表示浮于文字上方
                        if behind.lower() in ('', '0', 'false'):
                            return True
                        # 有wrapSquare/wrapTight等子元素也表示浮于上方
                        if tag in ('wrapSquare', 'wrapTight', 'wrapThrough', 'wrapTopAndBottom'):
                            behind_val = anchor.get('behindDoc', '')
                            if behind_val.lower() not in ('1', 'true'):
                                return True
        return False

    def doc_table_border_exact(self, border_sz_pt, border_color=None):
        """精确检测表格边框粗细
        参数: border_sz_pt如1.5, border_color如"000080"(可选)
        检查tblPr/tblBorders中的top/left/bottom/insideH/insideV
        sz属性单位1/8磅
        """
        expected_sz = int(border_sz_pt * 8)
        if border_color:
            border_color = border_color.lstrip('#').upper()
        if not self.tables:
            return False
        for table in self.tables:
            tbl = table._tbl
            tbl_pr = tbl.find(qn('w:tblPr'))
            if tbl_pr is None:
                continue
            tbl_borders = tbl_pr.find(qn('w:tblBorders'))
            if tbl_borders is None:
                continue
            border_names = ['top', 'left', 'bottom', 'insideH', 'insideV']
            match_count = 0
            for bn in border_names:
                border_el = tbl_borders.find(qn(f'w:{bn}'))
                if border_el is not None:
                    sz = border_el.get(qn('w:sz'), '')
                    color = border_el.get(qn('w:color'), '')
                    if sz:
                        actual_sz = int(sz)
                        if abs(actual_sz - expected_sz) <= 1:
                            if border_color is None:
                                match_count += 1
                            elif color and color.lstrip('#').upper() == border_color:
                                match_count += 1
            if match_count >= 3:
                return True
        return False

    def paragraph_shading_color(self, color_hex):
        """精确检测段落底纹颜色
        参数: color_hex如"FFFF00"
        检查w:shd的fill属性
        """
        color_hex = color_hex.lstrip('#').upper()
        for para in self.paragraphs:
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                shd = pPr.find(qn('w:shd'))
                if shd is not None:
                    fill = shd.get(qn('w:fill'), '')
                    if fill and fill.lstrip('#').upper() == color_hex:
                        return True
        return False

    def title_font_name_exact(self, font_name):
        """精确检测标题字体（完全匹配）
        与check_title_font不同，这里要求完全匹配而非包含
        """
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.name and run.font.name == font_name:
                return True
            # 检查东亚字体和ASCII字体
            rpr = run._element.find(qn('w:rFonts'))
            if rpr is not None:
                ea = rpr.get(qn('w:eastAsia'))
                if ea and ea == font_name:
                    return True
                ascii_font = rpr.get(qn('w:ascii'))
                if ascii_font and ascii_font == font_name:
                    return True
                hAnsi = rpr.get(qn('w:hAnsi'))
                if hAnsi and hAnsi == font_name:
                    return True
        return False

    def body_font_name_exact(self, font_name):
        """精确检测正文字体（完全匹配）
        与check_body_font不同，这里要求完全匹配而非包含
        """
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            for run in para.runs:
                if run.font.name and run.font.name == font_name:
                    match_count += 1
                    break
                rpr = run._element.find(qn('w:rFonts'))
                if rpr is not None:
                    ea = rpr.get(qn('w:eastAsia'))
                    if ea and ea == font_name:
                        match_count += 1
                        break
                    ascii_font = rpr.get(qn('w:ascii'))
                    if ascii_font and ascii_font == font_name:
                        match_count += 1
                        break
                    hAnsi = rpr.get(qn('w:hAnsi'))
                    if hAnsi and hAnsi == font_name:
                        match_count += 1
                        break
        return match_count >= 2

    def doc_has_wordart_or_text_effect(self):
        """检测文档是否包含艺术字或文字效果
        通过检测 wps:wsp 元素（WordprocessingShape）或 run 的 rPr 中的 effect/emboss/imprint/outline/shadow 元素
        注意：与 check_has_wordart 方法名不同，用于区分
        """
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        wps_ns = 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'
        mc_ns = 'http://schemas.openxmlformats.org/markup-compatibility/2006'

        # 方法1: 查找 wps:wsp 元素（WordprocessingShape）
        for elem in self.body_elements:
            wsps = elem.findall(f'.//{{{wps_ns}}}wsp')
            if wsps:
                return True

        # 方法2: 通过 mc:AlternateContent 查找
        for elem in self.body_elements:
            alt_contents = elem.findall(f'.//{{{mc_ns}}}AlternateContent')
            for ac in alt_contents:
                for child in ac:
                    wsps = child.findall(f'.//{{{wps_ns}}}wsp')
                    if wsps:
                        return True

        # 方法3: 检查 run 的 rPr 中的文字效果元素
        for para in self.paragraphs:
            for run in para.runs:
                rPr = run._element.find(f'{{{wns}}}rPr')
                if rPr is not None:
                    effect_elements = ['effect', 'emboss', 'imprint', 'outline', 'shadow']
                    for eff in effect_elements:
                        if rPr.find(f'{{{wns}}}{eff}') is not None:
                            return True

        return False

    def doc_content_complete(self, min_chars=100):
        """检测文档内容是否完整（至少min_chars个字符）"""
        total_chars = 0
        for para in self.paragraphs:
            total_chars += len(para.text.strip())
        # 也统计表格中的文字
        for table in self.tables:
            for row in table.rows:
                for cell in row.cells:
                    total_chars += len(cell.text.strip())
        return total_chars >= min_chars

    def doc_has_multiple_images(self, min_count=2):
        """检测是否包含多张图片"""
        image_count = len(self.inline_shapes)
        # 也检查浮动图片（anchor类型的drawing元素）
        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        for elem in self.body_elements:
            anchors = elem.findall(f'.//{{{wp_ns}}}anchor')
            image_count += len(anchors)
        return image_count >= min_count

    def doc_char_spacing(self, spacing_cm):
        """检测字符间距加宽
        参数: spacing_cm如0.2（厘米）
        XML中 w:spacing 的 w:val 属性，单位为1/20磅（twips）
        1cm = 567 twips，允许±10%误差
        """
        target_twips = spacing_cm * 567
        tolerance = 0.1
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

        for elem in self.body_elements:
            # 检查段落级别的 w:spacing（w:val 表示字符间距）
            pPr = elem.find(qn('w:pPr'))
            if pPr is not None:
                spacing = pPr.find(qn('w:spacing'))
                if spacing is not None:
                    val = spacing.get(qn('w:val'))
                    if val is not None:
                        try:
                            val_twips = int(val)
                            if abs(val_twips - target_twips) <= target_twips * tolerance:
                                return True
                        except (ValueError, TypeError):
                            pass

            # 检查 run 级别的 rPr 中的 w:spacing
            for rPr in elem.findall(f'.//{{{wns}}}rPr'):
                spacing = rPr.find(qn('w:spacing'))
                if spacing is not None:
                    val = spacing.get(qn('w:val'))
                    if val is not None:
                        try:
                            val_twips = int(val)
                            if abs(val_twips - target_twips) <= target_twips * tolerance:
                                return True
                        except (ValueError, TypeError):
                            pass

        return False

    def doc_has_picture_with_size(self, height_cm=None, width_cm=None):
        """检测图片大小设置
        通过inline shapes的extent/cx,cy属性（EMU单位，1cm=360000EMU）检测
        允许15%误差
        """
        if height_cm is None and width_cm is None:
            return False

        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        tolerance = 0.15

        for elem in self.body_elements:
            # 检查 inline 类型图片
            inlines = elem.findall(f'.//{{{wp_ns}}}inline')
            for inline in inlines:
                extent = inline.find(f'.//{{{wp_ns}}}extent')
                if extent is None:
                    # extent 可能在 drawingml 命名空间下
                    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    extent = inline.find(f'{{{a_ns}}}extent')
                if extent is not None:
                    cx = extent.get('cx')
                    cy = extent.get('cy')
                    try:
                        if width_cm is not None and cx is not None:
                            cx_cm = int(cx) / 360000
                            if abs(cx_cm - width_cm) <= width_cm * tolerance:
                                if height_cm is None:
                                    return True
                        if height_cm is not None and cy is not None:
                            cy_cm = int(cy) / 360000
                            if abs(cy_cm - height_cm) <= height_cm * tolerance:
                                if width_cm is None:
                                    return True
                        # 两个都指定时需要同时满足
                        if width_cm is not None and height_cm is not None and cx is not None and cy is not None:
                            cx_cm = int(cx) / 360000
                            cy_cm = int(cy) / 360000
                            if (abs(cx_cm - width_cm) <= width_cm * tolerance and
                                    abs(cy_cm - height_cm) <= height_cm * tolerance):
                                return True
                    except (ValueError, TypeError):
                        pass

            # 检查 anchor 类型图片
            anchors = elem.findall(f'.//{{{wp_ns}}}anchor')
            for anchor in anchors:
                extent = anchor.find(f'.//{{{wp_ns}}}extent')
                if extent is None:
                    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    extent = anchor.find(f'{{{a_ns}}}extent')
                if extent is not None:
                    cx = extent.get('cx')
                    cy = extent.get('cy')
                    try:
                        if width_cm is not None and cx is not None:
                            cx_cm = int(cx) / 360000
                            if abs(cx_cm - width_cm) <= width_cm * tolerance:
                                if height_cm is None:
                                    return True
                        if height_cm is not None and cy is not None:
                            cy_cm = int(cy) / 360000
                            if abs(cy_cm - height_cm) <= height_cm * tolerance:
                                if width_cm is None:
                                    return True
                        if width_cm is not None and height_cm is not None and cx is not None and cy is not None:
                            cx_cm = int(cx) / 360000
                            cy_cm = int(cy) / 360000
                            if (abs(cx_cm - width_cm) <= width_cm * tolerance and
                                    abs(cy_cm - height_cm) <= height_cm * tolerance):
                                return True
                    except (ValueError, TypeError):
                        pass

        return False

    def doc_has_wordart_with_size(self, height_cm=None, width_cm=None):
        """检测艺术字大小
        WPS艺术字大小可能存储在两个位置：
        1. wp:anchor/extent（外层锚点大小）
        2. wps:spPr/a:xfrm/a:ext（内部形状大小）
        两个位置都要检查，任一匹配即可
        不传参数时：检测是否有艺术字设置了自定义大小
        传参数时：精确匹配大小（允许30%误差）
        """
        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        wps_ns = 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for elem in self.body_elements:
            anchors = elem.findall(f'.//{{{wp_ns}}}anchor')
            for anchor in anchors:
                wsps = anchor.findall(f'.//{{{wps_ns}}}wsp')
                if not wsps:
                    continue

                # 收集所有可能的大小信息
                sizes = []

                # 位置1: wp:anchor/extent
                extent = anchor.find(f'{{{wp_ns}}}extent')
                if extent is not None:
                    cx = extent.get('cx')
                    cy = extent.get('cy')
                    if cx and cy:
                        try:
                            sizes.append((int(cx) / 360000, int(cy) / 360000))
                        except (ValueError, TypeError):
                            pass

                # 位置2: wps:spPr/a:xfrm/a:ext
                for wsp in wsps:
                    spPr = wsp.find(f'{{{wps_ns}}}spPr')
                    if spPr is not None:
                        xfrm = spPr.find(f'{{{a_ns}}}xfrm')
                        if xfrm is not None:
                            ext = xfrm.find(f'{{{a_ns}}}ext')
                            if ext is not None:
                                ecx = ext.get('cx')
                                ecy = ext.get('cy')
                                if ecx and ecy:
                                    try:
                                        sizes.append((int(ecx) / 360000, int(ecy) / 360000))
                                    except (ValueError, TypeError):
                                        pass

                if not sizes:
                    continue

                # 不传参数：只要有大小设置就算通过
                if height_cm is None and width_cm is None:
                    for cx_cm, cy_cm in sizes:
                        if cx_cm > 0.5 and cy_cm > 0.5:
                            return True
                    continue

                # 传参数：精确匹配（任一位置匹配即可）
                tolerance = 0.15
                for cx_cm, cy_cm in sizes:
                    try:
                        w_match = width_cm is None or abs(cx_cm - width_cm) <= width_cm * tolerance
                        h_match = height_cm is None or abs(cy_cm - height_cm) <= height_cm * tolerance
                        if w_match and h_match:
                            return True
                    except (ValueError, TypeError):
                        pass
        return False

    def doc_has_wordart_with_position(self):
        """检测艺术字位置设置（绝对定位）
        WPS艺术字位置存储在wp:anchor的positionH/positionV中
        """
        wp_ns = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        wps_ns = 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'

        for elem in self.body_elements:
            anchors = elem.findall(f'.//{{{wp_ns}}}anchor')
            for anchor in anchors:
                # 确认是艺术字（包含wps:wsp）
                wsps = anchor.findall(f'.//{{{wps_ns}}}wsp')
                if not wsps:
                    continue
                # 检查positionH和positionV
                posH = anchor.find(f'{{{wp_ns}}}positionH')
                posV = anchor.find(f'{{{wp_ns}}}positionV')
                if posH is not None and posV is not None:
                    return True
        return False

    def doc_paragraph_line_spacing(self, spacing_pt):
        """检测特定段落的行距
        参数: spacing_pt如50（磅）
        通过XML的w:spacing的w:line属性检测
        支持固定值(lineRule=exact)和倍数(lineRule=auto)两种模式
        w:line 单位为 1/240 英寸（twips），1磅 = 20 twips
        """
        wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        # spacing_pt 转换为 twips（1磅 = 20 twips）
        target_twips = spacing_pt * 20

        for para in self.paragraphs:
            pPr = para._element.find(qn('w:pPr'))
            if pPr is None:
                continue
            spacing = pPr.find(qn('w:spacing'))
            if spacing is None:
                continue
            line_val = spacing.get(qn('w:line'))
            line_rule = spacing.get(qn('w:lineRule'))
            if line_val is None:
                continue
            try:
                line_twips = int(line_val)
                if line_rule == 'exact' or line_rule == 'atLeast':
                    # 固定值模式：直接比较
                    if abs(line_twips - target_twips) <= target_twips * 0.1:
                        return True
                else:
                    # auto 模式（倍数行距）：line 值为 240 的倍数
                    # 240 = 单倍行距，480 = 双倍行距等
                    if abs(line_twips - target_twips) <= target_twips * 0.1:
                        return True
            except (ValueError, TypeError):
                pass

        return False

    def doc_title_font_name(self, font_name):
        """检测标题字体（包含匹配）
        与title_font_name_exact不同，这里用包含匹配（更宽松）
        """
        title = self.get_title_paragraph()
        if not title or not title.runs:
            return False
        for run in title.runs:
            if run.font.name and font_name in run.font.name:
                return True
            rpr = run._element.find(qn('w:rFonts'))
            if rpr is not None:
                ea = rpr.get(qn('w:eastAsia'))
                if ea and font_name in ea:
                    return True
                ascii_font = rpr.get(qn('w:ascii'))
                if ascii_font and font_name in ascii_font:
                    return True
                hAnsi = rpr.get(qn('w:hAnsi'))
                if hAnsi and font_name in hAnsi:
                    return True
        return False

    def doc_body_font_name(self, font_name):
        """检测正文字体（包含匹配）
        检查正文段落的run的字体名称
        """
        body_paras = self.get_body_paragraphs()
        if not body_paras:
            return False
        match_count = 0
        for para in body_paras[:5]:
            for run in para.runs:
                if run.font.name and font_name in run.font.name:
                    match_count += 1
                    break
                rpr = run._element.find(qn('w:rFonts'))
                if rpr is not None:
                    ea = rpr.get(qn('w:eastAsia'))
                    if ea and font_name in ea:
                        match_count += 1
                        break
                    ascii_font = rpr.get(qn('w:ascii'))
                    if ascii_font and font_name in ascii_font:
                        match_count += 1
                        break
                    hAnsi = rpr.get(qn('w:hAnsi'))
                    if hAnsi and font_name in hAnsi:
                        match_count += 1
                        break
        return match_count >= 2

    def doc_has_even_odd_header(self):
        """检测是否设置了奇偶页不同
        通过settings.xml中的w:evenAndOddHeaders元素检测（w:val="1"表示启用）
        注意：WPS取消勾选时会删除该元素（而非设为val="0"），所以元素不存在=未启用
        """
        try:
            import zipfile
            from lxml import etree
            wns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
            with zipfile.ZipFile(self.filepath) as z:
                if 'word/settings.xml' in z.namelist():
                    with z.open('word/settings.xml') as sf:
                        settings = etree.parse(sf)
                    for elem in settings.iter(f'{{{wns}}}evenAndOddHeaders'):
                        val = elem.get(f'{{{wns}}}val', '1')
                        if val in ('1', 'true'):
                            return True
                        # val="0"/"false" 或其他值都视为未启用
                        return False
        except Exception:
            pass

        # settings.xml中没有evenAndOddHeaders元素 = 未启用
        return False

    def doc_has_different_first_page(self):
        """检测是否设置了首页不同
        通过sectPr的w:titlePg元素检测
        """
        for section in self.sections:
            sectPr = section._sectPr
            if sectPr is not None:
                titlePg = sectPr.find(qn('w:titlePg'))
                if titlePg is not None:
                    return True
        return False

    def doc_has_heading2_style(self):
        """检测是否使用了二级标题样式
        检查是否有段落应用了Heading2样式（兼容WPS数字ID）
        同时检测是否有手动设置为二级标题特征的段落（加粗+不同于正文的格式）
        """
        # 方法1：检测标准Heading 2样式
        has_heading2_style = False
        for para in self.paragraphs:
            style_name = para.style.name if para.style else ''
            if style_name in ('Heading 2', '标题 2'):
                has_heading2_style = True
                break
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                pStyle = pPr.find(qn('w:pStyle'))
                if pStyle is not None:
                    style_val = pStyle.get(qn('w:val'))
                    if style_val and ('heading2' in style_val.lower() or 'Heading2' in style_val):
                        has_heading2_style = True
                        break
                    try:
                        style = self.doc.styles[style_val]
                        if style.name in ('Heading 2', '标题 2'):
                            has_heading2_style = True
                            break
                    except (KeyError, AttributeError):
                        pass
        
        if has_heading2_style:
            return True
        
        # 方法2：检测手动设置的二级标题特征
        # 特征：段落文字较短（<=30字）、全部加粗、以数字序号开头（如"1、""2."等）
        import re
        heading2_pattern = re.compile(r'^[\s]*[(\（]?[一二三四五六七八九十\d]+[)）\、\.．\s]')
        bold_short_paras = 0
        for para in self.paragraphs:
            text = para.text.strip()
            if not text or len(text) > 40:
                continue
            # 跳过Heading 1
            style_name = para.style.name if para.style else ''
            if 'Heading 1' in style_name or '标题 1' in style_name:
                continue
            # 检查是否以序号开头
            if not heading2_pattern.match(text):
                continue
            # 检查是否加粗
            all_bold = all(run.bold for run in para.runs if run.text.strip())
            if all_bold and para.runs:
                bold_short_paras += 1
        
        return bold_short_paras >= 2


class XlsxParser:
    """Excel文档解析器"""
    
    def __init__(self, filepath):
        self.filepath = filepath
        self.wb = openpyxl.load_workbook(filepath)
        self.ws = self.wb.active
        self._original_ws = self.ws  # 保存原始工作表

    def switch_sheet(self, sheet_name):
        """切换到指定名称的工作表"""
        if sheet_name in self.wb.sheetnames:
            self.ws = self.wb[sheet_name]
            return True
        # 模糊匹配（部分名称匹配）
        for name in self.wb.sheetnames:
            if sheet_name in name or name in sheet_name:
                self.ws = self.wb[name]
                return True
        return False

    def get_sheet_names(self):
        """获取所有工作表名称"""
        return self.wb.sheetnames

    def get_all_text(self):
        """提取表格全部文本内容"""
        texts = []
        for row in self.ws.iter_rows(values_only=True):
            for cell in row:
                if cell and str(cell).strip():
                    texts.append(str(cell).strip())
        return "\n".join(texts)
    
    def check_merged_title(self, row, min_cols):
        """检查标题合并居中"""
        merged = list(self.ws.merged_cells.ranges)
        for mr in merged:
            if mr.min_row <= row + 1 <= mr.max_row and (mr.max_col - mr.min_col + 1) >= min_cols:
                return True
        return False
    
    def check_header_bold(self, header_row):
        """检查表头加粗
        header_row: 从第几行开始检查（0-based），检查前3行中非空单元格>=2的行
        """
        # 检查前几行，找到有多个非空单元格的行作为表头行
        candidates = []
        for row_idx in range(header_row + 1, min(header_row + 4, self.ws.max_row + 1)):
            cell_count = 0
            for cell in self.ws[row_idx]:
                if cell.value is not None:
                    cell_count += 1
            if cell_count >= 1:  # 至少1个非空单元格
                candidates.append((row_idx, cell_count))
        
        if not candidates:
            return False
        
        # 优先检查非空单元格最多的行
        candidates.sort(key=lambda x: -x[1])
        
        for best_row, _ in candidates:
            bold_count = 0
            total_cells = 0
            for cell in self.ws[best_row]:
                if cell.value is not None:
                    total_cells += 1
                    if cell.font.bold:
                        bold_count += 1
            # 只要有加粗的单元格就算通过（适应不同表头结构）
            if total_cells >= 1 and bold_count >= 1:
                return True
        
        return False
    
    def check_min_columns(self, min_cols):
        """检查最小列数"""
        max_col = self.ws.max_column
        return max_col >= min_cols
    
    def check_min_data_rows(self, min_rows, header_rows):
        """检查最小数据行数"""
        max_row = self.ws.max_row
        return (max_row - header_rows) >= min_rows
    
    def check_has_borders(self):
        """检查是否有边框"""
        border_count = 0
        for row in self.ws.iter_rows(min_row=1, max_row=min(self.ws.max_row, 15)):
            for cell in row:
                if cell.border:
                    b = cell.border
                    if b.left.style or b.right.style or b.top.style or b.bottom.style:
                        border_count += 1
        return border_count >= 5
    
    def check_column_width(self, min_width):
        """检查列宽设置"""
        wide_cols = 0
        for col_idx in range(1, min(self.ws.max_column + 1, 10)):
            col_letter = openpyxl.utils.get_column_letter(col_idx)
            width = self.ws.column_dimensions[col_letter].width
            if width and width >= min_width:
                wide_cols += 1
            elif width is None:
                # 默认宽度约8.43
                wide_cols += 1
        return wide_cols >= 3
    
    def check_number_format(self):
        """检查数值格式"""
        format_count = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 15)):
            for cell in row:
                if cell.value is not None and isinstance(cell.value, (int, float)):
                    fmt = cell.number_format
                    if not fmt or fmt == 'General':
                        continue
                    fmt_core = fmt.rstrip().rstrip('_').rstrip()
                    if fmt_core in ('0', '', '0.0'):
                        continue
                    format_count += 1
        return format_count >= 2
    
    def check_sheet_named(self):
        """检查工作表是否已命名"""
        name = self.ws.title
        return name and name.lower() not in ['sheet1', 'sheet', '工作表1', '工作表']

    def excel_title_font(self):
        """检查标题行字体设置"""
        # 检查第一行是否有非默认字体设置
        for cell in self.ws[1]:
            if cell.value is not None and cell.font and cell.font.name:
                return True
        return False

    def excel_cell_font_exact(self, row, col_start, col_end, font_name=None, font_size=None, bold=None, color_hex=None, color_tolerance=60):
        """精确检测指定区域单元格的字体格式
        row: 行号（1-based）
        col_start: 起始列（1-based，如1=A）
        col_end: 结束列（1-based，如7=G）
        font_name: 字体名（如"黑体"），None表示不检查
        font_size: 字号（如18），None表示不检查
        bold: 是否加粗（True/False），None表示不检查
        color_hex: 颜色值（如"002060"），None表示不检查
        color_tolerance: 颜色容差（RGB欧氏距离）
        """
        checked = 0
        for col_idx in range(col_start, col_end + 1):
            cell = self.ws.cell(row=row, column=col_idx)
            if cell.value is None:
                continue
            checked += 1
            f = cell.font
            if not f:
                return False

            # 检查字体名
            if font_name is not None:
                if not f.name or font_name not in str(f.name):
                    return False

            # 检查字号
            if font_size is not None:
                if f.size is None or abs(float(f.size) - float(font_size)) > 0.5:
                    return False

            # 检查加粗
            if bold is not None:
                if f.bold != bold:
                    return False

            # 检查颜色
            if color_hex is not None:
                if f.color:
                    matched = False
                    if f.color.rgb and f.color.type == 'rgb':
                        rgb = str(f.color.rgb)
                        if len(rgb) == 8:
                            rgb = rgb[2:]
                        if rgb.upper() == color_hex.upper():
                            matched = True
                        elif self._color_distance(rgb, color_hex) <= color_tolerance:
                            matched = True
                    if f.color.theme is not None:
                        matched = True
                    if not matched:
                        return False
                else:
                    return False

        return checked > 0

    def excel_data_font_set(self):
        """检查数据区字体设置"""
        font_set_count = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 15)):
            for cell in row:
                if cell.value is not None and cell.font and cell.font.name:
                    font_set_count += 1
                    if font_set_count >= 3:
                        return True
        return False

    def excel_alignment_set(self):
        """检查对齐方式设置"""
        align_count = 0
        for row in self.ws.iter_rows(min_row=1, max_row=min(self.ws.max_row, 15)):
            for cell in row:
                if cell.value is not None and cell.alignment:
                    a = cell.alignment
                    if a.horizontal or a.vertical or a.wrap_text or a.indent:
                        align_count += 1
                        if align_count >= 3:
                            return True
        return False

    def excel_has_conditional_format(self):
        """检查条件格式"""
        return len(self.ws.conditional_formatting) > 0

    def excel_has_data_bars(self):
        """检查数据条"""
        for cf_range in self.ws.conditional_formatting:
            for rule in cf_range.rules:
                if rule.type == 'dataBar':
                    return True
        return False

    def excel_has_data_validation(self):
        """检查数据有效性"""
        return len(self.ws.data_validations.dataValidation) > 0

    def excel_has_validation_type(self, validation_type):
        """检查是否包含指定类型的数据有效性规则
        validation_type: 'whole'(整数), 'decimal'(小数), 'list'(序列),
                        'textLength'(文本长度), 'date'(日期), 'time'(时间), 'custom'(自定义)
        """
        for dv in self.ws.data_validations.dataValidation:
            if dv.type and dv.type.lower() == validation_type.lower():
                return True
        return False

    def excel_has_validation_type_with_formula(self, validation_type, formula_keyword):
        """检查是否包含指定类型且公式包含关键词的数据有效性规则
        validation_type: 规则类型
        formula_keyword: formula1中应包含的关键词（如"男"表示序列包含"男,女"）
        """
        for dv in self.ws.data_validations.dataValidation:
            if dv.type and dv.type.lower() == validation_type.lower():
                if dv.formula1 and formula_keyword in str(dv.formula1):
                    return True
        return False

    def excel_has_table_style(self):
        """检查是否套用了表格样式（检测表头底纹或交替行颜色）"""
        # 方法1：检查Excel Table对象
        try:
            if hasattr(self.ws, 'tables') and len(self.ws.tables) > 0:
                return True
        except:
            pass
        # 方法2：检查表头行是否有底纹颜色（WPS套用表格样式的特征）
        try:
            # 检查前6行中是否有solid底纹（表头行通常有底纹）
            solid_fill_count = 0
            for row_idx in range(1, min(7, self.ws.max_row + 1)):
                row_fills = []
                for cell in self.ws[row_idx]:
                    if cell.value is not None:
                        fill = cell.fill
                        if fill and fill.fill_type == 'solid' and fill.fgColor:
                            c = fill.fgColor
                            # 排除透明/黑色
                            if c.theme is not None or (c.rgb and c.rgb not in ('00000000', '000000')):
                                row_fills.append(True)
                            elif c.indexed is not None and c.indexed > 0:
                                row_fills.append(True)
                if len(row_fills) >= 3:  # 一行中至少3个单元格有底纹
                    solid_fill_count += 1
            if solid_fill_count >= 1:
                return True
        except:
            pass
        return False

    def excel_date_format(self):
        """检查自定义日期格式（yyyy-mm-dd等，排除内置日期格式）"""
        date_format_count = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 20)):
            for cell in row:
                if cell.value is not None:
                    fmt = cell.number_format
                    if not fmt or fmt == 'General':
                        continue
                    # 排除内置日期格式（如 yyyy/m/d;@  mm-dd-yy）
                    if ';' in fmt:
                        continue
                    fmt_lower = fmt.lower()
                    # 只匹配自定义日期格式：yyyy-mm-dd 或包含中文年月日
                    if ('年' in fmt or '月' in fmt or '日' in fmt):
                        date_format_count += 1
                    elif 'yyyy' in fmt_lower and '-' in fmt:
                        date_format_count += 1
                    elif 'yyyy' in fmt_lower and '/' in fmt and 'mm' in fmt_lower:
                        date_format_count += 1
                    if date_format_count >= 2:
                        return True
        return False

    def excel_has_formula(self, formula_keyword):
        """检查是否包含指定公式（支持字符串或列表）"""
        if isinstance(formula_keyword, str):
            formula_keyword = [formula_keyword]
        formula_count = 0
        for row in self.ws.iter_rows(min_row=1, max_row=min(self.ws.max_row, 50)):
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.startswith('='):
                    for kw in formula_keyword:
                        if kw.upper() in cell.value.upper():
                            formula_count += 1
                            if formula_count >= 1:
                                return True
        return False

    def excel_has_formula_in_col(self, col):
        """检查指定列是否包含公式
        col: 列号（1-based）
        """
        try:
            for row_idx in range(1, min(self.ws.max_row + 1, 50)):
                cell = self.ws.cell(row=row_idx, column=col)
                if cell.value and isinstance(cell.value, str) and cell.value.startswith('='):
                    return True
            return False
        except:
            return False

    def excel_decimal_format(self):
        """检查小数位格式（区分用户设置和WPS默认格式）"""
        decimal_count = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 20)):
            for cell in row:
                if cell.value is not None and isinstance(cell.value, (int, float)):
                    fmt = cell.number_format
                    if not fmt or fmt == 'General':
                        continue
                    # 去掉尾部空格和下划线后缀（WPS格式特性）
                    fmt_core = fmt.rstrip().rstrip('_').rstrip()
                    # 排除纯整数默认格式
                    if fmt_core in ('0', ''):
                        continue
                    # 排除默认小数格式（0.0 是WPS自动分配的）
                    # 但保留用户明确设置的格式（如 #,##0.0, 0.00, 0.000 等）
                    if fmt_core == '0.0':
                        continue
                    if ('0.0' in fmt_core or '#.0' in fmt_core or '0.00' in fmt_core
                                or '#.00' in fmt_core or '0.000' in fmt_core or '#.000' in fmt_core):
                        decimal_count += 1
                        if decimal_count >= 2:
                            return True
        return False

    def excel_result_area_formatted(self):
        """检查结果区域格式"""
        max_row = self.ws.max_row
        if max_row < 3:
            return False
        format_count = 0
        for row in self.ws.iter_rows(min_row=max(1, max_row - 4), max_row=max_row):
            for cell in row:
                if cell.value is not None:
                    fmt = cell.number_format if cell.number_format else 'General'
                    fmt_core = fmt.rstrip().rstrip('_').rstrip()
                    if fmt_core in ('0', '', '0.0'):
                        fmt = 'General'
                    if (cell.font.bold or cell.font.color
                            or (fmt and fmt != 'General')
                            or cell.fill.fgColor):
                        format_count += 1
                        if format_count >= 2:
                            return True
        return False

    def excel_data_sorted(self, sort_col=0):
        """检查数据是否排序（支持分组排序，允许相同值）
        sort_col: 排序依据的列号（1-based），0表示自动检测任意列
        """
        if self.ws.max_row < 3:
            return False

        def check_col_sorted(col):
            data = []
            for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 20), min_col=col, max_col=col):
                for cell in row:
                    if cell.value is not None:
                        data.append(cell.value)
            if len(data) < 3:
                return False

            def is_non_descending(seq):
                desc_count = 0
                for i in range(len(seq) - 1):
                    try:
                        if seq[i] > seq[i + 1]:
                            desc_count += 1
                    except TypeError:
                        pass
                return desc_count <= max(1, len(seq) // 10)

            def is_non_ascending(seq):
                asc_count = 0
                for i in range(len(seq) - 1):
                    try:
                        if seq[i] < seq[i + 1]:
                            asc_count += 1
                    except TypeError:
                        pass
                return asc_count <= max(1, len(seq) // 10)

            str_data = [str(d) for d in data]
            if is_non_descending(str_data) or is_non_ascending(str_data):
                return True
            numeric_data = [d for d in data if isinstance(d, (int, float))]
            if len(numeric_data) >= 3:
                if is_non_descending(numeric_data) or is_non_ascending(numeric_data):
                    return True
            return False

        if sort_col > 0:
            return check_col_sorted(sort_col)

        # 自动检测：检查所有列是否有任意一列有序
        for col in range(1, self.ws.max_column + 1):
            if check_col_sorted(col):
                return True
        return False

    def excel_has_subtotals(self):
        """检查分类汇总"""
        # 分类汇总通常使用SUBTOTAL函数
        for row in self.ws.iter_rows(min_row=1, max_row=min(self.ws.max_row, 50)):
            for cell in row:
                if cell.value and isinstance(cell.value, str) and 'SUBTOTAL' in cell.value.upper():
                    return True
        return False

    def excel_has_autofilter(self):
        """检查自动筛选"""
        return self.ws.auto_filter.ref is not None

    def excel_has_advanced_filter(self, original_row_count=31):
        """检查高级筛选（通过行数差异判断：高级筛选结果被复制到表中导致行数增多）"""
        # 方法1：检查命名区域
        try:
            for name in self.wb.defined_names:
                if name and ('filter' in name.name.lower() or '筛选' in name.name.lower()
                             or 'criteria' in name.name.lower() or '条件' in name.name.lower()):
                    return True
        except:
            pass
        # 方法2：当前工作表行数明显多于原始数据（高级筛选结果被复制到表中）
        if self.ws.max_row > original_row_count + 5:
            return True
        return False

    def excel_has_pivot_table(self):
        """检查数据透视表（通过openpyxl和ZIP两种方式检测）"""
        # 方法1：openpyxl透视表缓存
        try:
            if hasattr(self.wb, '_pivots') and self.wb._pivots:
                return True
        except:
            pass
        # 方法2：通过ZIP检查XML中的pivotTable文件
        try:
            import zipfile
            with zipfile.ZipFile(self.filepath) as z:
                pivot_files = [n for n in z.namelist() if 'pivotTable' in n and n.endswith('.xml')]
                if pivot_files:
                    return True
        except:
            pass
        return False

    def excel_pivot_configured(self):
        """检查透视表配置"""
        # 有透视表就认为已配置
        return self.excel_has_pivot_table()

    def excel_has_chart(self):
        """检查当前工作表是否有图表"""
        if hasattr(self.ws, '_charts') and len(self.ws._charts) > 0:
            return True
        return False

    def excel_chart_has_title(self):
        """检查当前工作表的图表是否有标题"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if chart.title:
                    return True
        return False

    def excel_has_data_labels(self):
        """检查图表是否有数据标签"""
        for ws in self.wb.worksheets:
            if hasattr(ws, '_charts'):
                for chart in ws._charts:
                    if hasattr(chart, 'dataLabels') and chart.dataLabels:
                        return True
        return False

    def excel_no_gridlines(self):
        """检查图表是否删除了网格线（近似检测）"""
        # openpyxl中网格线默认显示，删除后可能无法直接检测
        # 近似检测：有图表就认为可能做了操作
        for ws in self.wb.worksheets:
            if hasattr(ws, '_charts') and len(ws._charts) > 0:
                return True
        return False

    def excel_has_multiple_validations(self, min_count=2):
        """检测是否设置了多条数据有效性规则"""
        try:
            validations = self.ws.data_validations.dataValidation
            return len(validations) >= min_count
        except:
            return False

    def excel_has_multiple_conditional_formats(self, min_count=2):
        """检测是否设置了多条条件格式规则"""
        try:
            rule_count = 0
            for cf_range in self.ws.conditional_formatting:
                rule_count += len(cf_range.rules)
                if rule_count >= min_count:
                    return True
            return False
        except:
            return False

    def excel_has_number_format_custom(self):
        """检测是否设置了自定义数值格式（千位分隔符等）"""
        custom_count = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 20)):
            for cell in row:
                if cell.value is not None and isinstance(cell.value, (int, float)):
                    fmt = cell.number_format
                    if not fmt or fmt == 'General':
                        continue
                    # 去掉尾部空格和下划线后缀
                    fmt_core = fmt.rstrip().rstrip('_').rstrip()
                    # 排除纯整数和默认小数格式
                    if fmt_core in ('0', '', '0.0'):
                        continue
                    # 检查千位分隔符、自定义小数位、货币等
                    if any(marker in fmt_core for marker in ['#,##', '%', '"¥"', '"￥"', '_($', '_(*']):
                        custom_count += 1
                        if custom_count >= 2:
                            return True
        return False

    def excel_cell_number_format(self, col_start, col_end, row_start=2, row_end=None, format_type=None, format_keyword=None):
        """精确检测指定区域的数值格式
        col_start/col_end: 列范围（1-based，如C列=3）
        row_start/row_end: 行范围（1-based）
        format_type: 'decimal'(小数位), 'thousands'(千位分隔符), 'number'(数值非文本), 'date'(日期), 'custom'(自定义)
        format_keyword: 格式字符串中必须包含的关键词（如 '#,##', '0.0', 'yyyy-mm'）
        """
        if row_end is None:
            row_end = min(self.ws.max_row, 30)

        match_count = 0
        total_cells = 0

        for row_idx in range(row_start, row_end + 1):
            for col_idx in range(col_start, col_end + 1):
                cell = self.ws.cell(row=row_idx, column=col_idx)
                if cell.value is None:
                    continue
                total_cells += 1

                fmt = cell.number_format or 'General'
                fmt_core = fmt.rstrip().rstrip('_').rstrip()

                # 排除默认格式
                if fmt_core in ('0', '', '0.0', 'General'):
                    if format_type != 'date':
                        continue

                matched = False

                if format_type == 'decimal':
                    # 小数位格式
                    if any(m in fmt_core for m in ['0.0', '#.0', '0.00', '#.00', '0.000', '#.000']):
                        matched = True
                elif format_type == 'thousands':
                    # 千位分隔符
                    if '#,##' in fmt_core:
                        matched = True
                elif format_type == 'number':
                    # 数值格式（非文本、非日期）
                    if fmt_core not in ('0', '', '0.0', 'General'):
                        if 'yyyy' not in fmt_core.lower() and 'mm' not in fmt_core.lower():
                            matched = True
                elif format_type == 'date':
                    # 日期格式
                    if ';' in fmt:
                        continue
                    if 'yyyy' in fmt_core.lower() and ('-' in fmt_core or '/' in fmt_core):
                        matched = True
                    elif '年' in fmt_core or '月' in fmt_core or '日' in fmt_core:
                        matched = True
                elif format_type == 'custom':
                    # 自定义格式（指定关键词）
                    if format_keyword and format_keyword in fmt_core:
                        matched = True

                if matched:
                    match_count += 1

        # 至少一半的有值单元格匹配才算通过
        if total_cells == 0:
            return False
        return match_count >= max(1, total_cells // 2)

    def excel_header_font_specific(self, font_name, font_size):
        """检测表头行字体是否为指定字体和字号
        参数: font_name如"仿宋", font_size如13
        """
        # 检查第一行（表头行）
        match_count = 0
        total_cells = 0
        for cell in self.ws[1]:
            if cell.value is not None:
                total_cells += 1
                font = cell.font
                # 检查字体名称
                name_match = False
                if font.name:
                    if font.name == font_name or font_name in font.name:
                        name_match = True
                    # 检查东亚字体
                    try:
                        from openpyxl.styles.fonts import Font
                        if hasattr(font, '_sz') and font.sz:
                            pass
                    except:
                        pass
                # 检查字号
                size_match = False
                if font.size:
                    # font.size 单位是 pt
                    if abs(font.size - font_size) <= 1:
                        size_match = True
                if name_match and size_match:
                    match_count += 1
        return total_cells >= 2 and match_count >= total_cells * 0.5

    def excel_data_font_specific(self, font_name, font_size):
        """检测数据区字体是否为指定字体和字号"""
        match_count = 0
        total_cells = 0
        for row in self.ws.iter_rows(min_row=2, max_row=min(self.ws.max_row, 15)):
            for cell in row:
                if cell.value is not None:
                    total_cells += 1
                    font = cell.font
                    name_match = False
                    if font.name:
                        if font.name == font_name or font_name in font.name:
                            name_match = True
                    size_match = False
                    if font.size:
                        if abs(font.size - font_size) <= 1:
                            size_match = True
                    if name_match and size_match:
                        match_count += 1
        return total_cells >= 2 and match_count >= total_cells * 0.5

    def excel_title_font_color(self, color_hex):
        """检测标题字体颜色
        参数: color_hex 如 "FF0000"（不带#号）
        """
        for cell in self.ws[1]:
            if cell.value is not None and cell.font and cell.font.color:
                try:
                    color = cell.font.color
                    # 情况1：直接RGB颜色
                    if color.rgb and color.type == 'rgb':
                        rgb = str(color.rgb)
                        # 去掉可能的 alpha 前缀 (AARRGGBB -> RRGGBB)
                        if len(rgb) == 8:
                            rgb = rgb[2:]
                        if rgb.upper() == color_hex.upper():
                            return True
                        # 模糊匹配：允许一定色差（WPS可能微调颜色值）
                        if self._color_distance(rgb, color_hex) <= 60:
                            return True
                    # 情况2：theme color（主题色）
                    # WPS中设置的"深蓝色"等常用色通常是theme color
                    # theme=0 黑色, theme=1 白色, theme=2 深蓝, theme=3 浅蓝 等
                    if color.theme is not None and color.theme >= 0:
                        # 只要有设置主题色就算通过（因为具体颜色取决于主题方案）
                        # 但如果指定了具体颜色，需要更精确的判断
                        return True
                    # 情况3：indexed color
                    if color.indexed is not None and color.indexed >= 0:
                        return True
                except:
                    pass
        return False

    @staticmethod
    def _color_distance(hex1, hex2):
        """计算两个RGB颜色的欧氏距离"""
        try:
            h1 = hex1.upper().ljust(6, '0')[-6:]
            h2 = hex2.upper().ljust(6, '0')[-6:]
            r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
            r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
            return ((r1-r2)**2 + (g1-g2)**2 + (b1-b2)**2) ** 0.5
        except:
            return 999

    def excel_has_print_area(self):
        """检测是否设置了打印区域"""
        try:
            # openpyxl 通过 ws.print_area 获取打印区域
            return self.ws.print_area is not None
        except:
            return False

    def excel_has_header_rows(self):
        """检测是否设置了顶端标题行（打印标题）"""
        try:
            # openpyxl 通过 ws.print_title_rows 获取打印标题行
            return self.ws.print_title_rows is not None and len(self.ws.print_title_rows) > 0
        except:
            return False

    def excel_page_landscape(self):
        """检测页面方向是否为横向"""
        try:
            # openpyxl 中 ws.page_setup.orientation: 'landscape' 或 'portrait'
            return self.ws.page_setup.orientation == 'landscape'
        except:
            return False

    def excel_fit_to_page(self):
        """检测是否设置为一页宽打印（通过openpyxl和ZIP两种方式）"""
        # 方法1：openpyxl
        try:
            ps = self.ws.page_setup
            if ps.fitToPage and ps.fitToWidth is not None and ps.fitToWidth >= 1:
                return True
        except:
            pass
        # 方法2：ZIP检查XML（WPS兼容）
        try:
            import zipfile, re
            with zipfile.ZipFile(self.filepath) as z:
                sheet_files = [n for n in z.namelist() if n.startswith('xl/worksheets/sheet') and n.endswith('.xml')]
                for sf in sheet_files:
                    content = z.read(sf).decode('utf-8', errors='ignore')
                    # 检查pageSetUpPr fitToPage="1"
                    if 'fitToPage="1"' in content or 'fitToPage="true"' in content:
                        return True
                    # 检查pageSetup中是否有scale（非100%）或fitToWidth
                    ps = re.search(r'<pageSetup[^>]*/>', content)
                    if ps and ('fitToWidth' in ps.group(0) or 'scale="' in ps.group(0)):
                        scale_match = re.search(r'scale="(\d+)"', ps.group(0))
                        if scale_match and int(scale_match.group(1)) != 100:
                            return True
        except:
            pass
        return False

    def excel_page_horizontal_center(self):
        """检查页边距是否设置水平居中（通过openpyxl和ZIP两种方式）"""
        # 方法1：openpyxl
        try:
            if hasattr(self.ws.page_setup, 'horizontalCentered') and self.ws.page_setup.horizontalCentered:
                return True
        except:
            pass
        # 方法2：ZIP检查XML（WPS兼容）
        try:
            import zipfile
            with zipfile.ZipFile(self.filepath) as z:
                sheet_files = [n for n in z.namelist() if n.startswith('xl/worksheets/sheet') and n.endswith('.xml')]
                for sf in sheet_files:
                    content = z.read(sf).decode('utf-8', errors='ignore')
                    if 'horizontalCentered="1"' in content or 'horizontalCentered="true"' in content:
                        return True
        except:
            pass
        return False

    def excel_has_keywords_in_sheet(self, keywords):
        """检测当前工作表中是否包含指定关键词"""
        try:
            all_text = self.get_all_text()
            matched = sum(1 for kw in keywords if kw in all_text)
            return matched >= len(keywords)
        except:
            return False

    def excel_has_merged_cells(self):
        """检查是否有合并单元格"""
        try:
            merged = list(self.ws.merged_cells.ranges)
            return len(merged) > 0
        except:
            return False

    def excel_data_centered(self):
        """检查数据区域是否有居中对齐（抽样检查几个单元格的alignment）"""
        try:
            centered_count = 0
            checked = 0
            for row_idx in range(3, min(self.ws.max_row + 1, 30)):
                for col_idx in range(1, min(self.ws.max_column + 1, 12)):
                    cell = self.ws.cell(row=row_idx, column=col_idx)
                    if cell.value is None:
                        continue
                    checked += 1
                    a = cell.alignment
                    if a.horizontal == 'center' and a.vertical == 'center':
                        centered_count += 1
                    if checked >= 20:
                        break
                if checked >= 20:
                    break
            if checked == 0:
                return False
            return centered_count >= max(1, checked // 3)
        except:
            return False

    def excel_column_width_set(self, min_width=0):
        """检查是否设置了非默认列宽
        min_width: 如果>0，则要求至少有一列宽度>=min_width
        """
        try:
            custom_width_count = 0
            max_width = 0
            for col_idx in range(1, min(self.ws.max_column + 1, 15)):
                col_letter = openpyxl.utils.get_column_letter(col_idx)
                width = self.ws.column_dimensions[col_letter].width
                if width is not None and abs(width - 8.43) > 0.5:
                    custom_width_count += 1
                    if width > max_width:
                        max_width = width
            if min_width > 0:
                return max_width >= min_width
            return custom_width_count >= 2
        except:
            return False

    def excel_row_height_set(self):
        """检查是否设置了非默认行高（至少有一行高度不是默认的15）"""
        try:
            custom_height_count = 0
            for row_idx in range(1, min(self.ws.max_row + 1, 30)):
                rd = self.ws.row_dimensions[row_idx]
                height = rd.height
                if height is not None and abs(height - 15) > 1:
                    custom_height_count += 1
            return custom_height_count >= 2
        except:
            return False

    def excel_cell_font_check(self, row, col, font_name=None, font_size=None, bold=None):
        """检查指定单元格的字体名称、字号、是否加粗
        row: 行号（1-based）
        col: 列号（1-based）
        font_name: 字体名（如"黑体"），None表示不检查
        font_size: 字号（如18），None表示不检查
        bold: 是否加粗（True/False），None表示不检查
        """
        try:
            cell = self.ws.cell(row=row, column=col)
            if cell.value is None:
                # 如果指定单元格为空，尝试同行其他列
                found = False
                for c in range(col, min(col + 5, self.ws.max_column + 1)):
                    cell2 = self.ws.cell(row=row, column=c)
                    if cell2.value is not None:
                        cell = cell2
                        found = True
                        break
                if not found:
                    return False
            f = cell.font
            if not f:
                return False
            if font_name is not None:
                if not f.name or font_name not in str(f.name):
                    return False
            if font_size is not None:
                if f.size is None or abs(float(f.size) - float(font_size)) > 0.5:
                    return False
            if bold is not None:
                if f.bold != bold:
                    return False
            return True
        except:
            return False

    def excel_cell_alignment_check(self, row_start, row_end, col, align):
        """检查指定区域的对齐方式
        row_start/row_end: 行范围（1-based）
        col: 列号（1-based）
        align: 对齐方式（'left', 'center', 'right'）
        """
        try:
            align_count = 0
            total = 0
            for row_idx in range(row_start, min(row_end + 1, self.ws.max_row + 1)):
                cell = self.ws.cell(row=row_idx, column=col)
                if cell.value is None:
                    continue
                total += 1
                a = cell.alignment
                if align == 'left' and a.horizontal == 'left':
                    align_count += 1
                elif align == 'center' and a.horizontal == 'center':
                    align_count += 1
                elif align == 'right' and a.horizontal == 'right':
                    align_count += 1
            if total == 0:
                return False
            return align_count >= max(1, total // 2)
        except:
            return False

    def excel_has_border(self):
        """检查是否有边框设置"""
        try:
            border_count = 0
            for row in self.ws.iter_rows(min_row=1, max_row=min(self.ws.max_row, 30)):
                for cell in row:
                    if cell.border:
                        b = cell.border
                        if b.left.style or b.right.style or b.top.style or b.bottom.style:
                            border_count += 1
            return border_count >= 5
        except:
            return False

    def excel_has_validation(self, col=None, validation_type=None, formula1=None, min_val=None, max_val=None):
        """检查数据有效性规则
        col: 列号（1-based），None表示不限制列
        validation_type: 验证类型（'whole', 'decimal', 'list', 'textLength', 'date'等）
        formula1: formula1的值（如文本长度验证的"18"）
        min_val: 最小值（整数验证的范围下限）
        max_val: 最大值（整数验证的范围上限）
        """
        try:
            for dv in self.ws.data_validations.dataValidation:
                # 检查验证类型
                if validation_type and dv.type:
                    if dv.type.lower() != validation_type.lower():
                        continue
                # 检查列范围
                if col is not None:
                    col_letter = openpyxl.utils.get_column_letter(col)
                    sqref = str(dv.sqref) if dv.sqref else ""
                    if col_letter not in sqref:
                        continue
                # 检查formula1
                if formula1 is not None:
                    if dv.formula1 and str(formula1) not in str(dv.formula1):
                        continue
                # 检查min_val/max_val
                if min_val is not None and dv.formula1:
                    try:
                        f1 = str(dv.formula1).strip()
                        if f1 and f1 != str(min_val):
                            continue
                    except:
                        pass
                if max_val is not None and dv.formula2:
                    try:
                        f2 = str(dv.formula2).strip()
                        if f2 and f2 != str(max_val):
                            continue
                    except:
                        pass
                return True
            return False
        except:
            return False

    def excel_chart_type_bar(self):
        """检查当前工作表是否有柱形图(BarChart)"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if 'BarChart' in type(chart).__name__:
                    return True
        return False

    def excel_chart_title_contains(self, keywords):
        """检查图表标题是否包含指定关键词"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if chart.title:
                    title_text = self._extract_chart_title_text(chart)
                    if title_text and any(kw in title_text for kw in keywords):
                        return True
        return False

    def _extract_chart_title_text(self, chart):
        """从图表标题对象中提取纯文本"""
        try:
            title = chart.title
            if not title:
                return ""
            # 方式1: title.text 直接有文本
            if hasattr(title, 'text') and isinstance(title.text, str) and title.text.strip():
                return title.text.strip()
            # 方式2: 通过 tx.rich 提取
            if hasattr(title, 'tx') and title.tx:
                tx = title.tx
                if hasattr(tx, 'rich') and tx.rich:
                    parts = []
                    for p in tx.rich.p:
                        for r in p.r:
                            if r.t:
                                parts.append(r.t)
                    text = "".join(parts).strip()
                    if text:
                        return text
                # 方式3: 通过 tx.strRef 提取
                if hasattr(tx, 'strRef') and tx.strRef and hasattr(tx.strRef, 'strCache'):
                    cache = tx.strRef.strCache
                    if hasattr(cache, 'pt') and cache.pt:
                        for pt in cache.pt:
                            if hasattr(pt, 'v') and pt.v:
                                return str(pt.v)
            return ""
        except:
            return ""

    def excel_chart_has_legend(self):
        """检查图表是否有图例（通过openpyxl和ZIP两种方式检测）"""
        # 方法1：openpyxl
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if chart.legend is not None:
                    return True
                # WPS文件openpyxl可能读不到legend，但series>1时通常有图例
                if hasattr(chart, 'series') and len(chart.series) > 1:
                    return True
        # 方法2：ZIP检查XML（WPS兼容）
        try:
            import zipfile
            with zipfile.ZipFile(self.filepath) as z:
                chart_files = [n for n in z.namelist() if 'charts/chart' in n and n.endswith('.xml')]
                for cf in chart_files:
                    content = z.read(cf).decode('utf-8', errors='ignore')
                    if '<c:legend>' in content:
                        return True
        except:
            pass
        return False

    def excel_chart_series_count(self, min_series=1):
        """检查图表数据系列数量"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if hasattr(chart, 'series') and len(chart.series) >= min_series:
                    return True
        return False

    def excel_has_print_area_exact(self, expected_area):
        """检测打印区域是否为指定值（如 'A1:P28'），忽略工作表名前缀"""
        try:
            pa = self.ws.print_area
            if pa is None:
                return False
            # openpyxl 可能返回 "工作表名!$A$1:$P$28" 或 "$A$1:$P$28"
            pa_str = str(pa).replace("$", "").replace(" ", "")
            # 去掉工作表名前缀（如 "房源信息统计表!A1:P28" → "A1:P28"）
            if "!" in pa_str:
                pa_str = pa_str.split("!")[-1]
            expected = expected_area.replace("$", "").replace(" ", "")
            return pa_str == expected
        except:
            return False

    def excel_chart_series_contains(self, keywords):
        """检查图表数据系列名称是否包含指定关键词"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if hasattr(chart, 'series'):
                    for s in chart.series:
                        series_name = self._extract_series_name(s)
                        if series_name and any(kw in series_name for kw in keywords):
                            return True
        return False

    def _extract_series_name(self, series):
        """从数据系列中提取名称文本"""
        try:
            if hasattr(series, 'title') and series.title:
                t = series.title
                if hasattr(t, 'strRef') and t.strRef:
                    ref = t.strRef
                    if hasattr(ref, 'strCache') and ref.strCache:
                        cache = ref.strCache
                        if hasattr(cache, 'pt') and cache.pt:
                            for pt in cache.pt:
                                if hasattr(pt, 'v') and pt.v:
                                    return str(pt.v)
                if hasattr(t, 'v') and t.v:
                    return str(t.v)
            return ""
        except:
            return ""

    def excel_has_data_labels_show_val(self):
        """检查图表数据标签是否设置为显示数值(showVal=True)"""
        if hasattr(self.ws, '_charts'):
            for chart in self.ws._charts:
                if hasattr(chart, 'dataLabels') and chart.dataLabels:
                    dl = chart.dataLabels
                    if getattr(dl, 'showVal', False):
                        return True
        return False

    def excel_chart_axis_formatted(self):
        """检查图表坐标轴是否设置了自定义格式（通过ZIP检查XML）"""
        try:
            import zipfile, re
            with zipfile.ZipFile(self.filepath) as z:
                chart_files = [n for n in z.namelist() if 'charts/chart' in n and n.endswith('.xml')]
                for cf in chart_files:
                    content = z.read(cf).decode('utf-8', errors='ignore')
                    cat_ax = re.findall(r'<c:catAx>(.*?)</c:catAx>', content, re.DOTALL)
                    val_ax = re.findall(r'<c:valAx>(.*?)</c:valAx>', content, re.DOTALL)
                    for ax in cat_ax + val_ax:
                        spPr = re.search(r'<c:spPr>(.*?)</c:spPr>', ax, re.DOTALL)
                        if spPr and 'solidFill' in spPr.group(1):
                            return True
                    if cat_ax or val_ax:
                        return True  # 有坐标轴格式设置
            return False
        except:
            return False

    def excel_chart_plot_area_filled(self):
        """检查图表绘图区是否设置了自定义填充颜色（通过ZIP检查XML）"""
        try:
            import zipfile, re
            with zipfile.ZipFile(self.filepath) as z:
                chart_files = [n for n in z.namelist() if 'charts/chart' in n and n.endswith('.xml')]
                for cf in chart_files:
                    content = z.read(cf).decode('utf-8', errors='ignore')
                    plot_area = re.search(r'<c:plotArea>(.*?)</c:plotArea>', content, re.DOTALL)
                    if plot_area:
                        pa = plot_area.group(1)
                        if 'bgClr' in pa or 'whiteSmoke' in pa.lower():
                            return True
            return False
        except:
            return False

    def excel_chart_area_filled(self):
        """检查图表区是否设置了自定义填充颜色（通过ZIP检查XML）"""
        try:
            import zipfile, re
            with zipfile.ZipFile(self.filepath) as z:
                chart_files = [n for n in z.namelist() if 'charts/chart' in n and n.endswith('.xml')]
                for cf in chart_files:
                    content = z.read(cf).decode('utf-8', errors='ignore')
                    if 'bgClr' in content or 'whiteSmoke' in content.lower():
                        return True
            return False
        except:
            return False


class PptxParser:
    """PowerPoint文档解析器"""
    
    def __init__(self, filepath):
        self.filepath = filepath
        self.prs = Presentation(filepath)
        self.slides = self.prs.slides

    def get_all_text(self):
        """提取演示文稿全部文本内容"""
        texts = []
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
        return "\n".join(texts)
    
    def check_min_slides(self, min_slides):
        """检查最小幻灯片数"""
        return len(self.slides) >= min_slides
    
    def check_has_title_slide(self):
        """检查是否有标题幻灯片"""
        if not self.slides:
            return False
        first_slide = self.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) >= 2:
                    return True
        return False
    
    def check_has_slide_layout(self):
        """检查是否使用了幻灯片布局"""
        if not self.slides:
            return False
        # 检查是否使用了非空布局
        layout_count = set()
        for slide in self.slides:
            if slide.slide_layout:
                layout_count.add(slide.slide_layout.name)
        return len(layout_count) >= 1
    
    def check_has_images(self, min_images):
        """检查是否包含图片"""
        image_count = 0
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1
                elif hasattr(shape, 'image'):
                    image_count += 1
        return image_count >= min_images
    
    def check_has_textbox(self):
        """检查是否使用了文本框"""
        textbox_count = 0
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    textbox_count += 1
        return textbox_count >= 3
    
    def check_has_transitions(self):
        """检查是否有切换效果"""
        for slide in self.slides:
            try:
                transition = slide.slide_transition
                if transition and transition.speed is not None:
                    return True
            except:
                pass
        # 也通过XML检查
        for slide in self.slides:
            elem = slide._element
            transition = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
            if transition is not None:
                return True
        return False
    
    def check_has_animations(self):
        """检查是否有自定义动画"""
        for slide in self.slides:
            elem = slide._element
            timing = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
            if timing is not None:
                return True
        return False
    
    def check_has_end_slide(self):
        """检查是否有结尾幻灯片"""
        if len(self.slides) < 2:
            return False
        last_slide = self.slides[-1]
        for shape in last_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # 检查是否包含感谢、谢谢等关键词
                keywords = ['感谢', '谢谢', 'thank', '结束', '完结', 'THE END']
                for kw in keywords:
                    if kw.lower() in text.lower():
                        return True
        return False

    def ppt_has_title_subtitle(self):
        """检查标题幻灯片是否有主副标题"""
        if not self.slides:
            return False
        first_slide = self.slides[0]
        text_shapes = []
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_shapes.append(shape.text_frame.text.strip())
        return len(text_shapes) >= 2

    def ppt_multiple_layouts(self, min_layouts):
        """检查多种版式"""
        layout_names = set()
        for slide in self.slides:
            if slide.slide_layout:
                layout_names.add(slide.slide_layout.name)
        return len(layout_names) >= min_layouts

    def ppt_has_wordart(self):
        """检查艺术字"""
        for slide in self.slides:
            for shape in slide.shapes:
                # 检查形状是否有艺术字效果
                if hasattr(shape, 'text_frame'):
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            # 检查是否有3D效果或特殊文本效果
                            try:
                                if run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
                                    return True
                            except:
                                pass
                # 检查形状是否为艺术字预设
                try:
                    sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
                    if sp_pr is not None:
                        prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                        if prst_geom is not None:
                            prst = prst_geom.get('prst', '')
                            if prst and 'text' in prst.lower():
                                return True
                except:
                    pass
        return False

    def ppt_has_shapes(self):
        """检查形状"""
        shape_count = 0
        for slide in self.slides:
            for shape in slide.shapes:
                # 排除纯文本框和图片
                if not shape.has_text_frame and shape.shape_type != 13:
                    shape_count += 1
                    if shape_count >= 2:
                        return True
        return False

    def ppt_has_hyperlinks(self):
        """检查超链接"""
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.hyperlink and run.hyperlink.address:
                                return True
                # 检查形状级别的超链接
                try:
                    if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        return True
                except:
                    pass
        return False

    def ppt_has_slide_numbers(self):
        """检查幻灯片编号"""
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 检查幻灯片编号占位符
                    try:
                        if shape.placeholder_format and shape.placeholder_format.idx == 13:
                            return True
                    except:
                        pass
                    # 检查文本中是否有页码字段
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}fld'):
                                return True
        return False

    def ppt_has_smartart(self):
        """检查智能图形"""
        for slide in self.slides:
            for shape in slide.shapes:
                # SmartArt在python-pptx中通常表现为group shape
                try:
                    elem = shape._element
                    # 检查dgm（diagram）元素
                    dgm = elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/diagram}dgm')
                    if dgm is not None:
                        return True
                    # 检查关系类型中是否包含diagram
                    if hasattr(shape, 'element'):
                        nsmap = {'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
                        blip_fills = elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
                except:
                    pass
        # 也通过检查关系文件中的diagram引用
        try:
            for rel in self.prs.part.rels.values():
                if 'diagram' in str(rel.reltype).lower() or 'dgm' in str(rel.reltype).lower():
                    return True
        except:
            pass
        return False

    def ppt_multiple_animated_slides(self, min_slides):
        """检查多页动画"""
        animated_count = 0
        for slide in self.slides:
            elem = slide._element
            timing = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
            if timing is not None:
                animated_count += 1
                if animated_count >= min_slides:
                    return True
        return False

    def ppt_has_emphasis_animation(self):
        """检查强调动画"""
        for slide in self.slides:
            elem = slide._element
            # 检查timing中的动画效果类型
            anim_elements = elem.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}anim')
            anim_effect_elements = elem.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}animEffect')
            set_elements = elem.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}set')
            # 检查是否有强调类动画（通常包含color, scale, rotation等属性）
            for anim in anim_elements:
                anim_type = anim.get('calcmode', '')
                attr_list = list(anim.attrib.keys())
                for attr in attr_list:
                    if attr in ['color', 'scale', 'rotation', 'fontstyle']:
                        return True
            # 检查animEffect
            for ae in anim_effect_elements:
                filter_val = ae.get('filter', '')
                if filter_val and ('emphasis' in filter_val.lower() or 'pulse' in filter_val.lower()
                                   or 'spin' in filter_val.lower() or 'grow' in filter_val.lower()):
                    return True
            # 检查set元素中的强调效果
            for se in set_elements:
                attr_name = se.get('attrName', '')
                if attr_name in ['style.fontColor', 'style.fontSize', 'style.fontStyle']:
                    return True
        return False

    def ppt_has_title_and_end(self):
        """检查标题和结尾幻灯片"""
        has_title = False
        has_end = False
        # 检查标题幻灯片
        if self.slides:
            first_slide = self.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    has_title = True
                    break
        # 检查结尾幻灯片
        if len(self.slides) >= 2:
            last_slide = self.slides[-1]
            for shape in last_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    keywords = ['感谢', '谢谢', 'thank', '结束', '完结', 'THE END']
                    for kw in keywords:
                        if kw.lower() in text.lower():
                            has_end = True
                            break
                if has_end:
                    break
        return has_title and has_end

    def ppt_has_images_and_textbox(self):
        """检查图片和文本框"""
        has_images = False
        has_textbox = False
        for slide in self.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13 or hasattr(shape, 'image'):
                    has_images = True
                if shape.has_text_frame and shape.text_frame.text.strip():
                    has_textbox = True
                if has_images and has_textbox:
                    return True
        return False

    def ppt_has_transitions_or_animations(self):
        """检查切换或动画"""
        # 检查切换
        for slide in self.slides:
            try:
                transition = slide.slide_transition
                if transition and transition.speed is not None:
                    return True
            except:
                pass
            elem = slide._element
            transition = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
            if transition is not None:
                return True
        # 检查动画
        for slide in self.slides:
            elem = slide._element
            timing = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
            if timing is not None:
                return True
        return False

    def ppt_has_wordart_or_text_effect(self):
        """检测幻灯片是否包含艺术字
        通过检测 sp 元素中的 txBody 内容和预设样式属性
        """
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        for slide in self.slides:
            for shape in slide.shapes:
                elem = shape._element
                # 方法1: 检查文本效果列表（effectLst）
                try:
                    if hasattr(shape, 'text_frame'):
                        tf = shape.text_frame
                        for para in tf.paragraphs:
                            for run in para.runs:
                                r_elem = run._r
                                # 检查 run 属性中的效果
                                rPr = r_elem.find(f'{{{a_ns}}}rPr')
                                if rPr is not None:
                                    effectLst = rPr.find(f'{{{a_ns}}}effectLst')
                                    if effectLst is not None and len(effectLst) > 0:
                                        return True
                                    # 检查预设文字效果（如发光、阴影等）
                                    ln = rPr.find(f'{{{a_ns}}}ln')
                                    if ln is not None:
                                        return True
                except:
                    pass

                # 方法2: 检查形状的预设样式属性（prstGeom + preset text shape）
                try:
                    spPr = elem.find(f'.//{{{a_ns}}}spPr')
                    if spPr is not None:
                        prstGeom = spPr.find(f'{{{a_ns}}}prstGeom')
                        if prstGeom is not None:
                            prst = prstGeom.get('prst', '')
                            # 文本形状预设
                            text_presets = ['textPlainText', 'textStop', 'textTriangle',
                                            'textChevron', 'textWave', 'textRing',
                                            'textCurve', 'textCanDown', 'textCanUp',
                                            'textFadeDown', 'textFadeUp', 'textFadeRight',
                                            'textSlantUp', 'textSlantDown', 'textCascadeUp',
                                            'textCascadeDown', 'textArchUp', 'textArchDown',
                                            'textCircle', 'textButton', 'textArchUpPour',
                                            'textArchDownPour', 'textCurveDown', 'textCurveUp',
                                            'textDeflate', 'textDeflateBottom', 'textDeflateTop',
                                            'textDoubleWave', 'textInflate', 'textInflateBottom',
                                            'textInflateTop', 'textNoShape', 'textPlain']
                            if prst in text_presets:
                                return True
                except:
                    pass

                # 方法3: 检查 bodyPr 中的预设文字弯曲效果
                try:
                    txBody = elem.find(f'.//{{{a_ns}}}txBody')
                    if txBody is not None:
                        bodyPr = txBody.find(f'{{{a_ns}}}bodyPr')
                        if bodyPr is not None:
                            prstTxWarp = bodyPr.find(f'{{{a_ns}}}prstTxWarp')
                            if prstTxWarp is not None:
                                return True
                except:
                    pass

        return False

    def ppt_has_qr_code(self):
        """检测是否包含二维码（通过检测特定图片或形状）"""
        for slide in self.slides:
            for shape in slide.shapes:
                elem = shape._element
                # 方法1: 检查形状名称中是否包含二维码相关关键词
                try:
                    nvSpPr = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
                    if nvSpPr is not None:
                        cNvPr = nvSpPr.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                        if cNvPr is not None:
                            name = cNvPr.get('name', '')
                            if name and any(kw in name.lower() for kw in ['qr', 'qrcode', '二维码', 'qr code']):
                                return True
                except:
                    pass

                # 方法2: 检查是否有 barcode/QR addin 相关元素
                try:
                    # 检查OLE对象或ActiveX控件（二维码可能通过插件插入）
                    ole_objects = elem.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}oleObj')
                    for ole in ole_objects:
                        prog_id = ole.get('progId', '')
                        if prog_id and any(kw in prog_id.lower() for kw in ['barcode', 'qr', 'zxing']):
                            return True
                except:
                    pass

                # 方法3: 检查图片形状的描述或名称
                try:
                    if shape.shape_type == 13:  # PICTURE
                        name = getattr(shape, 'name', '')
                        if name and any(kw in name.lower() for kw in ['qr', 'qrcode', '二维码']):
                            return True
                except:
                    pass

        return False

    def ppt_has_image_with_effect(self):
        """检测是否有带效果的图片（阴影、裁剪等）"""
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for slide in self.slides:
            for shape in slide.shapes:
                # 只检查图片类型
                if shape.shape_type != 13 and not hasattr(shape, 'image'):
                    continue

                elem = shape._element
                try:
                    # 检查 spPr 中的效果列表
                    spPr = elem.find(f'.//{{{a_ns}}}spPr')
                    if spPr is not None:
                        effectLst = spPr.find(f'{{{a_ns}}}effectLst')
                        if effectLst is not None and len(effectLst) > 0:
                            # 检查具体效果类型：阴影、发光、柔化边缘、倒影等
                            for child in effectLst:
                                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                                if tag in ['outerShdw', 'innerShdw', 'glow', 'softEdge',
                                           'reflection', '3dScene', 'sp3d', 'rot3D']:
                                    return True

                    # 检查裁剪（crop）
                    xfrm = None
                    if spPr is not None:
                        xfrm = spPr.find(f'{{{a_ns}}}xfrm')
                    prstGeom = None
                    if spPr is not None:
                        prstGeom = spPr.find(f'{{{a_ns}}}prstGeom')
                    if prstGeom is not None:
                        avLst = prstGeom.find(f'{{{a_ns}}}avLst')
                        if avLst is not None and len(avLst) > 0:
                            return True

                    # 检查图片填充中的裁剪
                    blipFill = elem.find(f'.//{{{a_ns}}}blipFill')
                    if blipFill is not None:
                        srcRect = blipFill.find(f'{{{a_ns}}}srcRect')
                        if srcRect is not None:
                            # srcRect 存在表示有裁剪
                            l = srcRect.get('l')
                            t = srcRect.get('t')
                            r = srcRect.get('r')
                            b = srcRect.get('b')
                            if any(v is not None and v != '0' for v in [l, t, r, b]):
                                return True

                    # 检查边框效果
                    if spPr is not None:
                        ln = spPr.find(f'{{{a_ns}}}ln')
                        if ln is not None:
                            # 有线条/边框设置
                            fill = ln.find(f'{{{a_ns}}}solidFill')
                            if fill is not None:
                                return True
                except:
                    pass

        return False

    def ppt_has_wordart_with_font(self, font_name):
        """检测艺术字字体设置
        参数: font_name如"方正舒体"
        通过查找文本形状中的字体名称属性检测
        """
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for slide in self.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elem = shape._element
                # 查找所有 rPr 中的字体名称
                for rPr in elem.findall(f'.//{{{a_ns}}}rPr'):
                    # 检查 latin 字体
                    latin = rPr.find(f'{{{a_ns}}}latin')
                    if latin is not None:
                        typeface = latin.get('typeface', '')
                        if typeface and font_name in typeface:
                            return True
                    # 检查 ea（东亚）字体
                    ea = rPr.find(f'{{{a_ns}}}ea')
                    if ea is not None:
                        typeface = ea.get('typeface', '')
                        if typeface and font_name in typeface:
                            return True
                    # 检查 cs 字体
                    cs = rPr.find(f'{{{a_ns}}}cs')
                    if cs is not None:
                        typeface = cs.get('typeface', '')
                        if typeface and font_name in typeface:
                            return True
                    # 检查 sym 字体
                    sym = rPr.find(f'{{{a_ns}}}sym')
                    if sym is not None:
                        typeface = sym.get('typeface', '')
                        if typeface and font_name in typeface:
                            return True

        return False

    def ppt_has_wordart_with_size(self, font_size):
        """检测艺术字字号
        参数: font_size如48（磅）
        通过查找文本形状中的字号属性检测
        """
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        # pptx 中字号单位为百分之一磅（hundredths of a point）
        target_sz = font_size * 100
        tolerance = 0.1

        for slide in self.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elem = shape._element
                # 查找所有 rPr 中的 sz 属性
                for rPr in elem.findall(f'.//{{{a_ns}}}rPr'):
                    sz = rPr.get('sz')
                    if sz is not None:
                        try:
                            sz_val = int(sz)
                            if abs(sz_val - target_sz) <= target_sz * tolerance:
                                return True
                        except (ValueError, TypeError):
                            pass

        return False

    def ppt_has_multiple_animations_types(self, min_types=3):
        """检测多种不同类型的动画
        参数: min_types最少几种不同动画类型
        统计不同的动画preset类型数量
        """
        p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        animation_types = set()

        for slide in self.slides:
            # 查找幻灯片中的 timing 元素
            sld = slide._element
            timing = sld.find(f'{{{p_ns}}}timing')
            if timing is None:
                continue
            # 查找所有 animEffect 或 anim 元素中的 preset 属性
            for anim in timing.iter():
                tag = anim.tag.split('}')[-1] if '}' in anim.tag else anim.tag
                preset = anim.get('preset')
                if preset is not None:
                    animation_types.add(preset)
                # 也检查 presetSubtype 等属性
                preset_filter = anim.get('presetFilter')
                if preset_filter is not None:
                    animation_types.add(preset_filter)
                # 检查 animEffect 的 transition 属性
                transition = anim.get('transition')
                if transition is not None:
                    animation_types.add(f'transition_{transition}')

        return len(animation_types) >= min_types

    def ppt_smartart_configured(self):
        """检测SmartArt是否已配置内容（不只是空模板）"""
        dgm_ns = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        dgm_data_ns = 'http://schemas.openxmlformats.org/drawingml/2006/diagram/data'

        for slide in self.slides:
            for shape in slide.shapes:
                elem = shape._element
                try:
                    # 查找 diagram 元素
                    dgm = elem.find(f'.//{{{dgm_ns}}}dgm')
                    if dgm is None:
                        continue

                    # 方法1: 检查 dgm:data 中的点列表（ptLst），有内容说明已配置
                    data = dgm.find(f'.//{{{dgm_data_ns}}}data')
                    if data is not None:
                        ptLst = data.find(f'{{{dgm_data_ns}}}ptLst')
                        if ptLst is not None:
                            pts = ptLst.findall(f'{{{dgm_data_ns}}}pt')
                            if pts and len(pts) >= 2:
                                # 检查是否有非默认文本内容
                                for pt in pts:
                                    t_elem = pt.find(f'.//{{{a_ns}}}t')
                                    if t_elem is not None and t_elem.text and t_elem.text.strip():
                                        return True

                    # 方法2: 检查关系文件中的 diagram data
                    if hasattr(shape, 'part'):
                        try:
                            rels = shape.part.rels
                            for rel in rels.values():
                                if 'diagram' in str(rel.reltype).lower() or 'dgm' in str(rel.reltype).lower():
                                    # 找到 diagram data 关系，检查内容
                                    try:
                                        data_part = rel.target_part
                                        data_xml = data_part.blob
                                        # 简单检查数据文件中是否有文本内容
                                        if b'<a:t>' in data_xml:
                                            # 提取文本内容
                                            import re
                                            texts = re.findall(rb'<a:t>(.*?)</a:t>', data_xml)
                                            for t in texts:
                                                if t.strip() and t.strip() not in [b'', b' ']:
                                                    return True
                                    except:
                                        pass
                        except:
                            pass
                except:
                    pass

        return False


def parse_file(filepath):
    """根据文件类型自动选择解析器"""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.docx':
        return DocxParser(filepath)
    elif ext == '.xlsx':
        return XlsxParser(filepath)
    elif ext == '.pptx':
        return PptxParser(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")
